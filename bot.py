"""
AKADEMIK IŞLER BOT — ŞAHSY WERSIÝA (tölegsiz)
Bot → TÜRKMEN dilinde | Faýl → HEMIŞE RUS dilinde
Render.com | Python 3.11

🔧 DÜZEDILEN WERSIÝA:
- Accept-Encoding: identity → "invalid distance too far back" çözüldi
- Şablon global cache → çalt we durnukly
- F.video handler dubl düzedildi
- Token/API key diňe env-den
- Redis URL env-den
- _copy_template_title TÄZE şablon (30 paragraph) indekslerine görä täzelendi ✅
"""

import asyncio, base64, copy, io, json, logging, os, re, sqlite3, time
import httpx
from aiogram import Bot, Dispatcher, F, Router
from aiogram.filters import CommandStart
from aiogram.fsm.context import FSMContext
from aiogram.fsm.state import State, StatesGroup
from aiogram.fsm.storage.memory import MemoryStorage
try:
    from aiogram.fsm.storage.redis import RedisStorage
    _HAS_REDIS = True
except ImportError:
    _HAS_REDIS = False
from aiogram.types import (
    BufferedInputFile, CallbackQuery,
    InlineKeyboardButton, InlineKeyboardMarkup, Message,
)
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt as PPtx
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.text import PP_ALIGN
from aiogram.types import LabeledPrice, PreCheckoutQuery
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_LINE_SPACING
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm, Pt, RGBColor

logging.basicConfig(level=logging.INFO, format="%(asctime)s | %(levelname)s | %(message)s")
log = logging.getLogger(__name__)

# ── SAZLAMALAR ──────────────────────────────────────────────
# ✏️ DIŇE ŞU ÝERLERI ÜÝTGEDIN
BOT_TOKEN        = os.getenv("BOT_TOKEN", "")
DEEPSEEK_API_KEY = os.getenv("DEEPSEEK_API_KEY", "")
REDIS_URL        = os.getenv("REDIS_URL", None)

ADMIN_IDS        = [8512644114, 7404431806]
INTRO_VIDEO_URL  = "https://youtu.be/Lm3D5o-4gcM?si=rjMoWCbWourFm8C9"
DEEPSEEK_URL     = "https://api.deepseek.com/v1/chat/completions"
DEEPSEEK_MODEL   = "deepseek-chat"
PRICE            = {"referat": 299, "doklad": 299, "pptx": 299}  # rubl
PRICE_STARS      = {"referat": 5, "doklad": 5, "pptx": 5}  # Telegram Stars
CARD_NUMBER      = "2202 2084 5873 0067"
PHONE_NUMBER     = "+7 922 309 80 64"
CARD_HOLDER      = "Мекан Н"
CONTACT_PHONE    = "+7 922 309 80 64"   # ulanyjyny netije alanyndan soň habarlaşmak üçin
STABILITY_API_KEY = os.getenv("STABILITY_API_KEY", "")
DB_PATH           = os.getenv("DB_PATH", "bot.db")
STABILITY_URL    = "https://api.stability.ai/v2beta/stable-image/generate/core"


TEMPLATE_B64 = (
    "UEsDBBQABgAIAAAAIQBKvAJxbQEAACgGAAATAAgCW0NvbnRlbnRfVHlwZXNdLnhtbCCiBAIooAACAAAA"
    "AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAA"
    "AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAA"
    "AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAA"
    "AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAA"
    "AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAA"
    "AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAA"
    "AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAA"
    "AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAA"
    "AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAC0lMtqwzAQRfeF/oPRtthKuiilxMmij2UbaPoBijRO"
    "RGVJSJPX33ccO6aUJC5NvTHYM/feM2OY0WRbmmQNIWpnczbMBiwBK53SdpGzj9lLes+SiMIqYZyFnO0g"
    "ssn4+mo023mICaltzNkS0T9wHuUSShEz58FSpXChFEivYcG9kJ9iAfx2MLjj0lkEiylWHmw8eoJCrAwm"
    "z1v6XJMEMJElj3VjlZUz4b3RUiDV+dqqHylpk5CRct8Tl9rHG2pg/GhCVTkd0OjeaDVBK0imIuCrKKmL"
    "b1xQXDm5KkmZnbc5wumKQkto9ZWbD05CjLTz0mRtpRTaHvhPckTcGYj/T1H7dscDIgn6AGicOxE2MH/v"
    "jeKbeSdI4Rxah338jda6EwKs6onh4PyLPVCimBvoYw+NdScE0iWC+jm8mGNvcy6SOqfB+UiXLfxh7MPp"
    "qtQpDewhoD6/6TaRrC+eD6qrqEAdyeb7Oz/+AgAA//8DAFBLAwQUAAYACAAAACEAHpEat+8AAABOAgAA"
    "CwAIAl9yZWxzLy5yZWxzIKIEAiigAAIAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAA"
    "AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAA"
    "AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAA"
    "AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAA"
    "AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAA"
    "AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAA"
    "AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAA"
    "AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAA"
    "AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAKySwWrD"
    "MAxA74P9g9G9UdrBGKNOL2PQ2xjZBwhbSUwT29hq1/79PNjYAl3pYUfL0tOT0HpznEZ14JRd8BqWVQ2K"
    "vQnW+V7DW/u8eACVhbylMXjWcOIMm+b2Zv3KI0kpyoOLWRWKzxoGkfiImM3AE+UqRPblpwtpIinP1GMk"
    "s6OecVXX95h+M6CZMdXWakhbeweqPUW+hh26zhl+CmY/sZczLZCPwt6yXcRU6pO4Mo1qKfUsGmwwLyWc"
    "kWKsChrwvNHqeqO/p8WJhSwJoQmJL/t8ZlwSWv7niuYZPzbvIVm0X+FvG5xdQfMBAAD//wMAUEsDBBQA"
    "BgAIAAAAIQCVqf+Z1gsAAIWHAAARAAAAd29yZC9kb2N1bWVudC54bWzsXVtv28gVfi/Q/yDoqX1wxPtF"
    "WGdB8bIbIC2CJNs+LmiKsthIokpSlt2nxN5b4bQG2reii2TRRV8WBRxnvfEmTgLsLyD/Uc8ZkrJkyTZF"
    "y7KsUA+8zOXM5XznzJkzM9RHH2+2W6UN2/Mdt7Napm9R5ZLdsdy601lfLX/20FiRyiU/MDt1s+V27NXy"
    "lu2XP77961991K/WXavXtjtBCUh0/Gq/a62Wm0HQrVYqvtW026Z/q+1Ynuu7jeCW5bYrbqPhWHal73r1"
    "CkPRFHnqeq5l+z6Up5qdDdMvJ+SszWzU6p7Zh8xIkKtYTdML7M0TGvTURPiKXJHGCTE5CEELGXqcFDs1"
    "KaGCtRojxOUiBLUao8TnozShcUI+Ssw4JTEfJXackpSP0hic2uMAd7t2ByIbrtc2A3j11itt03vU664A"
    "4a4ZOGtOywm2gCYlpGRMp/MoR40g14BCm61PTUGstN263WLrKRV3tdzzOtUk/8ogP1a9GudPboMcditb"
    "sVCcXLE3g5YfpHm9LH0XZ9cSxUJ6reLZLehHt+M3ne5AO7TzUoPIZkpk47wO2Gi30nT9Lp1R1M5SbVrM"
    "hhOCWaqf8K7dimt+PkWaysBNJDHIkaUKo2WmNWkDgk8KztU1Q51LZ1Q+KQFmjIBg2RkHi5SGlNCoWCfS"
    "jXScjGKV0om5gnSck46lM+rA05UZIuDXg3pzKipM2q8VzGsGZtP0B0BHivZ0leIH5LbaQ33UXb+cIHzi"
    "ub3uCTXnctTunKjEPhonU9BKBGpYyP3LVeZB0+yCpmxb1TvrHdcz11pQIxCPEiC8RDiAVwAK3sijvUnC"
    "kdcl1DHl22BVrbn1Lbx3IY6rdk3PvAOgFFhVYxWDKZNQGJMCDBWTH4RWwYKr318tU5ShsVKNGgTd8zBQ"
    "kQSBkwaBmt0we61gPPm9oSBSi3se3vyuaUETIZHZCGykiBlaDnY6ww1e7vewzWYvcMsVzPYnCyI2TBg0"
    "LFDAtheHejFNz3A7gY8kfcsBIDx02rZf+r3dL91322YHadqmHyi+Y06MbCodf3I2yx8PJiWvkWvLJC0h"
    "9fJ6K/c/w8hKUq/KoM1rrvsIx/IHARgBkB51A2lpx2xDKz/vN5qNrXVOaAiCRUp1W4bj+dipyetdM36L"
    "y07I6Z36gFjcHRM5tfAd1a8Gt8N/h0fh2/AoehJth4fRY3I/CN+Vwnfhi+hxuB++gqcDuGOivVJ4VIKn"
    "/WgnfI2Pz8N3kOMJ0PgZrq8h6c+l8PvwMPwRiUG6ryDqCEsM4nJj9oxLB8/yNU6VhEI6LsP0jHKxiHi9"
    "CKmo4wmbgAFdz/Ztb8Mu3x7GWvgmegrYfBcelsKXBJc7ELU/wPQhRCbR++EBBL0jr8dp4CnAozwMkYx2"
    "oq8h+jD8CctDaYDADMAWKZmVDdoogF0A+wwVfBDtAkS/AVS9PFPxZkCaoHAqrwtqgbS5GRiLiLdsA/9k"
    "dfrLD+EzokyPyWh+hKM5jvY4iqO2THRstIsRaDM8Cd8QBTyiMQd5M4CW0lWDo7hCPRagzaIq3wPS3gDy"
    "0Fb9kmjGr+EpxRsM+DguH8SGbJJs+5c3GXDIGjqjGRpX4PAybPT/ktaHkdIQFaswEpYJp2dYU1qNVlSl"
    "YNOisGkx1ckMGggaw0StNKo0grVWckvSrLX+CCT6YH1xDI91Cra6wPv6pnke6yHbXXPL7aFTIs7QcDbt"
    "+kmk6z5K81HwQ8IN9Ezcd/sp6MzhNxKpuq1euzMUPxLQcT+tmZ2BB8T9Q/oWN3K4SZ94Th0f1+EONOIG"
    "MpQkxhUcCeZYlnRnTCLNGQxQMSJbQ5JMi5LKy9oF7rCHBFMSIzAsQzAVpMzxPrWd9SZ2YMxTXji7w7Fu"
    "aT4rvqZvCfdI405zD7MlCUe1EMOxMqXJOSzt8xXOUDsT4QIM+sk9bVbLbgSYq+sC+mWZTWuapMypoq5L"
    "Vsf1Sddz3YbuIYBibvhdu9UiLrykqotQX2KNfA8G8WswOtDmJabGiKFxVlP0DhH0WKUQjJ0DTBAvOTsw"
    "2Zpa02m2VgDzwwbmPwGUx4nrYBtdYPGcLTc8iQIltDHx+SbWJGXPyAIj1WTtfFxmV/YiM3NlD3SUlrPe"
    "mUxworSJOq9LkoC0rljaCtmZn+z8C+TlC3RqoBNuTgo9B/iomsLqbCpoVzwTuvGIXHIP8v/I0lu8mnEE"
    "z/u4TgdjwHviSD4gTrp4xYIEQbJoN/omCcah4TiNOApf4fIfoD9ZCAyPR2RgpsMCL3OCYrCIwlkMC6LE"
    "z28OwHEyo3H0Ms0BHLKu7WGXrpZXaCqZJhejSI5R5FuQuzHpm2Y0WaDGVCdogDFxMHSKVnilEIdCHCaI"
    "wzMcXMCsOsK58hKJQS5r78wxRWB5Q2X0HEI0E5tu6ef0y20FstItir1F0aVffhg3CDMtggk1lRdqDFvg"
    "r8Df9PgD2P03PIq+IPMI3A6wV4r+lnVikgGdlKGotKHgKkuBzgKdU6ITd6/Ek1zcr4L3o2gPprh4/bKE"
    "3lKYEONWWNw4uE0wTKAZ7U1UncnAj7e4gEGnTDUDpnhV5xjxAo17ratgeXxTYk1VVVm/JkldDnsfKpRd"
    "KC9hjeZxfNMUI9LKfI4yfCi+x0l6BW/xjoNTu+U5TdQ1ptgtfym2XP1uJaVGK5pUbG5cbDZRGq+yFJt7"
    "2i1JDK2jMpw9m1jhguGKbPy5S9JKPH2agWtu0JyWfbl3b66p/qz5NfkIF42VSY5wrW+IzJ85y0JbdVLF"
    "aU9wJV3mTeTozes4Yuk+Dw9xBZOcv9omG7hLZGv2cbRTHXMOnNpAG/2VZHxRAjsY8r6CWdrL6HG0B09v"
    "U4MYjyi8hsCd8Cco4e+YBfM+iW3raDeTr4HhaoymijlsxGnl66KJGy/Kg01ws/AYI3GwtXg2scWXeYBT"
    "VVpiLtrEVzBxoZnICCqnsGwOv0rBxIVhokDVeEPUZ74XqWDiHJnIGgbH8sbMt28WTJwjE2lBFRVJziGJ"
    "M5n0FZy9unkiJ9cERbyu6fx1c/ZkWkl8pEvE7v4SH2KafNC3NOUv/Ee0mxy+xOW6N1m2BPEyT3MMvdyj"
    "GbIWOODgtxZlecwBU0jKTZeU8Ltom3xI5DB8ix4R9IbsgCi8j3ZHz7enlxveC7iZ+D/hs/DbFYZdYcIX"
    "GQSd1gTakI3c20ZYmZJikSuk//oFfLF3Fc4O5t9FO2QTylSnopak9WcYBZl6YvmQ8Bzsmlfh8cmHTKZB"
    "xBm2D1PTFVGez16lwiAqDKIrkYzPh38ZDAFRZAydvbaPqMwP9QuG53ij0BDTk+XOEaYnYTcY1TNrJgL5"
    "N2Q6+yNc8TtWT3H5csjOD/d/mwHvFMtxHE3lwHvh1bsGXT2ZiYIq6oasYYcUTLypTOR1jpH4WrH8dZOZ"
    "yMo0VROu7eNeBWevTjxptiZReb6UUaycFNOfxfIHk6Os8TcmyUpIpo9AMzVZVaULvn9TTIwKpC/ARD/L"
    "TJ9XdZnljKXaM1ZM6pduUp+c94sD8bxp+o3qvSzTe1FWaNGgLvg8TQHyRQL5ZE5yLGgrJs/4W3BysThJ"
    "sYJGsSr2SsHJm81JieVlg1kq5+kHykmVETSZyrH/o+Dkgo2TMm/oqpLD4inmqTdBUClKqPF0jmPjcz96"
    "nOfk6qy5snaaK1fjHKVVXlN1cQ77ZAuuZOcKz6sizdNz8HAUXJmCK7REKQoZYq5jgCpYNYlV/dOOJVUU"
    "dEUgHbcEzUQ/UfgMzz7gf0AmfyOJXzcL30VfkdfX4f6oB/WsH7FcYspnIpyiaVHl52BML+afmF25CuFk"
    "rcYzau6tmRLFKxx6kj60Xr9Q7Bfon+RmJ/bxvyE+LTEUMyK96WU5GipM1Ey+bQX3BkzOKBEPIBNBg8Iz"
    "egL99QdYzT7YurRMxXuE8aCsFP91FCb4nYnlBG4X07DoFEq/xCbxpDwYOQO3fRIbr7HTIkW+BtO0zTqO"
    "xyLx1lcbrkuG5+R1vRckozUpzXJb2EPJ0i6mIcF118K/r0LaMIbfcwKrGX+AJ+6ruDPIY/wH7xWSpdcG"
    "Mb79fwAAAP//AwBQSwMEFAAGAAgAAAAhABEXoNkNAQAAOQQAABwACAF3b3JkL19yZWxzL2RvY3VtZW50"
    "LnhtbC5yZWxzIKIEASigAAEAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAA"
    "AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAA"
    "AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAA"
    "AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAA"
    "AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAArJPNTsMwEITvSLyDtXfipEBBqE4vqFKvEB7A"
    "TTY/IllH9hbI22MVpaRQRT34uGPtzOexvFp/da34QOsaQwqSKAaBlJuioUrBW7a5eQThWFOhW0OoYEAH"
    "6/T6avWCrWa/5Oqmd8K7kFNQM/dPUrq8xk67yPRI/qQ0ttPsR1vJXufvukK5iOOltFMPSE88xbZQYLfF"
    "LYhs6PESb1OWTY7PJt93SHwmQn7i7hWZ/eWct9W2QlYwESPvCPI8yENIEPa7+ItwGH/EZI5hEZLB/Wti"
    "VOYQkqAIPLQ4BTjMc/HLkPGlIc70rp28xFGag7gPCYFUkOFpC6Myh3AXtgfDfxiO0gghTz58+g0AAP//"
    "AwBQSwMEFAAGAAgAAAAhAK7IR2fCAgAAzQsAABIAAAB3b3JkL2Zvb3Rub3Rlcy54bWzUlktvozAQgO8r"
    "7X9A3FMDSZOAmlTqZrvqrWp3f4BrTEDFD9kmJP9+xzzT0q0IPS0HHmPP55nxzOCb2yPLnQNVOhN84/pX"
    "nutQTkSc8f3G/fP7frZ2HW0wj3EuON24J6rd2+33bzdllAhhuDBUO8DgOool2bipMTJCSJOUMqyvWEaU"
    "0CIxV0QwJJIkIxSVQsUo8HyvepNKEKo1LPgD8wPWboMjx3G0WOESlC1wgUiKlaHHnuFfDLlGIVoPQcEE"
    "EHgY+EPU/GLUElmrBqDFJBBYNSBdTyN94NxyGikYklbTSPMhaT2NNEgnNkxwISmHwUQohg18qj1iWL0W"
    "cgZgiU32kuWZOQHTW7YYnPHXCRaBVkdg8/hiwgoxEdN8HrcUsXELxaNGf9bpW9OjWr95dBo0H7csLBci"
    "ejS5Nq2uGhO7Wn0nSMEoN1XUkKI5xFFwnWay6w5sKg0G0xZy+CwAB5a380rpjyy1f7W2Xb0NPXCM+c3e"
    "sby2/HOi743YTYvoNMaY8HbN1hIGGdwvPCk0Z8H1RzafFhAMAEtCR/4sWsa6YSDSV7flZCPLquXUu2I5"
    "WR9Yf2QPfG/MGUDHJk4vogRtXJHVxQanWHeJbon0MqOuO9yJncVI7r9WCL+UKGRPy75Ge+hbYmlPJxew"
    "moI6L3L9NWOeUyyhUzISPey5UPglB4ugPBzIcKfaAXuHRLGP6pUeK7nda8f2GHd7dqxyysicJCA0lVhh"
    "I5QLIpugM7+aKEF5EdmxBxCu7tah74VwXLNS+GmZStpcVhXOePHTxvU8LwxWP3edaEcTXORmOPJoRfe7"
    "+frOqxd8VPahJSbgLkzCiaHQ1j2rkGd2A4JF9/FUWP9xYYSLtjeoU68ZrU/1kKonVPfW/w9jQQQ3GS+q"
    "/8Hz+7g0Vr4Jy3LlezsvDP+PsHzo3mchOvvQ278AAAD//wMAUEsDBBQABgAIAAAAIQB8juxPwAIAAMcL"
    "AAARAAAAd29yZC9lbmRub3Rlcy54bWzUlttuozAQhu9X2ndA3KcGcqKoSaU221Xvqu3uA7jGCVbxQbbJ"
    "4e13zDEt3YrQq81FgLHn8/j3zMDN7ZHn3p5qw6RY+eFV4HtUEJkysVv5f34/TGLfMxaLFOdS0JV/osa/"
    "XX//dnNIqEiFtNR4gBAmOSiy8jNrVYKQIRnl2FxxRrQ0cmusiORIbreMUHSQOkVREAblndKSUGNgvXss"
    "9tj4NY4ch9FSjQ/g7IAzRDKsLT12jPBiyBxdo7gPikaAYIdR2EdNL0YtkIuqB5qNAkFUPdJ8HOmDzS3G"
    "kaI+aTmONO2T4nGkXjrxfoJLRQUMbqXm2MKj3iGO9WuhJgBW2LIXljN7AmawaDCYidcREYFXS+DT9GLC"
    "EnGZ0nyaNhS58gstktp/0vq70JPKv760HjQftiwsd43o0ebGNr56iHaV+0aSglNhS9WQpjnoKIXJmGq7"
    "Ax9Lg8Gsgew/E2DP82beQYUDS+1frW1THUMHHBJ+fXY8ryL/nBgGA07TIVqPISG8XbOJhEMGdwuPkuZM"
    "3HBg82kAUQ+wIHTgy6JhxDUDka66HYcNLKuGU52K47BO2HBgD3wfzBnApDbNLqJEja7I+WKLM2zaRHdE"
    "ellQ8xZ34mcaqd3XCuGnloXqaOxrtMeuJR7cx8kFrLqgzovcfC2Y5wwr6JScJI87ITV+ySEiKA8PMtwr"
    "T8D9Q6K4S3lLj6XdnbXneoy/7r6qvENiTwoIhiqssZXaB5PLz0lYzlPgO0vc2CMY50Ecx/fL0C+t8M6y"
    "zrqsf84VvvDSXys/CILraPlj05o2dIuL3PZHnpzpYTON74JqwSftLkZhAruFSXhrKXT1wDnkzOkfzdqH"
    "X4XbPi6s9NH6BrXuFaPZUzWkqwnlf739j5QgUlgmivJl8PxelTrGN6Isojic3YUP/4coH27vE4G6e7P+"
    "CwAA//8DAFBLAwQUAAYACAAAACEA5aDr0vwGAAD6IAAAFQAAAHdvcmQvdGhlbWUvdGhlbWUxLnhtbOxZ"
    "W48bNRR+R+I/WPOeZmaSyaVqinKltLtt1d0W8ehNnBk3nvHIdnYbISTUPvECQiqIB5CAFx4QYiWKQAjE"
    "X1h+Q6VWXH4EtmeSGScObekWVWg30saX7xx/Puf4+GTmwmt3YgIOEeOYJh3HO+c6ACVjOsFJ2HFu7o8q"
    "LQdwAZMJJDRBHWeBuPPaxVdfuQDPiwjFCEj5hJ+HHScSIj1frfKxHIb8HE1RIuemlMVQyC4LqxMGj6Te"
    "mFR9121UY4gTByQwlmpPvjz5/uTnk2NwbTrFY+RcXOofEvkvEVwNjAnbU9rRUuiL3+6eHJ/8cvLg5Pi3"
    "d2X7F/n9gZadzDz1xRe8Txg4hKTjyKUn9Ggf3REOIJALOdFxXP3nVC9eqK6EiNgiW5Ib6b9cLheYzHwt"
    "x8KDlaA79Ft1b6VfA4jYxA1b6rPSpwFwPJY7z7iUsV7QcFt+ji2BsqZFd7vp1Ux8SX9tU3+70fPrBl6D"
    "smZ9c4+j9nAQGHgNyprBBr7r+r12zcBrUNZsbODrw27THxp4DYoITmab6Eaz1Wrk6BVkSsklK7zdaLjN"
    "QQ4vUNVStGXyiXja2IvhbcpGUkA7GwqcALFI0RSOpVw3FZSDAeYpgQsHpDChXA67vufJQKy7/uqjPQDP"
    "I1iSzobGfGNI8QN8zHAqOs5lqdUpQR79+OPDuw8e3v3h4b17D+9+C3ZwGAmL3CWYhGW5P7/68K/P3gV/"
    "fPf5n/c/suN5Gf/4m/ce//TrP6kXBq2Pjx8/OH70yfu/f33fAu8yeFCG7+MYcXAVHYEbNJYbtCyADtiz"
    "SexHEJcluknIYQKVjAU9FJGBvrqABFpwPWTa8RaT6cMGfH1+2yC8F7G5wBbglSg2gLuUkh5l1j1dUWuV"
    "rTBPQvvibF7G3YDw0LZ2f83Lw3kqzwG2qexHyKB5nUiXwxAlSAA1R2cIWcTewtiw6y4eM8rpVIC3MOhB"
    "bDXJPj4woqkQuoRj6ZeFjaD0t2Gb3VugR4lN/QAdmkh5NiCxqUTEMOPrcC5gbGUMY1JG7kAR2UjuLdjY"
    "MDgX0tMhIhQMJ4hzm8w1tjDoXoEyj1ndvksWsYlkAs9syB1IaRk5oLN+BOPUyhknURn7Bp/JEIXgOhVW"
    "EtQ8Iaov/QCTre6+hZHh7ief7ZsyDdkDRM3Mme1IIGqexwWZQmRT3mWxkWK7DFujozcPjdDeQYjAIzhB"
    "CNx8w4anqWHzgvTlSGaVS8hmm8vQjFXVTxBHQBc7FsdiboTsHgrpFj67i7XEs4BJDNk2zVdnZsgMD5g8"
    "jLZ4JeOZkUoxU4fWTuIaj439bdV6PYJGWKk+t8frghn+e5ozJmVu/wsZ9MwyMrE/tW32ITEWKAJmH2Kw"
    "Y0u3UsRwfyGijpMWm1vlpuahLdxQXSt6Ypw8oQL67yofWV88+vQzC/Z0qh078HnqnG2pZL262YZbr2n6"
    "lE3wy1/SDOA8uY7kLWKBnlU0ZxXN/76i2Xaez+qYszrmrI6xi7yAOqYoXfQDoeVjH60lfupnQFNMyJ5Y"
    "ELTDdRHEZS6YjOSg7mglq0dQaSSb+fIGLmRQtwGj4k0sor0IpnJZT68Q8lx1yEFKuSyk9LBVt5og83iX"
    "TvInfKri0k89pQAUxbgbrMZl0Say0UazeES6Uq97oX4suySgZJ+FRGkxk0TNQqK5HHwCCb2zU2HRtrBo"
    "KfVbWeiv3CvysgJQPUMP6hkjGX4yxCfKT5n80run7ultxjS37Vu211ZcT8fTBolSuJkkSmEYyctkffiU"
    "fd0uXGrQU6bYpNFsvQhfq6SylhtIYvbAkTxztUCqGcO040zlDyjZjFOpj6vMBUmYdJyxyA39bzJLyrgY"
    "QB5lMD2V7T/GAjFAcCxjvewGkhTcPL+p9viSkmu7L5/l9FfZyWg6RWOxZaToyrlMiXX2OcGqQ+eS9F40"
    "OQIHZM5uQGmooOkpA04wFytrTjArBXdhxbV0lR9F491QcUQhSSOY3yjlZJ7BdXtFp7QPzXR9V2Y/38xB"
    "qJz03Lfuk4XURClpbrlA1K1pzx8v7pIvsSryvsEqS93rua69zHXbbonnvxBK1IrFDGqKsYVaMWpSO8WC"
    "oLTcKjS33RGnfRusR626IJZ1pu5tvAanB7dl5A9k9Tongmuq8lcMg/3lC8ssE+jRZXa5I8Cc4Y7ztht0"
    "630/6FfcVjCs1Gt1t9IKurVKNwhq3jDw3EHPf0caRUSxF2Rrj+SPf7LIX/Tr8Y2X/fGy9D43pnGV6nf4"
    "VS2sX/Z7vvGyP3vHD/bVvAOwtMzb/tCr+12/X+kPvEal7g8alVaz1q30/cbA78ok1Bh133HAoQZ7vcFg"
    "NAr8SqMvcXW3G1S6vVq/0mgNe/7IG9YHrgTnyfBOnj5yWywNevFvAAAA//8DAFBLAwQUAAYACAAAACEA"
    "QldEqnkEAAAMDQAAEQAAAHdvcmQvc2V0dGluZ3MueG1stFfbbts4EH1fYP/B0PMqlmTLF6FO4Uu8TRFv"
    "F3GKfaZEyiZCigJJ2XGL/fcdUqJlJ94iaZGXmJwzc2Y0HM4wHz4+cdbZEamoKCZeeBV4HVJkAtNiM/G+"
    "Piz9kddRGhUYMVGQiXcgyvt4/ftvH/aJIlqDmuoARaESnk28rdZl0u2qbEs4UleiJAWAuZAcadjKTZcj"
    "+ViVfiZ4iTRNKaP60I2CYOA1NGLiVbJIGgqf00wKJXJtTBKR5zQjzY+zkK/xW5ssRFZxUmjrsSsJgxhE"
    "oba0VI6N/ywbgFtHsvvRR+w4c3r7MHjF5+6FxEeL14RnDEopMqIUHBBnLkBatI77L4iOvq/Ad/OJlgrM"
    "w8CuTiOP30YQvSAYZOTpbRyjhqMLlqc8FL+NZ3DkoW1iw8HPBXNCoLDG2zexRC6vXWOLNNoidawiw0je"
    "FlR8pDvwNkeKvaZqauiOphLJ+k42JcOz5HZTCIlSBuFA6XTg9Ds2OvMXkmh+7JI8WbnJg3cNPeKbELyz"
    "T0oiM7go0GCiwOsaAMpT5GuNNFAkqiSM2Y6TMYLA4z7ZSMShVziJtcEkRxXTDyhda1GC0g7Bhw0dZbZF"
    "EmWayHWJMmCbi0JLwZweFn8JPYe+I+Fa1Ba5ELoQmvwtT3dgYArKD8+VGrF11n1uSwr8YvOM51zqaM4M"
    "667YrtZ1hwWTAnFI/VnXXAlMTKYqSV9fI8bAZiOMm6RddCRgIkiKyYM58rU+MLKEZK7pNzIt8OdKaQqM"
    "tnf+QgQ/CoAUxvMXKNKHQ0mWBOkKju2dnNnKWDJarqiUQt4WGGr13ZzRPCcSHFCo/RWUM5Vib/P8iSAM"
    "g/id/FaK/APK0CN6D3BNHmdCa8E/Hcot5PrXTtLVclu+8JzAyi3u4aYcVYNoHsTTJlKDtkgwjoY3i0vI"
    "/9v0ggia6CVkGPWG8+lFZNGbzceXkOk0jm6abnKOLOeD8fhibMtFbzRzGWi+myfmMWDudL0yl6fDa4s5"
    "4qmkqLMyz4Wu0Ujl44wWDk8J9GRyiqyr1IG+XwOKI8aWcIwOsAHwBFNVLkhu12yF5KblbTTkRSl01s9H"
    "LtOpifxTiqqs0b1EZX0pnErY7zeWtNB3lDu5qtK1sypgipxAVYG/7KTNU5uefaKhyGxzuUO2WK2urPz7"
    "r6a8UoqhIJH0181JZkyuTV2SFSrLurzTTTjxGN1sdWhMNOwwPDLtJt1EDRZZLKoxu0GZ+VDQbhatLHKy"
    "E72ek/VaWd/J+q0sdrK4lQ2cbGBkW2hoEqbdI9w0tzTyXDAm9gR/avEXojoJaotKsqiHIVSbqAXNdFSd"
    "XUKeYNQSTDW83UuKOYJ3VhhEA2PeaDN0EJU+0zWYUS7PGcyrpOkt3TNjW/HPYjFDOqNQnesDT9vZe1UH"
    "zqiCvlTCmNZCOuwPi4X9BIvs1rwk+rU8XkTzm2B5U8OxHe/ati4493uSz5AiuMGcaVybfo9Gi8ViGQZ+"
    "PB30/f54GPij0Sz2lzfDcLYYj2aDQf/f5s66f2Ou/wMAAP//AwBQSwMEFAAGAAgAAAAhAM45jWxrEAAA"
    "/KoAAA8AAAB3b3JkL3N0eWxlcy54bWzsXdty20YSfd+q/QcUn3YfHEkkRUmuKClJttautR3FkjfPQ2Ao"
    "IgIxXAC0rHz9zg0gyMaA6EGLUVxbrrKIS58ZzOnTmG7cfvz52yIJvvIsj0V6Pjj64XAQ8DQUUZzenw++"
    "3F2/Oh0EecHSiCUi5eeDJ54Pfv7p73/78fF1XjwlPA8kQJq/XoTng3lRLF8fHOThnC9Y/oNY8lRunIls"
    "wQq5mN0fLFj2sFq+CsViyYp4Gidx8XQwPDycDCxM1gVFzGZxyN+IcLXgaaHtDzKeSESR5vN4mZdoj13Q"
    "HkUWLTMR8jyXB71IDN6CxWkFczQGQIs4zEQuZsUP8mBsjzSUND861L8WyRrgGAcwBACTkH/DYZxajANp"
    "WceJIxzOpMKJoxqOX2dqAHlURHMUyrAc1wNlywo2Z/m8jshxnTqu4J4WaowW4ev396nI2DSRSJL1QBIX"
    "aGD1vzx+9Uf/5N/0enUIg5+kFiIRvuEztkqKXC1mN5ldtEv6z7VIizx4fM3yMI7vZAdlK4tYNvjuIs3j"
    "gdzCWV5c5DFr3DhXPxq3hHlRW30ZR/HgQLWY/yE3fmXJ+WA4LNdcqR5srEtYel+u4+mrL7f1ntRWTSXu"
    "+YBlr24vtOHR+HUS37NilckooJY0ggkWWXQlD5Z/K1YsUTsf2FEwf2tjs9xe0r1csjDWnWKzgsuYcDQ5"
    "VD1IYhWChsdn5cLnlWKKrQphG9EA5m8FewDokaFCBo5bE7/kVj77IMIHHt0WcsP5QLclV355f5PFIpMx"
    "6nxwptuUK2/5In4XRxFPazum8zjiv815+iXn0Xr9r9c6ztgVoVil8vfoZKJdJsmjt99CvlRRS25NmSLw"
    "kzJI1N6reN24Nv9vCXZkaWuyn3OmQndwtA2hu4+CGCqLvHa0zZirrWPXe6EaGu2rofG+GjreV0OTfTV0"
    "sq+GTvfVkIZ5zobiNJJnCb0/bAag7sJxqBGN4xAbGsehJTSOQypoHIcS0DgOR0fjOPwYjeNwUwROIUKX"
    "F9acfeTw9nbc3ecIP9zdpwQ/3N1nAD/c3QHfD3d3fPfD3R3O/XB3R28/3N3BGo9rplrBeymztOitspkQ"
    "RSoKHqhJb280lkosnc/S4KmTHs9IDpIAxkQ2eyLujRYyvbzbQ7RI/c/nhUoLAzELZvG9Snl6d5ynX3ki"
    "ljxgUSTxCAEzLpMyx4j4+HTGZzzjacgpHZsOVGWCQbpaTAl8c8nuybB4GhEPX4lIEhQqh5b581yJJCZw"
    "6gULM9G/a4KRxYcPcd5/rBRIcLlKEk6E9YnGxTRW/9xAw/RPDTRM/8xAw/RPDGqcUQ2RRSMaKYtGNGAW"
    "jWjcjH9SjZtFIxo3i0Y0bhat/7jdxUWiQ3x91nHUvXZ3lQh1BaJ3P27j+1RXZXsj2ZppcMMydp+x5TxQ"
    "Jexm2PoxY9u5FNFTcEdxTquQqOb12kVULTtOV/0HdAONSlwVHpG8KjwigVV4/SX2UU6T1QTtHU0+c7ua"
    "Fo2i1UidRHvLkpWZ0PZXGyv6e9haANdxlpPJoBmWwIM/qemsopMi8q172b9ja6z+stqOSqTds5AEvUxE"
    "+EATht89LXkm07KH3kjXIknEI4/oEG+LTBhfq0t+qCnpJPm3i+Wc5bHOlTYgup/qy3sXgo9s2fuAbhIW"
    "pzS8vX21YHES0M0g3t19/BDciaVKM9XA0ABeiqIQCzJMWwn8x298+k+aDl7IJDh9IjraC6LykAa7iglO"
    "MgZJRERIcpoZpzHJOVTj/Zs/TQXLIhq0m4yb24UKToR4yxZLM+kg0JaMi48y/hDMhjTef1gWq7oQlaju"
    "SMBqZcN8Nf2dh/1D3ScRkFSGflkVuv6op7ramg6u/zRhA67/FEGzKU8Pyn8JDnYDrv/BbsBRHexVwvI8"
    "dl5C9cajOtwSj/p4+yd/Fk8kIputEroBLAHJRrAEJBtCkawWaU55xBqP8IA1HvXxErqMxiMoyWm8f2Vx"
    "REaGBqNiQoNR0aDBqDjQYKQE9L9DpwbW/zadGlj/e3UMGNEUoAZG5Wekp3+iqzw1MCo/02BUfqbBqPxM"
    "g1H52ehNwGczOQmmO8XUIKl8rgZJd6JJC75YioxlT0SQbxN+zwgKpAbtJhMz9RyJSM1N3ASQqkadEE62"
    "DRwVyb/xKVnXFBZlvwgqoixJhCCqra1PONpy8961XWb6sY/eXbhJWMjnIol45jgmt63Ml2/NYxnb3dfd"
    "6FT2/BDfz4vgdl5V++swk8OdlmXCvmG2u8GmMZ+UD780mX3kUbxalB2FD1NMRt2NtUdvGI93G69nEhuW"
    "xx0tYZuT3ZbrWfKG5UlHS9jmaUdLrdMNyzY9vGHZQ6MjnLT5T5XjOZzvpM2LKuPGZtscqbJscsGTNi/a"
    "kEpwEYbqagFkp5tm3PbdxOO2x6jIjYKRkxuls67cEG0C+8y/xurMjgmaur3q7gkQ9/UkulPk/HUlTN1+"
    "44JT94e63suJU5rzoBFn1P3C1UaUcY9j53Djhugcd9wQnQOQG6JTJHKao0KSG6VzbHJDdA5Sbgh0tIJn"
    "BFy0gva4aAXtfaIVRPGJVj1mAW6IztMBNwRaqBACLdQeMwU3BEqowNxLqBAFLVQIgRYqhEALFU7AcEKF"
    "9jihQnsfoUIUH6FCFLRQIQRaqBACLVQIgRYqhEAL1XNu7zT3EipEQQsVQqCFCiHQQtXzxR5ChfY4oUJ7"
    "H6FCFB+hQhS0UCEEWqgQAi1UCIEWKoRACxVCoIQKzL2EClHQQoUQaKFCCLRQzaOG/kKF9jihQnsfoUIU"
    "H6FCFLRQIQRaqBACLVQIgRYqhEALFUKghArMvYQKUdBChRBooUIItFD1xcIeQoX2OKFCex+hQhQfoUIU"
    "tFAhBFqoEAItVAiBFiqEQAsVQqCECsy9hApR0EKFEGihQog2/7SXKF232R/hq57OO/a7X7qynfpcf5S7"
    "DjXqDlX2yo3V/VmESyEegsYHD0c63+gGEk+TWOgSteOyeh1X3xKBuvD5y1X7Ez519J4vXbLPQuhrpgB8"
    "3NUS1FTGbS5ftwRJ3rjN0+uWYNY5bou+dUtwGhy3BV2ty/KmFHk6AsZtYaZmfOQwb4vWNXM4xG0xumYI"
    "R7gtMtcM4QC3xeOa4XGggvO29XHHcZpU95cChDZ3rCGcuBHa3BJyVYZjKIyupLkRurLnRuhKoxsBxacT"
    "Bk+sGwrNsBvKj2ooMyzV/kJ1I2CphgheVAMYf6ohlDfVEMqPahgYsVRDBCzV/sHZjeBFNYDxpxpCeVMN"
    "ofyohqcyLNUQAUs1RMBS3fOE7ITxpxpCeVMNofyohpM7LNUQAUs1RMBSDRG8qAYw/lRDKG+qIZQf1SBL"
    "RlMNEbBUQwQs1RDBi2oA4081hPKmGkK1Ua2rKBtUoxiumeMmYTVD3Am5ZogLzjVDj2ypZu2ZLdUQPLMl"
    "yFXJOS5bqpPmRujKnhuhK41uBBSfThg8sW4oNMNuKD+qcdlSE9X+QnUjYKnGZUtOqnHZUivVuGyplWpc"
    "tuSmGpctNVGNy5aaqPYPzm4EL6px2VIr1bhsqZVqXLbkphqXLTVRjcuWmqjGZUtNVPc8ITth/KnGZUut"
    "VOOyJTfVuGypiWpcttRENS5baqIaly05qcZlS61U47KlVqpx2ZKbaly21EQ1LltqohqXLTVRjcuWnFTj"
    "sqVWqnHZUivVuGzpozSJCV4BdbtgWRHQvS/uHcvnBev/csIvacZzkXzlUUB7qB9QR3nwuPH5K4WtP+Qn"
    "9y/kmKk3oNceV4rMG2AtoN7xfVR9pkoZq54E9uthdrXusL1cq39nucyp7T6Hh9dvRqeX2j8PLOSOTlTN"
    "2qvIR6Dh9WevdHtTJo/3FzVOoFupemViw3rlKuX6spmrOcvM1rUTl/tYmbqPcnh1eHxh+2M/a/bA+fKT"
    "bF+vUwuSOZ7rpfUXz6bqbWNyBEbmk2f2A2inVs/CvM/pw9ekaskOpG2j9Vtz7PeWb82pjW/tOrV943Nz"
    "G5brz82p1ZfV5+ZCpf+qX9fjk4n2Gr2zjg3nA6Yjw3q1ul1FAl1eG4T1B+vKy9D1D9aZdbVPyfk4z9Dp"
    "PDY40TjPsIPzrAVr9tuQ6zO7l/2i3k73KmPGd+ZeI0t23b3Mup7uNXK6l70RhMa9Rt+Je5VD7nCvXU60"
    "D1cZ2jndxqcz9bqerjJ2uoq984fGVcYv3FVO655Shn3oKVo+9J4Sm/+vTO/6+k1Pjzh2eoS9o4vGI46/"
    "D4/QKnl5saOnD5iPwzb5gM1vaXxg8sJ9YFz3AacLaFnsNSgcn6l/2w6hvse0doe7WH3n90Lz1dMbTpze"
    "YGsVNN5w8l14QzngzxkQ9sz/qZN/Oyuh4f/0hfK/i3Etgr3qf3ii/nXh/w3FHPHMyb9lhYb/s78o/+UQ"
    "P6fi6RkP5WCz0L6y3VFhs59eqt4dpD+8tO0Lju8zOXi047+LR3e/C1XnbemzrgO3lgZNqdjpaJ09rZgm"
    "hmr5432qHO1ReUnV0+gbM1By+xVPko/M7C2W7l0TPlNykVuPDvWbOre2T81HJ5z2mb464QQ42OyMWWz3"
    "E/MZytg8NuOsxKoSfMNw62e4+o50Rx8OV7kcmlu1w3b/Nmqp2720G4OjYB1/tgJaow5cYcx6uDOEuYPS"
    "/8umaEpNhdNF6ZCIUlun63pW+v4Z7lO5RDJsiowuhkdEDNu6KD3Df1YBoM5Wn+Ihki1T53OxNSZiy5Ym"
    "Xw5b+y7gIVkxtTYXK8dErNjy4PejIXIeTL3LxcOEiAdbovtLqIO+koGkxBSdXJScEFFi62QvVBp/Ogmm"
    "8uMi4ZSIBHsW/Evo4pnz/d2UmGKMi5IzIkrsyL9QXeyrzGZenbE91mZt0xBj62saaU1YQ1HGJmyo2hko"
    "kJkrZqo4JofOFMvVwueVcjK2KkQ5xKkawhVL7Kv8zci9gHs71kekj/pVOSwPPKvGfj2XLtcc2/NtfXZt"
    "1tGJcs1go5f0VWPN1dzO8TKz2v1z1qzh6ive2wRVGyiUXIK1itmyghJzulqYH3ECb7uyG5+5xI2dhQDu"
    "j2wCst/Ed4MSF/l9BbrpRG7OX/jE8Zkpa1am+drANjNmLYUmNVKbIId2MuN5dq3fzab3+D0sLVXuynW7"
    "QJstc8vxofrXhTXqNHg9VI109FVJjVM3CzslsteRa3ZZddVk/cGO7bHSjzusN+/yYTgUI1s/QzlkrK9w"
    "qetT6uV71hXb5nId3aU6aPtGuuo1eduHDd6jh3OUBo9AnSh3e8ceb9OyY9Ec2jY/s7LLPbqEuHpzbZFu"
    "5JNHLC8j/ddcF9X75dKT7Ae7/1C32Kkf0r9UPNHq08PuWRavrqA+c0tKBvbIdj1XoZaMW9U0djrRvdEX"
    "dM2S3qVv8P9Ty6DAj1pdt+/pYEMkOzz2xem+NUauX9rpGsD1Hn2jZHmpDxUlp6ZVO1q5DCrJFVvSjB2Y"
    "RJb3X3pFUlWI4tAR1Q1ORlytwXO7pMWztvBY3vDQcC+J6+G0StkFm+rH5OTfrQAgF5ciV499li/rrO2j"
    "w0e1y9nkVHfhYI3XcJtTsLOIgzzF76g22kHbpsBsoig0lkS6GXENP9abroVQx7x9KDOzGuNNBun/3oTy"
    "ptqgbVNgNvX1Jstvb28qf+U//Q8AAP//AwBQSwMEFAAGAAgAAAAhAO8KKU5OAQAAfgMAABQAAAB3b3Jk"
    "L3dlYlNldHRpbmdzLnhtbJzTX2vCMBAA8PfBvkPJu6bKFClWYQzHXsZg2weI6dWGJbmSi6vu0+/aqXP4"
    "YveS//fjLiHz5c7Z5BMCGfS5GA1TkYDXWBi/ycX722owEwlF5Qtl0UMu9kBiubi9mTdZA+tXiJFPUsKK"
    "p8zpXFQx1pmUpCtwioZYg+fNEoNTkadhI50KH9t6oNHVKpq1sSbu5ThNp+LAhGsULEuj4QH11oGPXbwM"
    "YFlET5Wp6ag112gNhqIOqIGI63H2x3PK+BMzuruAnNEBCcs45GIOGXUUh4/SbuTsLzDpB4wvgKmGXT9j"
    "djAkR547pujnTE+OKc6c/yVzBlARi6qXMj7eq2xjVVSVoupchH5JTU7c3rV35HT2tPEY1NqyxK+e8MMl"
    "Hdy2XH/bdUPYdettCWLBHwLraJz5ghWG+4ANQZDtsrIWm5fnR57IP79m8Q0AAP//AwBQSwMEFAAGAAgA"
    "AAAhAFU6fTT3AQAALgcAABIAAAB3b3JkL2ZvbnRUYWJsZS54bWzck99umzAUxu8n7R0s3zcYEtIUlVTr"
    "2ki72cXUPYBjTLDmP8jHCcnbzzaEpoqmlUmTpoEEh+9wfjrnfHD/cFQSHbgFYXSJ0xnBiGtmKqF3Jf7+"
    "srlZYQSO6opKo3mJTxzww/rjh/uuqI12gHy9hkKxEjfOtUWSAGu4ojAzLdc+WRurqPOPdpcoan/s2xtm"
    "VEud2Aop3CnJCFniAWPfQzF1LRh/MmyvuHaxPrFceqLR0IgWzrTuPbTO2Kq1hnEAP7OSPU9RoUdMurgC"
    "KcGsAVO7mR9m6CiifHlKYqTkKyCfBsiuAEvGj9MYq4GR+MpLjqimcZYjR1QXnD9r5gIAlauaSZTsvNck"
    "1FJHGwrNJZFPayofcScVdqRY8WWnjaVb6UnedeSNQxEcrn7+cIshP0Y9jIDXw6+AukJT5Ss/tc5AlFlD"
    "LfCQOVBZYkJwEt+mSsjTWYVOAPSJVjjWnPUDtSK00qdA7HxiD1tSYv+FEpKtbnGvpIEcj/mgZKNCBmX+"
    "VmGREx/Tu82gvHJin0k/1vV4vi0Zx2upNsDTcbyMPJIlWZBgVX8u+ub/0h6ewxqeN7H/fg+fvXK7yh+v"
    "9nD3+z30nAl7CDajJwGtpKf/1+4XoTigr7xD34yi+hfGZ974Ocm9+bmP55OMt5H7Dxs/BLD+CQAA//8D"
    "AFBLAwQUAAYACAAAACEA+ifkVXwBAAARAwAAEQAIAWRvY1Byb3BzL2NvcmUueG1sIKIEASigAAEAAAAA"
    "AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAA"
    "AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAA"
    "AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAA"
    "AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAA"
    "AAAAAAAAAAAAAAAAjJLfT4MwEIDfTfwfSN+hBXUawtCo8UkTE2c0vtX2xuroj7TdkP/eAhuTaKJvPe67"
    "r8ddi8tPWUdbsE5oNUdpQlAEimkuVDVHz4u7+AJFzlPFaa0VzFELDl2Wx0cFMznTFh6tNmC9ABcFk3I5"
    "M3O08t7kGDu2AkldEggVkkttJfUhtBU2lK1pBTgjZIYleMqpp7gTxmY0op2Ss1FpNrbuBZxhqEGC8g6n"
    "SYoPrAcr3a8FfeYbKYVvDfyK7pMj/enECDZNkzQnPRr6T/Hrw/1T/6uxUN2sGKCy4Cz3wtdQFvhwDCe3"
    "ef8A5ofPYxDOzAL12pZqY1ta8baFbSJhTRU5v6okFXXCtOyr9mS3gzW0jbbcBd8kChgHx6wwPmx2uG3y"
    "IdA1df4hrHopgF+3f138s6BzWNiK7u2UWU+MYbFbxNAs8CgMMB/Gvc+8nNzcLu5QmZFsFpPTOM0W5CLP"
    "znJC3rp+J/UHodw18H/jbGrcC4aRTR9x+QUAAP//AwBQSwMEFAAGAAgAAAAhAPU2xKVyAQAAxgIAABAA"
    "CAFkb2NQcm9wcy9hcHAueG1sIKIEASigAAEAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAA"
    "AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAA"
    "AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAA"
    "AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAA"
    "AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAnFLLTsMwELwj8Q9R7q3TihZU"
    "bV2hVogDj0pN27PlbBILx7Zsg9q/Z9NACOJGTjuzu6PZiWF1anTygT4oa5bpZJylCRppC2WqZbrPH0Z3"
    "aRKiMIXQ1uAyPWNIV/z6CrbeOvRRYUhIwoRlWsfoFowFWWMjwpjahjql9Y2IBH3FbFkqiRsr3xs0kU2z"
    "bM7wFNEUWIxcL5h2iouP+F/RwsrWXzjkZ0d6HHJsnBYR+Uu7qYH1BOQ2Cp2rBvmE6B7AVlQYWq4r4Gh9"
    "QXgyBdaVsK6FFzJSeHx+Q4MDDPfOaSVFpFj5s5LeBlvG5PXiNWn3gQ1HgPzvUL57Fc88AzaE8KQMGZgB"
    "6wpy5kXlhau/7PUIdlJoXNPlvBQ6ILAfAta2ccKQHOsr0nsLe5fbTZvE18pvcnDkUcV654QkC7ez6fDc"
    "QQd2xGJB/nsLPQGP9De8bvVp11RYfM/8bbQBHrpnySfzcUbfJbFvju7u3wv/BAAA//8DAFBLAQItABQA"
    "BgAIAAAAIQBKvAJxbQEAACgGAAATAAAAAAAAAAAAAAAAAAAAAABbQ29udGVudF9UeXBlc10ueG1sUEsB"
    "Ai0AFAAGAAgAAAAhAB6RGrfvAAAATgIAAAsAAAAAAAAAAAAAAAAApgMAAF9yZWxzLy5yZWxzUEsBAi0A"
    "FAAGAAgAAAAhAJWp/5nWCwAAhYcAABEAAAAAAAAAAAAAAAAAxgYAAHdvcmQvZG9jdW1lbnQueG1sUEsB"
    "Ai0AFAAGAAgAAAAhABEXoNkNAQAAOQQAABwAAAAAAAAAAAAAAAAAyxIAAHdvcmQvX3JlbHMvZG9jdW1l"
    "bnQueG1sLnJlbHNQSwECLQAUAAYACAAAACEArshHZ8ICAADNCwAAEgAAAAAAAAAAAAAAAAAaFQAAd29y"
    "ZC9mb290bm90ZXMueG1sUEsBAi0AFAAGAAgAAAAhAHyO7E/AAgAAxwsAABEAAAAAAAAAAAAAAAAADBgA"
    "AHdvcmQvZW5kbm90ZXMueG1sUEsBAi0AFAAGAAgAAAAhAOWg69L8BgAA+iAAABUAAAAAAAAAAAAAAAAA"
    "+xoAAHdvcmQvdGhlbWUvdGhlbWUxLnhtbFBLAQItABQABgAIAAAAIQBCV0SqeQQAAAwNAAARAAAAAAAA"
    "AAAAAAAAACoiAAB3b3JkL3NldHRpbmdzLnhtbFBLAQItABQABgAIAAAAIQDOOY1saxAAAPyqAAAPAAAA"
    "AAAAAAAAAAAAANImAAB3b3JkL3N0eWxlcy54bWxQSwECLQAUAAYACAAAACEA7wopTk4BAAB+AwAAFAAA"
    "AAAAAAAAAAAAAABqNwAAd29yZC93ZWJTZXR0aW5ncy54bWxQSwECLQAUAAYACAAAACEAVTp9NPcBAAAu"
    "BwAAEgAAAAAAAAAAAAAAAADqOAAAd29yZC9mb250VGFibGUueG1sUEsBAi0AFAAGAAgAAAAhAPon5FV8"
    "AQAAEQMAABEAAAAAAAAAAAAAAAAAETsAAGRvY1Byb3BzL2NvcmUueG1sUEsBAi0AFAAGAAgAAAAhAPU2"
    "xKVyAQAAxgIAABAAAAAAAAAAAAAAAAAAxD0AAGRvY1Byb3BzL2FwcC54bWxQSwUGAAAAAA0ADQBAAwAA"
    "bEAAAAAA"
)


# ═══════════════════════════════════════════════════════════════════════
# ✅ SQLite DATABASE — ulanyjylary we sargytlary yzarlamak
# Her ulanyjynyň 1-nji sargydy MUGT, galany tölegli
# ═══════════════════════════════════════════════════════════════════════
def _db_init():
    """DB-ny başlat, tablisany döret (ýok bolsa)"""
    conn = sqlite3.connect(DB_PATH)
    cur = conn.cursor()
    cur.execute("""
        CREATE TABLE IF NOT EXISTS users (
            user_id       INTEGER PRIMARY KEY,
            username      TEXT,
            full_name     TEXT,
            created_at    INTEGER NOT NULL,
            free_used     INTEGER DEFAULT 0,
            total_orders  INTEGER DEFAULT 0,
            total_paid    INTEGER DEFAULT 0
        )
    """)
    cur.execute("""
        CREATE TABLE IF NOT EXISTS orders (
            id          INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id     INTEGER NOT NULL,
            service     TEXT NOT NULL,
            theme       TEXT,
            price       INTEGER NOT NULL,
            is_free     INTEGER DEFAULT 0,
            status      TEXT DEFAULT 'pending',
            created_at  INTEGER NOT NULL
        )
    """)
    conn.commit()
    conn.close()
    log.info(f"✅ SQLite DB taýyn: {DB_PATH}")

def _db_get_or_create_user(user_id: int, username: str = "", full_name: str = "") -> dict:
    """Ulanyjyny DB-den al, ýok bolsa döret. Gaýtarýar: {user_id, free_used, total_orders, ...}"""
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    cur = conn.cursor()
    cur.execute("SELECT * FROM users WHERE user_id = ?", (user_id,))
    row = cur.fetchone()
    if row is None:
        cur.execute(
            "INSERT INTO users (user_id, username, full_name, created_at) VALUES (?, ?, ?, ?)",
            (user_id, username or "", full_name or "", int(time.time()))
        )
        conn.commit()
        cur.execute("SELECT * FROM users WHERE user_id = ?", (user_id,))
        row = cur.fetchone()
        log.info(f"🆕 Täze ulanyjy DB-ä goşuldy: uid={user_id}")
    conn.close()
    return dict(row)

def _db_is_first_free(user_id: int) -> bool:
    """Ulanyjy mugt hyzmata eýe barmy (entek ulanmadyk bolsa True)"""
    conn = sqlite3.connect(DB_PATH)
    cur = conn.cursor()
    cur.execute("SELECT free_used FROM users WHERE user_id = ?", (user_id,))
    row = cur.fetchone()
    conn.close()
    if row is None:
        return True  # täze ulanyjy
    return row[0] == 0

def _db_mark_free_used(user_id: int):
    """Ulanyjynyň mugt hyzmatyny ulanyldy diýip belle"""
    conn = sqlite3.connect(DB_PATH)
    cur = conn.cursor()
    cur.execute("UPDATE users SET free_used = 1 WHERE user_id = ?", (user_id,))
    conn.commit()
    conn.close()

def _db_add_order(user_id: int, service: str, theme: str, price: int, is_free: bool, status: str = "pending") -> int:
    """Täze sargydy DB-ä goş. order_id gaýtaryar."""
    conn = sqlite3.connect(DB_PATH)
    cur = conn.cursor()
    cur.execute(
        "INSERT INTO orders (user_id, service, theme, price, is_free, status, created_at) "
        "VALUES (?, ?, ?, ?, ?, ?, ?)",
        (user_id, service, theme[:200], price, 1 if is_free else 0, status, int(time.time()))
    )
    order_id = cur.lastrowid
    # ulanyjynyň statistikasyny täzele
    cur.execute("UPDATE users SET total_orders = total_orders + 1 WHERE user_id = ?", (user_id,))
    conn.commit()
    conn.close()
    return order_id

def _db_update_order_status(order_id: int, status: str):
    """Sargydyň statusyny üýtgetmek (paid, delivered, cancelled)"""
    conn = sqlite3.connect(DB_PATH)
    cur = conn.cursor()
    cur.execute("UPDATE orders SET status = ? WHERE id = ?", (status, order_id))
    conn.commit()
    conn.close()

def _db_mark_paid(user_id: int, amount: int):
    """Ulanyjynyň töleg mukdaryny jemle"""
    conn = sqlite3.connect(DB_PATH)
    cur = conn.cursor()
    cur.execute("UPDATE users SET total_paid = total_paid + ? WHERE user_id = ?", (amount, user_id))
    conn.commit()
    conn.close()

def _db_stats() -> dict:
    """Admin statistikasy"""
    conn = sqlite3.connect(DB_PATH)
    cur = conn.cursor()
    cur.execute("SELECT COUNT(*) FROM users")
    users_total = cur.fetchone()[0]
    cur.execute("SELECT COUNT(*) FROM orders")
    orders_total = cur.fetchone()[0]
    cur.execute("SELECT COUNT(*) FROM orders WHERE is_free = 1")
    free_orders = cur.fetchone()[0]
    cur.execute("SELECT COUNT(*) FROM orders WHERE is_free = 0 AND status = 'delivered'")
    paid_orders = cur.fetchone()[0]
    cur.execute("SELECT COALESCE(SUM(total_paid), 0) FROM users")
    revenue = cur.fetchone()[0]
    conn.close()
    return {
        "users": users_total, "orders": orders_total,
        "free": free_orders, "paid": paid_orders, "revenue": revenue
    }

_db_init()


# ✅ DÜZEDIŞ #2: Şablon ÝÜKLENIŞI — iki çeşmeli
_TEMPLATE_BYTES_CACHE = None

_TEMPLATE_FILE_PATHS = ["template.docx", "/app/template.docx", "./template.docx"]
for _path in _TEMPLATE_FILE_PATHS:
    try:
        if os.path.exists(_path):
            with open(_path, "rb") as _f:
                _candidate = _f.read()
            if _candidate[:4] == b'PK\x03\x04':
                _TEMPLATE_BYTES_CACHE = _candidate
                log.info(f"✅ Şablon faýldan ýüklendi: {_path} ({len(_candidate)} baýt)")
                break
            else:
                log.warning(f"⚠️ {_path} ZIP däl, indiki çeşme synanyşýar")
    except Exception as e:
        log.warning(f"⚠️ {_path} okalmady: {e}")

if _TEMPLATE_BYTES_CACHE is None:
    try:
        _clean_b64 = "".join(TEMPLATE_B64.split())
        _TEMPLATE_BYTES_CACHE = base64.b64decode(_clean_b64)
        if _TEMPLATE_BYTES_CACHE[:4] != b'PK\x03\x04':
            log.error(f"❌ TEMPLATE_B64 ZIP faýly däl! Başlangyç baýtlar: {_TEMPLATE_BYTES_CACHE[:10]}")
            _TEMPLATE_BYTES_CACHE = None
        else:
            try:
                import zipfile
                with zipfile.ZipFile(io.BytesIO(_TEMPLATE_BYTES_CACHE)) as _zf:
                    _names = _zf.namelist()
                    if "word/document.xml" not in _names:
                        log.error(f"❌ Şablonda word/document.xml ýok! Faýllar: {_names[:5]}")
                        _TEMPLATE_BYTES_CACHE = None
                    else:
                        log.info(f"✅ Şablon base64-den ýüklendi ({len(_TEMPLATE_BYTES_CACHE)} baýt, {len(_names)} içki faýl)")
            except Exception as zex:
                log.error(f"❌ Şablon ZIP hökmünde açylmaýar: {zex}")
                _TEMPLATE_BYTES_CACHE = None
    except Exception as e:
        log.error(f"❌ TEMPLATE_B64 açylmady: {e}")
        _TEMPLATE_BYTES_CACHE = None

if _TEMPLATE_BYTES_CACHE is None:
    log.error("🔴 ŞABLON ÝÜKLENMEDI — Referat/Doklad işlemez! template.docx faýlyny repositora goş.")


def _get_template_bytes() -> bytes:
    if _TEMPLATE_BYTES_CACHE is None:
        raise RuntimeError(
            "Şablon ýüklenmedi! template.docx faýly ýok ýa-da TEMPLATE_B64 bozuk. "
            "Çözgüt: template.docx faýlyny repositoryň köküne goşuň."
        )
    return bytes(_TEMPLATE_BYTES_CACHE)


class St(StatesGroup):
    lang = State()  # UI dili saýlamak
    s01  = State()
    s02  = State()
    s02b = State()
    s03  = State()
    s04  = State()
    s05  = State()
    s06  = State()
    s07  = State()
    s08  = State()
    s09  = State()
    s09b = State()
    s10  = State()
    s11  = State()
    s11b = State()
    s12  = State()
    s13  = State()
    s02b_more = State()
    # Pptx states
    sp1  = State()  # tema
    sp1b = State()  # at-familiýa
    sp1c = State()  # kurs
    sp1d = State()  # gruppa
    sp2  = State()  # dil
    sp3  = State()  # slayd sany


PENDING: dict[int, dict] = {}
SEEN_USERS: set = set()
REQ_ITEMS: dict[int, list] = {}
PAYMENT_PENDING: dict[int, dict] = {}
ACTIVE_GENERATES: set = set()
CANCELLED_GENERATES: set = set()

# UI dil tekstleri
UI_LANG = {
    "tk": {
        "welcome_svc":  "📌 Haýsy hyzmaty isleýärsiňiz?",
        "req_yes":      "✅  Talaplы — öz talaplarym bar",
        "req_no":       "➡️  Talapsyz — adaty GOST",
        "ref_lang_ask": "📌 Referat/Doklad haýsy dilde bolsun?",
        "ref_lang_ru":  "🇷🇺 Rusça",
        "ref_lang_en":  "🇬🇧 Iňlisçe",
        "gen_start":    "Başlanýar...",
        "pay_stars":    "⭐ Stars bilen töle",
        "free_msg":     "🎁 <b>Birinji sargydyňyz MUGT!</b>",
        "paid_msg":     "✅ <b>Taýar boldy!</b>\n\nStars bilen tölemek üçin aşakdaky düwmä basyň:",
        "stars_btn_ref":"⭐ {n} Stars — töle",
        "delivered":    "✅ <b>Taýar!</b>\n\nTäze sargyt: /start",
    },
    "ru": {
        "welcome_svc":  "📌 Какую услугу вы хотите?",
        "req_yes":      "✅  С требованиями",
        "req_no":       "➡️  Без требований — ГОСТ",
        "ref_lang_ask": "📌 На каком языке реферат/доклад?",
        "ref_lang_ru":  "🇷🇺 Русский",
        "ref_lang_en":  "🇬🇧 Английский",
        "gen_start":    "Начинаем...",
        "pay_stars":    "⭐ Оплатить Stars",
        "free_msg":     "🎁 <b>Ваш первый заказ БЕСПЛАТНО!</b>",
        "paid_msg":     "✅ <b>Готово!</b>\n\nДля оплаты нажмите кнопку ниже:",
        "stars_btn_ref":"⭐ {n} Stars — оплатить",
        "delivered":    "✅ <b>Готово!</b>\n\nНовый заказ: /start",
    },
    "en": {
        "welcome_svc":  "📌 Which service do you need?",
        "req_yes":      "✅  With requirements",
        "req_no":       "➡️  Standard GOST",
        "ref_lang_ask": "📌 In which language should it be?",
        "ref_lang_ru":  "🇷🇺 Russian",
        "ref_lang_en":  "🇬🇧 English",
        "gen_start":    "Starting...",
        "pay_stars":    "⭐ Pay with Stars",
        "free_msg":     "🎁 <b>Your first order is FREE!</b>",
        "paid_msg":     "✅ <b>Ready!</b>\n\nPress below to pay with Stars:",
        "stars_btn_ref":"⭐ {n} Stars — pay",
        "delivered":    "✅ <b>Ready!</b>\n\nNew order: /start",
    },
}

def ui(lang, key, **kw):
    t = UI_LANG.get(lang, UI_LANG["tk"]).get(key, key)
    return t.format(**kw) if kw else t

SVC_RU = {"referat": "Реферат", "doklad": "Доклад", "pptx": "Презентация"}
SVC_TM = {"referat": "Referat 📄", "doklad": "Doklad 🎤", "pptx": "Prezentasiýa 📊"}


def kb(*rows) -> InlineKeyboardMarkup:
    return InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text=t, callback_data=c) for t, c in row]
        for row in rows
    ])

# KB_SVC - UI diline görä dinamik
def kb_svc(lang="tk"):
    labels = {
        "tk": ["📄  Referat", "🎤  Doklad", "📊  Prezentasiýa"],
        "ru": ["📄  Реферат", "🎤  Доклад", "📊  Презентация"],
        "en": ["📄  Essay",   "🎤  Report", "📊  Presentation"],
    }
    L = labels.get(lang, labels["tk"])
    return kb([(L[0],"svc:referat")],[(L[1],"svc:doklad")],[(L[2],"svc:pptx")])

KB_SVC = kb_svc()  # default

# Referat dili saýlaw KB
_BACK_LABELS = {"tk":"🔙  Yza — başa gaýt","ru":"🔙  Назад","en":"🔙  Back"}
_SKIP_LABELS = {"tk":"⏭️  Mugallym ýok — geç","ru":"⏭️  Пропустить","en":"⏭️  Skip"}
_CONT_LABELS = {"tk":"✅  Talaplarym taýar — dowam et","ru":"✅  Требования готовы — продолжить","en":"✅  Requirements ready — continue"}

_CRS_LABELS = {
    "tk": [("1-nji","crs:1"),("2-nji","crs:2"),("3-nji","crs:3"),("4-nji","crs:4"),("5-nji","crs:5")],
    "ru": [("1-й","crs:1"),("2-й","crs:2"),("3-й","crs:3"),("4-й","crs:4"),("5-й","crs:5")],
    "en": [("1st","crs:1"),("2nd","crs:2"),("3rd","crs:3"),("4th","crs:4"),("5th","crs:5")],
}
_SEC_LABELS = {
    "tk": [("2 bölüm","sec:2"),("3 bölüm","sec:3")],
    "ru": [("2 раздела","sec:2"),("3 раздела","sec:3")],
    "en": [("2 sections","sec:2"),("3 sections","sec:3")],
}
_SPC_LABELS = {
    "tk": [("📋  Adaty 1.5","spc:default"),("✏️  Üýtgetmek","spc:custom")],
    "ru": [("📋  По умолчанию 1.5","spc:default"),("✏️  Изменить","spc:custom")],
    "en": [("📋  Default 1.5","spc:default"),("✏️  Custom","spc:custom")],
}

def kb_ref_lang(lang="tk"):
    # Referat/doklad diňe rusça
    back = _BACK_LABELS.get(lang,"🔙  Yza — başa gaýt")
    return kb(
        [(ui(lang,"ref_lang_ru"),"reflang:ru")],
        [(back,"back:start")],
    )

def kb_req(lang="tk"):
    back = _BACK_LABELS.get(lang,"🔙  Yza — başa gaýt")
    return kb(
        [(ui(lang,"req_yes"),"req:yes")],
        [(ui(lang,"req_no"),"req:no")],
        [(back,"back:start")]
    )

KB_REQ = kb_req()  # default

def kb_crs(lang="tk"):
    crs = _CRS_LABELS.get(lang, _CRS_LABELS["tk"])
    back = _BACK_LABELS.get(lang,"🔙")
    return kb([crs[0],crs[1],crs[2]],[crs[3],crs[4]],[(back,"back:start")])

def kb_sec(lang="tk"):
    sec  = _SEC_LABELS.get(lang, _SEC_LABELS["tk"])
    back = _BACK_LABELS.get(lang,"🔙")
    return kb(sec,[(back,"back:start")])

def kb_spc(lang="tk"):
    spc  = _SPC_LABELS.get(lang, _SPC_LABELS["tk"])
    back = _BACK_LABELS.get(lang,"🔙")
    return kb(spc,[(back,"back:start")])

def kb_spcv(lang="tk"):
    return kb([("1.0","spv:1.0"),("1.25","spv:1.25"),("1.5","spv:1.5")])

def kb_skip(lang="tk"):
    skip = _SKIP_LABELS.get(lang,"⏭️")
    back = _BACK_LABELS.get(lang,"🔙")
    return kb([(skip,"skip:teacher")],[(back,"back:start")])

def kb_back(lang="tk"):
    return kb([(_BACK_LABELS.get(lang,"🔙  Yza"),"back:start")])

def kb_req_done(lang="tk"):
    cont = _CONT_LABELS.get(lang,"✅")
    back = _BACK_LABELS.get(lang,"🔙")
    return kb([(cont,"req:done")],[(back,"back:start")])

# Yza gabat gelme üçin global-lar saklansyn
KB_CRS  = kb_crs()
KB_SEC  = kb_sec()
KB_SPC  = kb_spc()
KB_SPCV = kb_spcv()
KB_SKIP = kb_skip()
KB_BACK = kb_back()
KB_REQ_DONE = kb_req_done()

def kb_src() -> InlineKeyboardMarkup:
    rows, row = [], []
    for i in range(8, 21):
        row.append((str(i), f"src:{i}"))
        if len(row) == 5:
            rows.append(row); row = []
    if row: rows.append(row)
    return kb(*rows)


async def ask(obj, text: str, markup=None) -> Message:
    kw = dict(parse_mode="HTML", reply_markup=markup)
    if isinstance(obj, CallbackQuery):
        try:
            await obj.message.edit_text(text, **kw); return obj.message
        except Exception:
            return await obj.message.answer(text, **kw)
    return await obj.answer(text, **kw)

def md_clean(t: str) -> str:
    t = re.sub(r"\*\*(.*?)\*\*", r"\1", t)
    t = re.sub(r"\*(.*?)\*", r"\1", t)
    t = re.sub(r"^#{1,6}\s*", "", t, flags=re.M)
    return t.strip()

def spc_float(k: str) -> float:
    return {"default":1.5,"1.0":1.0,"1.25":1.25,"1.5":1.5}.get(k, 1.5)

def spc_str(k: str) -> str:
    return {"default":"1.5","1.0":"1.0","1.25":"1.25","1.5":"1.5"}.get(k, "1.5")

def lv(spc) -> int:
    return {1.0:240, 1.25:300, 1.5:360}.get(spc, 360)


def build_prompt(d: dict) -> str:
    doc_lang = d.get("doc_lang","ru")
    doc_lang_str = {"ru":"РУССКОМ","en":"АНГЛИЙСКОМ"}.get(doc_lang,"РУССКОМ")
    doc_lang_name = {"ru":"русском языке","en":"English language"}.get(doc_lang,"русском языке")

    svc   = SVC_RU.get(d["service"], "Реферат")
    pages = int(d["pages"])
    secs  = int(d["sections"])
    spc   = spc_str(d.get("spacing", "default"))
    # Word-da 1 sahypa (A4, TNR 14pt, 1.5 interval) ≈ 1500 simwol
    # DeepSeek 8000 token ≈ 20000-22000 simwol (rus dili üçin)
    # Takyk sahypa sany üçin: her bölümiň simwol sanyny kesgitle
    cpp   = 1500
    total = pages * cpp
    # Giriş we netije has gysga, esasy bölümler köp
    ic    = int(total * 0.12)          # Giriş: 12%
    cc    = int(total * 0.68 / secs)   # Her bölüm: 68% / bölüm sany
    nc    = int(total * 0.10)          # Netije: 10%
    # max_tokens bilen sazla - 1 token ≈ 2.5 simwol (rus)
    _needed_tokens = int(total / 2.5) + 500  # goşmaça 500 token
    _max_tok = min(8000, max(3000, _needed_tokens))
    chs   = ""
    for i in range(1, secs + 1):
        chs += (f"\n##ГЛАВА_{i}##\n"
                f"Первая строка — заголовок: «{i}. [Название главы]»\n"
                f"Текст: ~{cc} символов, сплошными абзацами. БЕЗ подразделов.\n")
    teacher = d.get("teacher","______________")
    tp      = d.get("teacher_position","")
    tline   = f"{teacher}" if not tp else f"{tp} {teacher}"
    req_text = d.get("req_text", "").strip()
    if req_text:
        extra_block = (
            f"═══ СТУДЕНТИҢ ÝÖRITE TALAPLARY (HÖKMAN ÝERINE ÝETIR) ═══\n"
            f"{req_text}\n"
            f"Ýokardaky talaplary doly ýerine ýetir — olary äsgermezlik etme!\n"
            f"═══════════════════════════════════════════════════════════\n\n"
        )
    else:
        extra_block = ""
    # Dil maglumaty
    _lang_note_ru = f"Напиши {svc} строго на РУССКОМ языке."
    _lang_note_en = f"Write the {svc} STRICTLY IN ENGLISH LANGUAGE. All headings, text, sources — in English."
    _lang_note = _lang_note_en if doc_lang == "en" else _lang_note_ru

    _struct_ru = (
        f"##ВВЕДЕНИЕ##\n~{ic} символов. Актуальность, цель, {secs} задачи, методология.\n"
        f"{chs}\n"
        f"##ЗАКЛЮЧЕНИЕ##\n~{nc} символов. Выводы, итоги, значение.\n\n"
        f"##СПИСОК_ЛИТЕРАТУРЫ##\nРовно {d['sources']} источников, ГОСТ Р 7.0.5-2008.\n"
    )
    _struct_en = (
        f"##INTRODUCTION##\n~{ic} characters. Relevance, goal, {secs} objectives, methodology.\n"
        f"{chs.replace('ГЛАВА','CHAPTER').replace('Первая строка — заголовок','First line is the heading')}\n"
        f"##CONCLUSION##\n~{nc} characters. Findings, summary, significance.\n\n"
        f"##REFERENCES##\nExactly {d['sources']} sources, APA or GOST format.\n"
    )
    _struct = _struct_en if doc_lang == "en" else _struct_ru

    # max_tokens sahypa sanyna görä
    _needed = int(total / 2.5) + 500
    d["_max_tokens_calc"] = min(8000, max(3000, _needed))

    return (
        f"{extra_block}"
        f"Ты академический автор. {_lang_note}\n\n"
        f"ПАРАМЕТРЫ:\n"
        f"Тип: {svc} | Дисциплина: {d['subject']} | Тема: {d['theme']}\n"
        f"Университет: {d['university']}\n"
        f"Студент: {d['fullname']}, {d['course']} курс, гр. {d['group']}\n"
        f"Преподаватель: {tline}\n"
        f"Глав: {secs} | Страниц: {pages} | Источников: {d['sources']} | Интервал: {spc}\n\n"
        f"СТРОГИЙ ОБЪЁМ: РОВНО {pages} страниц A4.\n"
        f"1 страница A4 (TNR 14пт, интервал 1.5) = 1500 символов.\n"
        f"ИТОГО: ~{total} символов. НЕ БОЛЬШЕ и НЕ МЕНЬШЕ.\n"
        f"• Введение: строго ~{ic} символов ({round(ic/1500,1)} стр.)\n"
        f"• Каждая глава: строго ~{cc} символов ({round(cc/1500,1)} стр.)\n"
        f"• Заключение: строго ~{nc} символов ({round(nc/1500,1)} стр.)\n\n"
        f"СТРУКТУРА — строго эти маркеры:\n\n"
        f"{_struct}"
        f"Формат: «N. Автор. Название. — Город : Издательство, Год. — Стр. с.»\n"
        f"Только 1. 2. 3. — НЕТ 1.1/1.2\n\n"
        f"ПРАВИЛА:\n"
        f"• только русский язык\n"
        f"• НЕТ подразделов 1.1/1.2\n"
        f"• Все перечисления нумеруй: 1. 2. 3. (не используй маркеры •/—)\n"
        f"• ОБЯЗАТЕЛЬНО напиши ВСЕ разделы до конца — включая ЗАКЛЮЧЕНИЕ и СПИСОК_ЛИТЕРАТУРЫ\n"
        f"• Не обрывай текст на середине\n"
        f"• Начинай с ##ВВЕДЕНИЕ##:"
    )


def get_stages(lang="tk"):
    S = {
        "tk": [
            (5,  "🔍 Tema seljerilýär..."),
            (15, "📚 Çeşmeler gözlenýär..."),
            (30, "📖 Giriş ýazylýar..."),
            (50, "📝 Esasy bölümler..."),
            (70, "📝 Dowam edýär..."),
            (85, "🔚 Netije ýazylýar..."),
            (93, "📑 Çeşmeler sanawy..."),
            (97, "🔧 Jemlenýär..."),
        ],
        "ru": [
            (5,  "🔍 Анализ темы..."),
            (15, "📚 Поиск источников..."),
            (30, "📖 Написание введения..."),
            (50, "📝 Основные разделы..."),
            (70, "📝 Продолжаем..."),
            (85, "🔚 Написание заключения..."),
            (93, "📑 Список литературы..."),
            (97, "🔧 Финализация..."),
        ],
        "en": [
            (5,  "🔍 Analyzing topic..."),
            (15, "📚 Searching sources..."),
            (30, "📖 Writing introduction..."),
            (50, "📝 Main sections..."),
            (70, "📝 Continuing..."),
            (85, "🔚 Writing conclusion..."),
            (93, "📑 References list..."),
            (97, "🔧 Finalizing..."),
        ],
    }
    return S.get(lang, S["tk"])

STAGES = get_stages("tk")  # default

async def call_deepseek(d: dict, on_progress) -> str:
    prompt  = build_prompt(d)
    headers = {
        "Authorization": f"Bearer {DEEPSEEK_API_KEY}",
        "Content-Type":  "application/json",
        "Accept-Encoding": "identity",
        "Accept": "application/json",
    }
    req_items = d.get("req_items", [])
    extra_texts = []
    for item in req_items:
        if item["type"] == "text":
            extra_texts.append(item["content"])
        elif item["type"] == "image":
            extra_texts.append("[Ulanyjy surat iberdi — suradyň ýanyndaky tekst talabyna görä hereket et]")
    if extra_texts:
        user_content_final = "\n\n".join(extra_texts) + "\n\n" + prompt
    else:
        user_content_final = prompt
    _doc_lang = d.get("doc_lang","ru")
    if _doc_lang == "en":
        system_prompt = (
            "You are a professional academic author. "
            "Write ONLY in the English language. "
            "If there is a special requirements section at the beginning — follow them completely. "
            "Number all lists: 1. 2. 3. — no bullet points •/—. "
            "No preambles or explanations."
        )
    else:
        system_prompt = (
            "Ты профессиональный академический автор. "
            "Пиши ТОЛЬКО на русском языке. "
            "СТРОГО ОБЯЗАТЕЛЬНО: каждый раздел начинай ТОЧНО с этих маркеров на отдельной строке:\n"
            "##ВВЕДЕНИЕ##\n"
            "##ГЛАВА_1##\n"
            "##ГЛАВА_2##\n"  
            "##ГЛАВА_3##\n"
            "##ЗАКЛЮЧЕНИЕ##\n"
            "##СПИСОК_ЛИТЕРАТУРЫ##\n"
            "Маркеры писать ТОЧНО так, без изменений. "
            "СПИСОК ЛИТЕРАТУРЫ: только 1. 2. 3. — без ##ГЛАВА## внутри. "
            "Все перечисления: 1. 2. 3. — без маркеров •/—."
        )
    _max_tok = d.get("_max_tokens_calc", 8000)
    body = {
        "model":       DEEPSEEK_MODEL,
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user",   "content": user_content_final},
        ],
        "max_tokens":  _max_tok,
        "temperature": 0.7,
        "stream":      False,
    }
    result: dict = {}
    error:  dict = {}
    done:   dict = {"flag": False}

    async def fetch():
        max_retries = 5
        last_exc    = None
        for attempt in range(1, max_retries + 1):
            try:
                async with httpx.AsyncClient(
                    timeout=httpx.Timeout(connect=60.0, read=600.0, write=60.0, pool=30.0),
                    limits=httpx.Limits(max_keepalive_connections=5, max_connections=10),
                    http2=False,
                ) as cl:
                    r = await cl.post(DEEPSEEK_URL, headers=headers, json=body)
                    if r.status_code != 200:
                        raise RuntimeError(f"HTTP {r.status_code}: {r.text[:300]}")
                    resp = r.json()
                    text = resp["choices"][0]["message"]["content"]
                    if not text or not text.strip():
                        raise RuntimeError("DeepSeek boş jogap iberdi")
                    result["text"] = text.strip()
                    return
            except (httpx.ConnectError, httpx.ConnectTimeout,
                    httpx.ReadTimeout, httpx.RemoteProtocolError) as e:
                last_exc = e
                wait = min(attempt * 10, 60)
                log.warning(f"Bağlantı {attempt}/{max_retries}: {type(e).__name__} — {wait}s garaşylýar")
                await asyncio.sleep(wait)
            except RuntimeError as e:
                error["exc"] = e
                done["flag"] = True
                return
            except Exception as e:
                last_exc = e
                wait = min(attempt * 10, 60)
                log.warning(f"Ýalňyşlyk {attempt}/{max_retries}: {e} — {wait}s garaşylýar")
                await asyncio.sleep(wait)
        error["exc"] = last_exc
        done["flag"] = True

    async def ticker():
        stage_idx = 0
        elapsed   = 0
        _stages = get_stages(d.get("ui_lang","tk"))
        while not done["flag"]:
            await asyncio.sleep(6)
            elapsed += 6
            if stage_idx < len(_stages):
                pct, status = _stages[stage_idx]
                stage_idx += 1
            else:
                pct    = 97
                mins   = elapsed // 60
                secs_e = elapsed % 60
                _wait_lbl = {"tk":f"⏳ Taýarlanýar... {mins}:{secs_e:02d}","ru":f"⏳ Создаётся... {mins}:{secs_e:02d}","en":f"⏳ Creating... {mins}:{secs_e:02d}"}.get(d.get("ui_lang","tk"),"⏳")
                status = _wait_lbl
            try:
                await on_progress(pct, status)
            except Exception:
                pass

    done["flag"] = False
    ft = asyncio.create_task(fetch())
    tt = asyncio.create_task(ticker())
    try:
        await ft
    finally:
        done["flag"] = True
        tt.cancel()
        try:
            await tt
        except asyncio.CancelledError:
            pass
    if "exc" in error:
        raise error["exc"]
    await on_progress(100, "✅ Taýar!")
    return result["text"]


def parse_ai(raw: str, secs: int) -> dict:
    def _between(text, start, *ends):
        s = text.find(start)
        if s == -1: return ""
        s += len(start); best = len(text)
        for e in ends:
            p = text.find(e, s)
            if p != -1 and p < best: best = p
        return text[s:best].strip()

    intro_raw = _between(raw, "##ВВЕДЕНИЕ##", "##ГЛАВА_1##", "##ЗАКЛЮЧЕНИЕ##", "##СПИСОК_ЛИТЕРАТУРЫ##")
    chapters: list = []
    for i in range(1, secs + 1):
        nxt    = f"##ГЛАВА_{i+1}##" if i < secs else "##ЗАКЛЮЧЕНИЕ##"
        ch_raw = _between(raw, f"##ГЛАВА_{i}##", nxt, "##СПИСОК_ЛИТЕРАТУРЫ##")
        if not ch_raw: continue
        lines = [ln.strip() for ln in ch_raw.splitlines() if ln.strip()]
        if lines and re.match(r"^\d+\.", lines[0]):
            title = md_clean(lines[0]); body = lines[1:]
        else:
            title = f"{i}. Глава {i}"; body = lines
        chapters.append({"title": title, "lines": body})
    conc_raw = _between(raw, "##ЗАКЛЮЧЕНИЕ##", "##СПИСОК_ЛИТЕРАТУРЫ##")
    src_raw  = _between(raw, "##СПИСОК_ЛИТЕРАТУРЫ##")
    raw_srcs = [ln.strip() for ln in src_raw.splitlines() if ln.strip() and not ln.strip().startswith("##")]
    sources: list = []
    for ln in raw_srcs:
        ln = re.sub(r"^(\d+)\.\d+\.", r"\1.", ln)
        if re.match(r"^\d+\.", ln): sources.append(ln)
        elif sources: sources[-1] += " " + ln
    return dict(
        intro      = [ln.strip() for ln in intro_raw.splitlines() if ln.strip()],
        chapters   = chapters,
        conclusion = [ln.strip() for ln in conc_raw.splitlines() if ln.strip()],
        sources    = sources,
    )


def _sf(run, size_pt=14, bold=False, italic=False):
    run.font.name = "Times New Roman"; run.font.size = Pt(size_pt)
    run.bold = bold; run.italic = italic
    rpr = run._r.get_or_add_rPr()
    rf  = OxmlElement("w:rFonts")
    for a in ("w:ascii","w:hAnsi","w:cs","w:eastAsia"): rf.set(qn(a), "Times New Roman")
    old = rpr.find(qn("w:rFonts"))
    if old is not None: rpr.remove(old)
    rpr.insert(0, rf)
    for tag in ("w:sz","w:szCs"):
        el = rpr.find(qn(tag))
        if el is None: el = OxmlElement(tag); rpr.append(el)
        el.set(qn("w:val"), str(size_pt * 2))

def _spf(para, align, line, first_line=0, left=0, hanging=0, space_after=0):
    pPr = para._p.get_or_add_pPr()
    sp  = pPr.find(qn("w:spacing"))
    if sp is None: sp = OxmlElement("w:spacing"); pPr.append(sp)
    sp.set(qn("w:after"), str(space_after))
    sp.set(qn("w:line"), str(line))
    sp.set(qn("w:lineRule"), "auto")
    ind_old = pPr.find(qn("w:ind"))
    if ind_old is not None: pPr.remove(ind_old)
    if first_line or left or hanging:
        ind = OxmlElement("w:ind")
        if first_line: ind.set(qn("w:firstLine"), str(first_line))
        if left:       ind.set(qn("w:left"), str(left))
        if hanging:    ind.set(qn("w:hanging"), str(hanging))
        pPr.append(ind)
    jc = pPr.find(qn("w:jc"))
    if jc is None: jc = OxmlElement("w:jc"); pPr.append(jc)
    jc.set(qn("w:val"), align)

def _para(doc, text, *, bold=False, italic=False, center=False,
          size_pt=14, line=360, first_line=851, space_after=0):
    p = doc.add_paragraph(); r = p.add_run(text)
    _sf(r, size_pt, bold, italic)
    _spf(p, "center" if center else "both", line, first_line=first_line, space_after=space_after)

def _page_break(doc):
    p = doc.add_paragraph(); r = p.add_run()
    br = OxmlElement("w:br"); br.set(qn("w:type"), "page"); r._r.append(br)
    _spf(p, "both", 360)

def _init_heading(doc):
    try:    st = doc.styles["Heading 1"]
    except: st = doc.styles.add_style("Heading 1", 1)
    st.font.name = "Times New Roman"; st.font.size = Pt(14)
    st.font.bold = True; st.font.color.rgb = RGBColor(0, 0, 0)
    pf = st.paragraph_format
    pf.alignment = WD_ALIGN_PARAGRAPH.CENTER
    pf.first_line_indent = Cm(0)
    pf.space_before = Pt(0); pf.space_after = Pt(0)
    pf.line_spacing_rule = WD_LINE_SPACING.ONE_POINT_FIVE
    el = st.element; pPr = el.get_or_add_pPr()
    ol = pPr.find(qn("w:outlineLvl"))
    if ol is None: ol = OxmlElement("w:outlineLvl"); pPr.append(ol)
    ol.set(qn("w:val"), "0")

def _heading(doc, text: str, line_v: int):
    p = doc.add_paragraph(style="Heading 1")
    for run in p.runs: run.text = ""
    r = p.add_run(text); _sf(r, 14, True)
    r.font.color.rgb = RGBColor(0, 0, 0)
    p.paragraph_format.alignment         = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.first_line_indent = Cm(0)
    p.paragraph_format.space_before      = Pt(0)
    p.paragraph_format.space_after       = Pt(0)
    _spf(p, "center", line_v, first_line=0)

def _set_p_text(new_p, text: str):
    r_els = new_p.findall(f".//{qn('w:r')}")
    written = False
    for r_el in r_els:
        for t_el in r_el.findall(qn("w:t")):
            if not written:
                t_el.text = text
                t_el.set("{http://www.w3.org/XML/1998/namespace}space", "preserve")
                written = True
            else:
                t_el.text = ""

# ═══════════════════════════════════════════════════════════════════════
# ✅ TÄZE ŞABLONA GÖRÄ DÜZEDILEN _copy_template_title (30 paragraph)
# ═══════════════════════════════════════════════════════════════════════
# Şablon strukturasy (indeks 0..29):
#   [3]   «УНИВЕРСИТЕТ»                       → uniwersitet ady
#   [4]   'высшего профессионального...'      → boş edilýär
#   [8]   'РЕФЕРАТ по дисциплине...'          → görnüş + sapak + tema
#   [14]  'Студент 1 курса, группы X'         → kurs/grupp
#   [15]  'Ady Familiýa'                      → FIO
#   [18]  'Проверил:'                         → mugallym bar bolsa
#   [19]  'Mugallymyň ady'                    → mugallym ady
#   [29]  '2026г.'                            → ýyl (soňky)
def _copy_template_title(doc: Document, d: dict):
    svc_ru      = SVC_RU.get(d["service"], "Реферат")
    teacher     = d.get("teacher", "").strip()
    t_position  = d.get("teacher_position", "").strip()
    has_teacher = bool(teacher and teacher != "______________")
    subject     = d.get("subject", "").strip()
    university  = d.get("university", "").strip()

    try:
        tmpl_bytes = _get_template_bytes()
        tmpl = Document(io.BytesIO(tmpl_bytes))
    except Exception as e:
        log.error(f"❌ Şablon açylmady: {e}")
        raise RuntimeError(f"Şablon açylmady: {e}")

    total = len(tmpl.paragraphs)
    log.info(f"Şablonda {total} paragraph bar")

    body   = doc.element.body
    sectPr = body.find(qn("w:sectPr"))

    # Mugallym ýok bolsa — "Проверил:" (18) we mugallym ady (19) goýulmaýar
    skip = set()
    if not has_teacher:
        skip.update({18, 19})

    # Soňky paragraph (ýyl) — aýratyn ýerleşdireris, loop-da geçmesin
    year_idx = total - 1  # 29
    skip.add(year_idx)

    def _append(p_el):
        if sectPr is not None:
            body.insert(list(body).index(sectPr), p_el)
        else:
            body.append(p_el)

    for i, tp in enumerate(tmpl.paragraphs):
        if i in skip:
            continue
        new_p = copy.deepcopy(tp._element)

        if i == 3:
            # «УНИВЕРСИТЕТ» → hakyky uniwersitet ady
            _set_p_text(new_p, f"«{university}»")
        elif i == 4:
            # 'высшего профессионального образования' → boş
            _set_p_text(new_p, "")
        elif i == 8:
            # РЕФЕРАТ по дисциплине / на тему: «Тема работы»
            if subject:
                new_text = (f"{svc_ru} по дисциплине {subject}\n"
                            f"на тему: «{d['theme']}»")
            else:
                new_text = f"{svc_ru} на тему: «{d['theme']}»"
            _set_p_text(new_p, new_text)
        elif i == 14:
            # Студент 1 курса, группы X → täze kurs/grupp
            _set_p_text(new_p, f"Студент {d['course']} курса, группы {d['group']}")
        elif i == 15:
            # Ady Familiýa → hakyky FIO
            _set_p_text(new_p, d["fullname"])
        elif i == 18 and has_teacher:
            # Проверил: → wezipe bar bolsa goşulýar
            проверил = f"Проверил: {t_position}" if t_position else "Проверил:"
            _set_p_text(new_p, проверил)
        elif i == 19 and has_teacher:
            # Mugallymyň ady → hakyky ady
            _set_p_text(new_p, teacher)

        _append(new_p)

    # Soňky: ýyl paragraphy
    if 0 <= year_idx < total:
        p_year = copy.deepcopy(tmpl.paragraphs[year_idx]._element)
        _set_p_text(p_year, "2026г.")
        _append(p_year)


def _auto_toc(doc, line_v: int, doc_lang: str = "ru"):
    _contents_h = {"ru":"Содержание","en":"Contents"}.get(doc_lang,"Содержание")
    _para(doc, _contents_h, bold=True, center=True, size_pt=14,
          line=line_v, first_line=0, space_after=0)
    p = doc.add_paragraph()
    r1 = p.add_run(); fc = OxmlElement("w:fldChar")
    fc.set(qn("w:fldCharType"), "begin"); r1._r.append(fc)
    r2 = p.add_run(); ins = OxmlElement("w:instrText")
    ins.set(qn("xml:space"), "preserve")
    ins.text = ' TOC \\o "1-1" \\h \\z \\u '
    r2._r.append(ins)
    r3 = p.add_run(); fs = OxmlElement("w:fldChar")
    fs.set(qn("w:fldCharType"), "separate"); r3._r.append(fs)
    r4 = p.add_run(); fe = OxmlElement("w:fldChar")
    fe.set(qn("w:fldCharType"), "end"); r4._r.append(fe)
    settings = doc.settings.element
    uf = OxmlElement("w:updateFields")
    uf.set(qn("w:val"), "true")
    settings.append(uf)

def _add_page_numbers(doc: Document) -> None:
    from docx.opc.part import Part
    from docx.opc.packuri import PackURI
    FOOTER_XML = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<w:ftr xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"'
        '       xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
        '<w:p><w:pPr><w:jc w:val="center"/><w:spacing w:before="0" w:after="0"/></w:pPr>'
        '<w:r><w:rPr><w:rFonts w:ascii="Times New Roman" w:hAnsi="Times New Roman" w:cs="Times New Roman"/>'
        '<w:sz w:val="28"/><w:szCs w:val="28"/></w:rPr><w:fldChar w:fldCharType="begin"/></w:r>'
        '<w:r><w:rPr><w:rFonts w:ascii="Times New Roman" w:hAnsi="Times New Roman" w:cs="Times New Roman"/>'
        '<w:sz w:val="28"/><w:szCs w:val="28"/></w:rPr><w:instrText xml:space="preserve"> PAGE </w:instrText></w:r>'
        '<w:r><w:rPr><w:rFonts w:ascii="Times New Roman" w:hAnsi="Times New Roman" w:cs="Times New Roman"/>'
        '<w:sz w:val="28"/><w:szCs w:val="28"/></w:rPr><w:fldChar w:fldCharType="separate"/></w:r>'
        '<w:r><w:rPr><w:rFonts w:ascii="Times New Roman" w:hAnsi="Times New Roman" w:cs="Times New Roman"/>'
        '<w:sz w:val="28"/><w:szCs w:val="28"/></w:rPr><w:fldChar w:fldCharType="end"/></w:r>'
        '</w:p></w:ftr>'
    )
    FOOTER_EMPTY_XML = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<w:ftr xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
        '<w:p><w:pPr><w:jc w:val="center"/></w:pPr></w:p></w:ftr>'
    )
    FOOTER_CT  = 'application/vnd.openxmlformats-officedocument.wordprocessingml.footer+xml'
    FOOTER_REL = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/footer'
    footer_part = Part(PackURI('/word/footer_pnum.xml'), FOOTER_CT,
                       FOOTER_XML.encode('utf-8'), doc.part.package)
    rId = doc.part.relate_to(footer_part, FOOTER_REL)
    footer_first_part = Part(PackURI('/word/footer_first.xml'), FOOTER_CT,
                             FOOTER_EMPTY_XML.encode('utf-8'), doc.part.package)
    rId_first = doc.part.relate_to(footer_first_part, FOOTER_REL)
    body   = doc.element.body
    sectPr = body.find(qn('w:sectPr'))
    if sectPr is None:
        sectPr = OxmlElement('w:sectPr')
        body.append(sectPr)
    for old in sectPr.findall(qn('w:footerReference')):
        sectPr.remove(old)
    fr = OxmlElement('w:footerReference')
    fr.set(qn('w:type'), 'default')
    fr.set('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id', rId)
    sectPr.insert(0, fr)
    fr_first = OxmlElement('w:footerReference')
    fr_first.set(qn('w:type'), 'first')
    fr_first.set('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id', rId_first)
    sectPr.insert(0, fr_first)
    title_pg = sectPr.find(qn('w:titlePg'))
    if title_pg is None:
        title_pg = OxmlElement('w:titlePg')
        sectPr.append(title_pg)
    pg_num = sectPr.find(qn('w:pgNumType'))
    if pg_num is None:
        pg_num = OxmlElement('w:pgNumType')
        sectPr.append(pg_num)
    pg_num.set(qn('w:start'), '1')


def make_word(raw_text: str, d: dict) -> bytes:
    spc    = spc_float(d.get("spacing","default"))
    secs   = int(d["sections"])
    parsed = parse_ai(raw_text, secs)
    lv_    = lv(spc)
    doc = Document()
    sec = doc.sections[0]
    sec.page_width = Cm(21.001); sec.page_height = Cm(29.700)
    sec.left_margin = Cm(3.000); sec.right_margin = Cm(1.499)
    sec.top_margin  = Cm(2.000); sec.bottom_margin = Cm(2.000)
    _init_heading(doc)
    _copy_template_title(doc, d)
    _page_break(doc)
    _doc_lang_mw = d.get("doc_lang","ru")
    _intro_h   = {"ru":"Введение",         "en":"Introduction"}.get(_doc_lang_mw,"Введение")
    _concl_h   = {"ru":"ЗАКЛЮЧЕНИЕ",       "en":"CONCLUSION"}.get(_doc_lang_mw,"ЗАКЛЮЧЕНИЕ")
    _sources_h = {"ru":"Список литературы","en":"References"}.get(_doc_lang_mw,"Список литературы")
    _auto_toc(doc, lv_, _doc_lang_mw)
    _page_break(doc)
    _heading(doc, _intro_h, lv_)
    for ln in parsed["intro"]:
        t = md_clean(ln)
        if t: _para(doc, t, size_pt=14, line=lv_, first_line=851, space_after=0)
    for ch in parsed["chapters"]:
        _page_break(doc)
        _heading(doc, ch["title"], lv_)
        for ln in ch["lines"]:
            t = md_clean(ln)
            if t: _para(doc, t, size_pt=14, line=lv_, first_line=851, space_after=0)
    _page_break(doc)
    _heading(doc, _concl_h, lv_)
    for ln in parsed["conclusion"]:
        t = md_clean(ln)
        if t: _para(doc, t, size_pt=14, line=lv_, first_line=851, space_after=0)
    _page_break(doc)
    _heading(doc, _sources_h, lv_)
    for src in parsed["sources"]:
        t = md_clean(src)
        if not t or not re.match(r"^\d+\.", t): continue
        p = doc.add_paragraph(); r = p.add_run(t)
        _sf(r, 14, False); _spf(p, "both", lv_, first_line=851, space_after=0)
    _add_page_numbers(doc)
    buf = io.BytesIO(); doc.save(buf); buf.seek(0)
    return buf.getvalue()


def t_progress(d: dict, pct: int, status: str) -> str:
    bar = "█"*(pct//10) + "░"*(10-pct//10)
    lang = d.get("ui_lang","tk")
    def _e(t): return str(t).replace('&','&amp;').replace('<','&lt;').replace('>','&gt;')
    svc_label = {
        "tk": {"referat":"Referat 📄","doklad":"Doklad 🎤","pptx":"Prezentasiýa 📊"},
        "ru": {"referat":"Реферат 📄","doklad":"Доклад 🎤","pptx":"Презентация 📊"},
        "en": {"referat":"Essay 📄",  "doklad":"Report 🎤", "pptx":"Presentation 📊"},
    }.get(lang,{}).get(d.get("service","referat"), "Iş")
    prog_lbl = {"tk":"taýarlanýar","ru":"создаётся","en":"creating"}.get(lang,"taýarlanýar")
    return (f"⚙️ <b>{svc_label} {prog_lbl}...</b>\n\n"
            f"📝 <i>{_e(d.get('theme',''))}</i> \n👤 {_e(d.get('fullname',''))}\n\n"
            f"<code>[{bar}]</code> <b>{pct}%</b>\n<i>{_e(status)}</i>")

def t_summary(d: dict) -> str:
    mug = d.get('teacher','—')
    tp  = d.get('teacher_position','')
    mug_line = f"{tp} {mug}".strip() if tp else mug
    lang = d.get("ui_lang","tk")
    def _e(t): return str(t).replace('&','&amp;').replace('<','&lt;').replace('>','&gt;')
    svc_names = {
        "tk":{"referat":"Referat 📄","doklad":"Doklad 🎤"},
        "ru":{"referat":"Реферат 📄","doklad":"Доклад 🎤"},
        "en":{"referat":"Essay 📄",  "doklad":"Report 🎤"},
    }
    svc_nm = svc_names.get(lang,svc_names["tk"]).get(d.get("service","referat"),"?")
    labels = {
        "tk": ("📋 <b>Sargyt jemi:</b>","Görnüş","Sapak","Tema","Uni.","Talyp","Kurs","Mugallym","Bölüm","Aralyk","Sahypa","Çeşme"),
        "ru": ("📋 <b>Итого заказ:</b>","Тип","Предмет","Тема","Универ.","Студент","Курс","Преподаватель","Разделов","Интервал","Страниц","Источников"),
        "en": ("📋 <b>Order summary:</b>","Type","Subject","Topic","Uni.","Student","Year","Teacher","Sections","Spacing","Pages","Sources"),
    }
    lb = labels.get(lang, labels["tk"])
    return (f"{lb[0]}\n\n"
            f"📄 {lb[1]}  : <b>{svc_nm}</b>\n"
            f"📖 {lb[2]}  : {_e(d.get('subject','—'))}\n"
            f"📝 {lb[3]}   : <i>{_e(d.get('theme','—'))}</i>\n"
            f"🏫 {lb[4]}   : {_e(d.get('university','—'))}\n"
            f"👤 {lb[5]}  : {_e(d.get('fullname','—'))}\n"
            f"📚 {lb[6]}   : {_e(d.get('course','—'))}, {_e(d.get('group','—'))}\n"
            f"👨\u200d🏫 {lb[7]}: {_e(mug_line)}\n"
            f"📑 {lb[8]}  : {_e(d.get('sections','—'))}\n"
            f"📏 {lb[9]} : {_e(spc_str(d.get('spacing','default')))}\n"
            f"📄 {lb[10]} : {_e(d.get('pages','—'))}\n"
            f"🔗 {lb[11]}  : {_e(d.get('sources','—'))}\n")

async def deliver(uid: int, bot: Bot):
    if uid in CANCELLED_GENERATES:
        CANCELLED_GENERATES.discard(uid)
        PENDING.pop(uid, None)
        log.info(f"deliver: uid={uid} cancelled")
        return
    if uid not in PENDING:
        await bot.send_message(uid, "❌ Faýl tapylmady. Admin bilen habarlaşyň."); return
    info = PENDING.pop(uid)
    d    = info["data"]
    svc  = d.get("service", "referat")
    price = PRICE_STARS.get(svc, 159)
    svc_nm = {"referat":"Referat 📄","doklad":"Doklad 🎤","pptx":"Prezentasiýa 📊"}.get(svc, svc)

    # Ulanyjyny DB-ä goş (ýok bolsa)
    _db_get_or_create_user(uid)

    # BU ULANYJYNYŇ 1-NJI SARGYDYMY? → MUGT
    is_free = _db_is_first_free(uid)

    if is_free:
        # Mugt hyzmat — göni faýl iber
        theme  = re.sub(r"[^\w\s-]", "", d.get("theme", "work"))[:20].replace(" ", "_")
        fname  = f"{svc}_{theme}.docx"
        fb     = info["bytes"]
        # DB-ä ýaz
        _db_add_order(uid, svc, d.get("theme",""), 0, is_free=True, status="delivered")
        _db_mark_free_used(uid)
        lang_d = d.get("ui_lang","tk")
        free_caps = {
            "tk": (f"🎁 <b>Birinji sargydyňyz MUGT!</b>\n\n✅ <b>{svc_nm}</b> taýar\n"
                   f"📝 {d.get('theme','')}\n\n"
                   f"❓ Sorag: <code>{CONTACT_PHONE}</code>\n"
                   f"Täze sargyt: /start"),
            "ru": (f"🎁 <b>Первый заказ БЕСПЛАТНО!</b>\n\n✅ <b>{svc_nm}</b> готов\n"
                   f"📝 {d.get('theme','')}\n\n"
                   f"❓ Вопросы: <code>{CONTACT_PHONE}</code>\n"
                   f"Новый заказ: /start"),
            "en": (f"🎁 <b>First order FREE!</b>\n\n✅ <b>{svc_nm}</b> ready\n"
                   f"📝 {d.get('theme','')}\n\n"
                   f"❓ Questions: <code>{CONTACT_PHONE}</code>\n"
                   f"New order: /start"),
        }
        await bot.send_document(
            uid,
            BufferedInputFile(fb, filename=fname),
            caption=free_caps.get(lang_d, free_caps["tk"]),
            parse_mode="HTML",
            request_timeout=120
        )
        log.info(f"🎁 uid={uid} MUGT hyzmat aldy ({svc})")
        return

    # Galan sargytlar — Stars töleg
    stars = PRICE_STARS.get(svc, 159)
    oid   = _db_add_order(uid, svc, d.get("theme",""), stars, is_free=False, status="pending")
    PAYMENT_PENDING[uid] = {**info, "order_id": oid}
    lang  = d.get("ui_lang","tk")
    title_m = {"tk":svc_nm,"ru":SVC_RU.get(svc,svc_nm),"en":svc_nm}
    pay_hint = {
        "tk": f"💳 <b>{svc_nm} taýar! Eger faylda содержание gorunmeyan bolsa kompyuterde achsan chykar</b>\n\n📝 {d.get('theme','')[:60]}\n\nTölemek üçin aşakdaky düwmä basyň:",
        "ru": f"💳 <b>{svc_nm} готов! Если содержимое файла не отображается, оно появится при открытии файла на вашем компьютере.</b>\n\n📝 {d.get('theme','')[:60]}\n\nДля оплаты нажмите кнопку ниже:",
        "en": f"💳 <b>{svc_nm} ready! If the file contents are not visible, they will appear when you open them on your computer.</b>\n\n📝 {d.get('theme','')[:60]}\n\nPress the button below to pay:",
    }
    try:
        await bot.send_message(uid, pay_hint.get(lang, pay_hint["tk"]),
            parse_mode="HTML",
            reply_markup=kb_pay(lang))
        await bot.send_invoice(
            chat_id=uid,
            title=title_m.get(lang,svc_nm),
            description=d.get("theme","")[:100] or svc_nm,
            payload=f"{svc}_{uid}",
            currency="XTR",
            prices=[LabeledPrice(label=svc_nm, amount=stars)]
        )
    except Exception as e:
        log.error(f"Stars invoice: {e}")
        await bot.send_message(uid, f"❌ Stars töleg şowsuz: {e}")


router = Router()



@router.callback_query(St.lang, F.data.startswith("setlang:"))
async def h_setlang(cb: CallbackQuery, state: FSMContext):
    lang = cb.data.split(":")[1]
    await state.update_data(ui_lang=lang)
    uid = cb.from_user.id
    is_first  = _db_is_first_free(uid)
    if is_first:
        welcome = {
            "tk": f"🎓 <b>Akademik Işler Boty</b>\n\nSalam! Referat, Doklad we Prezentasiýa taýarlaýaryn.\n\n🎁 <b>Birinji sargydyňyz MUGT!</b>\n\n❓ Sorag: <code>{CONTACT_PHONE}</code>",
            "ru": f"🎓 <b>Академический Бот</b>\n\nПривет! Создаю рефераты, доклады и презентации.\n\n🎁 <b>Первый заказ БЕСПЛАТНО!</b>\n\n❓ Вопросы: <code>{CONTACT_PHONE}</code>",
            "en": f"🎓 <b>Academic Bot</b>\n\nHello! I create essays, reports and presentations.\n\n🎁 <b>First order FREE!</b>\n\n❓ Questions: <code>{CONTACT_PHONE}</code>",
        }
    else:
        welcome = {
            "tk": f"🎓 <b>Akademik Işler Boty</b>\n\n❓ Sorag: <code>{CONTACT_PHONE}</code>",
            "ru": f"🎓 <b>Академический Бот</b>\n\n❓ Вопросы: <code>{CONTACT_PHONE}</code>",
            "en": f"🎓 <b>Academic Bot</b>\n\n❓ Questions: <code>{CONTACT_PHONE}</code>",
        }
    svc_q = {"tk":"\n\n📌 Haýsy hyzmaty isleýärsiňiz?","ru":"\n\n📌 Какую услугу хотите?","en":"\n\n📌 Which service?"}
    await cb.message.edit_text(
        welcome.get(lang, welcome["tk"]) + svc_q.get(lang,"\n\n📌"),
        parse_mode="HTML", reply_markup=kb_svc(lang))
    await state.set_state(St.s01)
    await cb.answer()

@router.callback_query(F.data == "back:start")
async def h_back(cb: CallbackQuery, state: FSMContext):
    await state.clear()
    uid_b = cb.from_user.id
    PENDING.pop(uid_b, None)
    ACTIVE_GENERATES.discard(uid_b)
    CANCELLED_GENERATES.add(uid_b)
    await cb.message.edit_text(
        "🌐 <b>Dil saýlaň / Выберите язык / Select language:</b>",
        parse_mode="HTML", reply_markup=KB_LANG_SEL)
    await state.set_state(St.lang)
    await cb.answer()

# Dil saýlaw düwmeleri
KB_LANG_SEL = InlineKeyboardMarkup(inline_keyboard=[
    [InlineKeyboardButton(text="🇹🇲 Türkmençe", callback_data="setlang:tk")],
    [InlineKeyboardButton(text="🇷🇺 Русский",   callback_data="setlang:ru")],
    [InlineKeyboardButton(text="🇬🇧 English",   callback_data="setlang:en")],
])

@router.message(CommandStart())
async def h_start(msg: Message, bot: Bot, state: FSMContext):
    await state.clear()
    uid = msg.from_user.id
    PENDING.pop(uid, None)
    ACTIVE_GENERATES.discard(uid)
    CANCELLED_GENERATES.add(uid)
    try:
        _db_get_or_create_user(uid,
            username=(msg.from_user.username or ""),
            full_name=(msg.from_user.full_name or ""))
    except Exception as e:
        log.warning(f"DB register: {e}")

    if uid not in SEEN_USERS:
        SEEN_USERS.add(uid)
        if INTRO_VIDEO_URL:
            try:
                if INTRO_VIDEO_URL.startswith("http"):
                    await bot.send_message(uid,
                        f"🎓 <b>Akademik Işler Botyna hoş geldiňiz!</b>\n\n"
                        f"📹 {INTRO_VIDEO_URL}", parse_mode="HTML")
                else:
                    await bot.send_video(uid, INTRO_VIDEO_URL,
                        caption="🎓 <b>Akademik Işler Botyna hoş geldiňiz!</b>",
                        parse_mode="HTML")
            except Exception as e:
                log.warning(f"Wideo: {e}")

    await msg.answer(
        "🌐 <b>Dil saýlaň / Выберите язык / Select language:</b>",
        parse_mode="HTML", reply_markup=KB_LANG_SEL)
    await state.set_state(St.lang)

@router.callback_query(St.s01, F.data.startswith("svc:"))
async def h01(cb: CallbackQuery, state: FSMContext):
    svc = cb.data.split(":")[1]
    await state.update_data(service=svc)
    d    = await state.get_data()
    lang = d.get("ui_lang", "tk")
    if svc == "pptx":
        pptx_q = {"tk":"📌 <b>1/3:</b> Prezentasiýaňyzyň temasyny ýazyň\n\n<i>Mysal: Türkmenistanyň ykdysady ösüşi</i>",
                  "ru":"📌 <b>1/3:</b> Введите тему презентации\n\n<i>Пример: Экономика Туркменистана</i>",
                  "en":"📌 <b>1/3:</b> Enter the presentation topic\n\n<i>Example: Economy of Turkmenistan</i>"}
        await ask(cb, "✅ <b>Prezentasiýa 📊</b>\n\n" + pptx_q.get(lang, pptx_q["tk"]), kb_back(lang))
        await state.set_state(St.sp1)
    else:
        svc_name = SVC_TM[svc]
        req_q = {"tk":f"✅ <b>{svc_name}</b> saýlandy!\n\n📌 Ýörite talaplar barmy?\n\n• <b>Talaplы</b> — faýl ýa surat ugradyp bilersiňiz\n• <b>Talapsyz</b> — adaty GOST",
                 "ru":f"✅ <b>{svc_name}</b> выбран!\n\n📌 Есть особые требования?\n\n• <b>С требованиями</b> — можно прислать файл/фото\n• <b>Без требований</b> — стандарт ГОСТ",
                 "en":f"✅ <b>{svc_name}</b> selected!\n\n📌 Any special requirements?\n\n• <b>With requirements</b> — send files/photos\n• <b>Without</b> — standard GOST"}
        await ask(cb, req_q.get(lang, req_q["tk"]), kb_req(lang))
        await state.set_state(St.s02)
    await cb.answer()


@router.callback_query(St.s02, F.data == "req:yes")
async def h02_yes(cb: CallbackQuery, state: FSMContext):
    uid = cb.from_user.id; REQ_ITEMS[uid] = []
    await state.update_data(has_req=True, req_items=[])
    d    = await state.get_data(); lang = d.get("ui_lang","tk")
    q    = {"tk":"✅ <b>Talaplы görnüş.</b>\n\n📎 Islendik zat iberiň — tekst, surat, faýl.\n✅ Gutaransoň «Talaplarym taýar» düwmesine basyň:",
            "ru":"✅ <b>С требованиями.</b>\n\n📎 Отправьте текст, фото или файл.\n✅ После нажмите «Требования готовы»:",
            "en":"✅ <b>With requirements.</b>\n\n📎 Send text, photo or file.\n✅ Then press «Requirements ready»:"}
    await ask(cb, q.get(lang,q["tk"]), kb_req_done(lang))
    await state.set_state(St.s02b); await cb.answer()

@router.message(St.s02b)
async def h02b(msg: Message, bot: Bot, state: FSMContext):
    uid = msg.from_user.id
    if uid not in REQ_ITEMS:
        REQ_ITEMS[uid] = []
    added = []
    txt = (msg.caption or msg.text or "").strip()
    if txt:
        REQ_ITEMS[uid].append({"type": "text", "content": txt})
        added.append("tekst")
    if msg.photo:
        try:
            fobj  = await bot.get_file(msg.photo[-1].file_id)
            fbuf  = await bot.download_file(fobj.file_path)
            raw   = fbuf.read()
            b64   = base64.b64encode(raw).decode()
            REQ_ITEMS[uid].append({"type": "image", "mime": "image/jpeg", "content": b64})
            added.append("surat")
            await msg.answer(
                "🖼 <b>Surat kabul edildi!</b>\n\n"
                "⚠️ <i>Suradyň mazmunyny ýa-da talaplaryny tekst bilen hem ýazyň — "
                "şeýdip AI has gowy düşünýär.</i>\n\n"
                "<i>Mysal: «Şu suratdaky tablisany goş», «Şeýle görnüşde ýaz» we ş.m.</i>",
                parse_mode="HTML", reply_markup=kb_req_done(lang))
            return
        except Exception as e:
            log.warning(f"Surat okalmady: {e}")
    if msg.document:
        try:
            fobj  = await bot.get_file(msg.document.file_id)
            fbuf  = await bot.download_file(fobj.file_path)
            raw   = fbuf.read()
            mime  = msg.document.mime_type or "application/octet-stream"
            if "pdf" in mime:
                from pdfminer.high_level import extract_text as _ext
                import io as _io
                pdf_txt = _ext(_io.BytesIO(raw))
                if pdf_txt.strip():
                    REQ_ITEMS[uid].append({"type": "text",
                                           "content": f"[Faýldan: {msg.document.file_name}]\n{pdf_txt[:3000]}"})
                else:
                    REQ_ITEMS[uid].append({"type": "image", "mime": "image/jpeg",
                                           "content": base64.b64encode(raw).decode()})
            elif mime.startswith("text"):
                REQ_ITEMS[uid].append({"type": "text",
                                       "content": f"[Faýldan: {msg.document.file_name}]\n{raw.decode('utf-8','ignore')[:3000]}"})
            else:
                REQ_ITEMS[uid].append({"type": "text",
                                       "content": f"[{msg.document.file_name} faýly goşuldy]"})
            added.append("faýl")
        except Exception as e:
            log.warning(f"Faýl okalmady: {e}")
    count = len(REQ_ITEMS[uid])
    d2 = await state.get_data(); lang = d2.get("ui_lang","tk")
    if added:
        ok_msg = {"tk":f"✅ kabul edildi! (Jemi: {count} zat)\n\nÝene iberip bilersiňiz ýa-da:\n✅ «Talaplarym taýar» düwmesine basyň.",
                  "ru":f"✅ принято! (Всего: {count})\n\nМожно отправить ещё или:\n✅ Нажмите «Требования готовы».",
                  "en":f"✅ accepted! (Total: {count})\n\nYou can send more or:\n✅ Press «Requirements ready»."}
        await msg.answer(ok_msg.get(lang,ok_msg["tk"]), parse_mode="HTML", reply_markup=kb_req_done(lang))
    else:
        err_msg = {"tk":"⚠️ Hiç zat tapylmady, täzeden iberiň.","ru":"⚠️ Ничего не найдено, попробуйте ещё раз.","en":"⚠️ Nothing found, try again."}
        await msg.answer(err_msg.get(lang,err_msg["tk"]), reply_markup=kb_req_done(lang))

@router.callback_query(St.s02b, F.data == "req:done")
async def h02b_done(cb: CallbackQuery, state: FSMContext):
    uid = cb.from_user.id
    items = REQ_ITEMS.pop(uid, [])
    texts = [i["content"] for i in items if i["type"] == "text"]
    req_text = "\n".join(texts) if texts else ""
    await state.update_data(req_text=req_text, req_items=items)
    d_pre = await state.get_data()
    lang_p = d_pre.get("ui_lang","tk")
    uni_q = {"tk":"📌 <b>3/13:</b> Uniwersitetiňiziň doly adyny ýazyň\n<i>Mysal: Türkmenistanyň Oguz han adyndaky Inžener-Tehnologiýalar Uniwersiteti</i>",
              "ru":"📌 <b>3/13:</b> Полное название вашего университета\n<i>Пример: Туркменский государственный энергетический институт</i>",
              "en":"📌 <b>3/13:</b> Full name of your university\n<i>Example: Turkmen State Energy Institute</i>"}
    await ask(cb,
        f"✅ {'Talaplar kabul edildi!' if lang_p=='tk' else 'Требования приняты!' if lang_p=='ru' else 'Requirements accepted!'} ({len(items)} zat)\n\n"
        + uni_q.get(lang_p, uni_q["tk"]))
    await state.set_state(St.s03); await cb.answer()

@router.callback_query(St.s02, F.data == "req:no")
async def h02_no(cb: CallbackQuery, state: FSMContext):
    await state.update_data(has_req=False, req_text="")
    d_pre2 = await state.get_data()
    lang_p2 = d_pre2.get("ui_lang","tk")
    uni_q2 = {"tk":"✅ <b>Adaty GOST görnüşi.</b>\n\n📌 <b>3/13:</b> Uniwersitetiňiziň doly adyny ýazyň\n<i>Mysal: Пермский национальный исследовательский политехнический университет</i>",
               "ru":"✅ <b>Стандартный ГОСТ.</b>\n\n📌 <b>3/13:</b> Полное название университета\n<i>Пример: Пермский национальный исследовательский политехнический университет</i>",
               "en":"✅ <b>Standard GOST format.</b>\n\n📌 <b>3/13:</b> Full university name\n<i>Example: Пермский национальный исследовательский политехнический университет</i>"}
    await ask(cb, uni_q2.get(lang_p2, uni_q2["tk"]))
    await state.set_state(St.s03); await cb.answer()

@router.message(St.s03)
async def h03(msg: Message, state: FSMContext):
    d = await state.get_data(); lang = d.get("ui_lang","tk")
    if len(msg.text.strip()) < 4:
        errs = {"tk":"❌ Iň az 4 harp ýazyň.","ru":"❌ Минимум 4 символа.","en":"❌ At least 4 chars."}
        await msg.answer(errs.get(lang,"❌")); return
    await state.update_data(university=msg.text.strip())
    q = {"tk":"✅ Kabul edildi!\n\n📌 <b>4/13:</b> Sapakyňyzyň adyny ýazyň\n<i>Mysal: Ykdysadyýet nazaryýeti</i>",
         "ru":"✅ Принято!\n\n📌 <b>4/13:</b> Название предмета\n<i>Пример: Экономическая теория</i>",
         "en":"✅ Accepted!\n\n📌 <b>4/13:</b> Subject name\n<i>Example: Economic Theory</i>"}
    await msg.answer(q.get(lang,q["tk"]), parse_mode="HTML")
    await state.set_state(St.s04)

@router.message(St.s04)
async def h04(msg: Message, state: FSMContext):
    d = await state.get_data(); lang = d.get("ui_lang","tk")
    if len(msg.text.strip()) < 3:
        errs = {"tk":"❌ Iň az 3 harp.","ru":"❌ Минимум 3 символа.","en":"❌ At least 3 chars."}
        await msg.answer(errs.get(lang,"❌")); return
    await state.update_data(subject=msg.text.strip())
    q = {"tk":"✅ Kabul edildi!\n\n📌 <b>5/13:</b> Işiňiziň temasyny doly ýazyň\n<i>Mysal: Türkmenistanda ykdysady ösüş</i>",
         "ru":"✅ Принято!\n\n📌 <b>5/13:</b> Полная тема работы\n<i>Пример: Экономическое развитие Туркменистана</i>",
         "en":"✅ Accepted!\n\n📌 <b>5/13:</b> Full topic of the work\n<i>Example: Economic development of Turkmenistan</i>"}
    await msg.answer(q.get(lang,q["tk"]), parse_mode="HTML")
    await state.set_state(St.s05)

@router.message(St.s05)
async def h05(msg: Message, state: FSMContext):
    theme = msg.text.strip()
    d = await state.get_data(); lang = d.get("ui_lang","tk")
    if len(theme) < 8:
        errs = {"tk":"❌ Tema gysga! Iň az 8 harp.","ru":"❌ Тема слишком короткая! Минимум 8 символов.","en":"❌ Topic too short! At least 8 chars."}
        await msg.answer(errs.get(lang,"❌"), parse_mode="HTML"); return
    await state.update_data(theme=theme)
    q = {"tk":f"✅ Tema: {theme}\n\n📌 <b>6/13:</b> Adyňyzy we Familiýaňyzy ýazyň\n<i>Mysal: Myrat Mämmedow</i>",
         "ru":f"✅ Тема: {theme}\n\n📌 <b>6/13:</b> Ваше имя и фамилия\n<i>Пример: Иван Иванов</i>",
         "en":f"✅ Topic: {theme}\n\n📌 <b>6/13:</b> Your full name\n<i>Example: John Smith</i>"}
    await msg.answer(q.get(lang,q["tk"]), parse_mode="HTML")
    await state.set_state(St.s06)

@router.message(St.s06)
async def h06(msg: Message, state: FSMContext):
    d = await state.get_data(); lang = d.get("ui_lang","tk")
    if len(msg.text.strip()) < 3:
        errs = {"tk":"❌ Dogry ýazyň.","ru":"❌ Введите корректно.","en":"❌ Enter correctly."}
        await msg.answer(errs.get(lang,"❌")); return
    await state.update_data(fullname=msg.text.strip())
    q = {"tk":"✅ Kabul edildi!\n\n📌 <b>7/13:</b> Haýsy kursda okaýarsyňyz?",
         "ru":"✅ Принято!\n\n📌 <b>7/13:</b> На каком курсе вы учитесь?",
         "en":"✅ Accepted!\n\n📌 <b>7/13:</b> Which year are you in?"}
    await msg.answer(q.get(lang,q["tk"]), parse_mode="HTML", reply_markup=kb_crs(lang))
    await state.set_state(St.s07)

@router.callback_query(St.s07, F.data.startswith("crs:"))
async def h07(cb: CallbackQuery, state: FSMContext):
    n = cb.data.split(":")[1]; await state.update_data(course=n)
    d = await state.get_data(); lang = d.get("ui_lang","tk")
    q = {"tk":f"✅ {n}-nji kurs!\n\n📌 <b>8/13:</b> Toparыňyzyň adyny ýazyň\n<i>Mysal: EHM-22</i>",
         "ru":f"✅ {n}-й курс!\n\n📌 <b>8/13:</b> Название вашей группы\n<i>Пример: ЭВМ-22</i>",
         "en":f"✅ Year {n}!\n\n📌 <b>8/13:</b> Your group name\n<i>Example: CS-22</i>"}
    await ask(cb, q.get(lang,q["tk"]))
    await state.set_state(St.s08); await cb.answer()

@router.message(St.s08)
async def h08(msg: Message, state: FSMContext):
    await state.update_data(group=msg.text.strip())
    d = await state.get_data(); lang = d.get("ui_lang","tk")
    q = {"tk":"✅ Kabul edildi!\n\n📌 <b>9/13:</b> Mugallymyň adyny ýazyň\n<i>Mugallym ýok bolsa — «Geç» basyň</i>",
         "ru":"✅ Принято!\n\n📌 <b>9/13:</b> Имя преподавателя\n<i>Если нет — нажмите «Пропустить»</i>",
         "en":"✅ Accepted!\n\n📌 <b>9/13:</b> Teacher's name\n<i>If none — press «Skip»</i>"}
    await msg.answer(q.get(lang,q["tk"]), parse_mode="HTML", reply_markup=kb_skip(lang))
    await state.set_state(St.s09)

@router.message(St.s09)
async def h09_text(msg: Message, state: FSMContext):
    teacher = msg.text.strip()
    await state.update_data(teacher=teacher)
    d = await state.get_data(); lang = d.get("ui_lang","tk")
    q = {"tk":"✅ Mugallym kabul edildi!\n\n📌 <b>9б/13:</b> Mugallymyň wezipesini ýazyň\n<i>Mysal: доцент, профессор</i>",
         "ru":"✅ Преподаватель принят!\n\n📌 <b>9б/13:</b> Должность преподавателя\n<i>Пример: доцент, профессор</i>",
         "en":"✅ Teacher accepted!\n\n📌 <b>9б/13:</b> Teacher's position\n<i>Example: associate professor</i>"}
    await msg.answer(q.get(lang,q["tk"]), parse_mode="HTML")
    await state.set_state(St.s09b)

@router.message(St.s09b)
async def h09b(msg: Message, state: FSMContext):
    d = await state.get_data(); lang = d.get("ui_lang","tk")
    if len(msg.text.strip()) < 2:
        await msg.answer({"tk":"❌ Iň az 2 harp.","ru":"❌ Минимум 2 символа.","en":"❌ Min 2 chars."}.get(lang,"❌")); return
    await state.update_data(teacher_position=msg.text.strip())
    q = {"tk":"✅ Kabul edildi!\n\n📌 <b>10/13:</b> Näçe esasy bölüm bolmaly?",
         "ru":"✅ Принято!\n\n📌 <b>10/13:</b> Сколько основных разделов?",
         "en":"✅ Accepted!\n\n📌 <b>10/13:</b> How many main sections?"}
    await msg.answer(q.get(lang,q["tk"]), parse_mode="HTML", reply_markup=kb_sec(lang))
    await state.set_state(St.s10)

@router.callback_query(St.s09, F.data == "skip:teacher")
async def h09_skip(cb: CallbackQuery, state: FSMContext):
    await state.update_data(teacher="", teacher_position="")
    d = await state.get_data(); lang = d.get("ui_lang","tk")
    q = {"tk":"✅ Geçildi!\n\n📌 <b>10/13:</b> Näçe esasy bölüm bolmaly?",
         "ru":"✅ Пропущено!\n\n📌 <b>10/13:</b> Сколько основных разделов?",
         "en":"✅ Skipped!\n\n📌 <b>10/13:</b> How many main sections?"}
    await ask(cb, q.get(lang,q["tk"]), kb_sec(lang))
    await state.set_state(St.s10); await cb.answer()

@router.callback_query(St.s10, F.data.startswith("sec:"))
async def h10(cb: CallbackQuery, state: FSMContext):
    n = int(cb.data.split(":")[1]); await state.update_data(sections=n)
    d = await state.get_data(); lang = d.get("ui_lang","tk")
    q = {"tk":f"✅ {n} bölüm!\n\n📌 <b>11/13:</b> Setirler aralygy näçe?\n• <b>Adaty (1.5)</b> — GOST\n• <b>Üýtgetmek</b>",
         "ru":f"✅ {n} раздела!\n\n📌 <b>11/13:</b> Межстрочный интервал?\n• <b>По умолчанию (1.5)</b> — ГОСТ\n• <b>Изменить</b>",
         "en":f"✅ {n} sections!\n\n📌 <b>11/13:</b> Line spacing?\n• <b>Default (1.5)</b> — GOST\n• <b>Custom</b>"}
    await ask(cb, q.get(lang,q["tk"]), kb_spc(lang))
    await state.set_state(St.s11); await cb.answer()

@router.callback_query(St.s11, F.data == "spc:default")
async def h11_def(cb: CallbackQuery, state: FSMContext):
    await state.update_data(spacing="default")
    d_11 = await state.get_data(); lang_11 = d_11.get("ui_lang","tk")
    q11 = {"tk":"✅ Setirler aralygy <b>1.5</b> saýlandy!\n\n📌 <b>12/13:</b> Näçe sahypa?\n\n💡 <i>7—15 maslahat</i>\nSan ýazyň <i>(mysal: 12)</i>:",
           "ru":"✅ Интервал <b>1.5</b> выбран!\n\n📌 <b>12/13:</b> Сколько страниц?\n\n💡 <i>7—15 рекомендуется</i>\nВведите число <i>(пример: 12)</i>:",
           "en":"✅ Spacing <b>1.5</b> selected!\n\n📌 <b>12/13:</b> How many pages?\n\n💡 <i>7—15 recommended</i>\nEnter number <i>(example: 12)</i>:"}
    await ask(cb, q11.get(lang_11,q11["tk"]))
    await state.set_state(St.s12); await cb.answer()

@router.callback_query(St.s11, F.data == "spc:custom")
async def h11_cus(cb: CallbackQuery, state: FSMContext):
    d_11c = await state.get_data(); lang_11c = d_11c.get("ui_lang","tk")
    q11c = {"tk":"✏️ <b>Setirler aralygyny saýlaň:</b>\n\n• <b>1.0</b> — gysga\n• <b>1.25</b> — orta\n• <b>1.5</b> — GOST",
            "ru":"✏️ <b>Выберите межстрочный интервал:</b>\n\n• <b>1.0</b> — узкий\n• <b>1.25</b> — средний\n• <b>1.5</b> — ГОСТ",
            "en":"✏️ <b>Select line spacing:</b>\n\n• <b>1.0</b> — compact\n• <b>1.25</b> — medium\n• <b>1.5</b> — GOST"}
    await ask(cb, q11c.get(lang_11c,q11c["tk"]), KB_SPCV)
    await state.set_state(St.s11b); await cb.answer()

@router.callback_query(St.s11b, F.data.startswith("spv:"))
async def h11b(cb: CallbackQuery, state: FSMContext):
    val = cb.data.split(":")[1]; await state.update_data(spacing=val)
    d_11b = await state.get_data(); lang_11b = d_11b.get("ui_lang","tk")
    q11b = {"tk":f"✅ Setirler aralygy <b>{val}</b> saýlandy!\n\n📌 <b>12/13:</b> Näçe sahypa?\n\n💡 <i>7—17 maslahat</i>\nSan ýazyň <i>(mysal: 12)</i>:",
            "ru":f"✅ Интервал <b>{val}</b> выбран!\n\n📌 <b>12/13:</b> Сколько страниц?\n\n💡 <i>7—17 рекомендуется</i>\nВведите число <i>(пример: 12)</i>:",
            "en":f"✅ Spacing <b>{val}</b> selected!\n\n📌 <b>12/13:</b> How many pages?\n\n💡 <i>7—17 recommended</i>\nEnter number <i>(example: 12)</i>:"}
    await ask(cb, q11b.get(lang_11b,q11b["tk"]))
    await state.set_state(St.s12); await cb.answer()

@router.message(St.s12)
async def h12(msg: Message, state: FSMContext):
    d = await state.get_data(); lang = d.get("ui_lang","tk")
    try: n = int(msg.text.strip())
    except:
        await msg.answer({"tk":"❌ Diňe san ýazyň! Mysal: 12","ru":"❌ Введите число! Пример: 12","en":"❌ Enter a number! Example: 12"}.get(lang,"❌"), parse_mode="HTML"); return
    if not (7 <= n <= 15):
        await msg.answer({"tk":"⚠️ Diňe 7—15 arasynda!","ru":"⚠️ Только от 7 до 15!","en":"⚠️ Only 7—15!"}.get(lang,"⚠️"), parse_mode="HTML"); return
    await state.update_data(pages=n)
    q = {"tk":f"✅ {n} sahypa!\n\n📌 <b>13/13:</b> Çeşme sany? (8-20)",
         "ru":f"✅ {n} страниц!\n\n📌 <b>13/13:</b> Количество источников? (8-20)",
         "en":f"✅ {n} pages!\n\n📌 <b>13/13:</b> Number of sources? (8-20)"}
    await msg.answer(q.get(lang,q["tk"]), parse_mode="HTML", reply_markup=kb_src())
    await state.set_state(St.s13)

async def _run_generate(cb: CallbackQuery, state: FSMContext, d: dict = None):
    if d is None:
        d = await state.get_data()
    required = ["service","university","subject","theme","fullname","course","group","sections","pages","sources"]
    missing = [k for k in required if not d.get(k)]
    if missing:
        await cb.answer("❌ Ýetmeýän maglumat!", show_alert=True)
        await cb.message.answer(
            f"❌ Ýetmeýänler: <code>{chr(44).join(missing)}</code>\nTäzeden başlamak: /start",
            parse_mode="HTML")
        await state.clear(); return
    uid_gen = cb.from_user.id
    if uid_gen in ACTIVE_GENERATES:
        await cb.answer("⏳ Sargydyňyz eýýäm taýarlanýar, sabyr ediň!", show_alert=True); return
    ACTIVE_GENERATES.add(uid_gen)
    CANCELLED_GENERATES.discard(uid_gen)
    _lang_gs = d.get("ui_lang","tk")
    _start_lbl = {"tk":"Başlanýar...","ru":"Начинаем...","en":"Starting..."}.get(_lang_gs,"Başlanýar...")
    prog = await cb.message.edit_text(t_progress(d, 0, _start_lbl), parse_mode="HTML")
    await state.set_state(St.s13); await cb.answer()
    bot = cb.bot; cid = cb.message.chat.id; mid = prog.message_id
    async def pcb(pct, status):
        try: await bot.edit_message_text(t_progress(d, pct, status), chat_id=cid, message_id=mid, parse_mode="HTML")
        except: pass
    try:
        raw       = await call_deepseek(d, pcb)
        doc_bytes = await asyncio.get_running_loop().run_in_executor(None, make_word, raw, d)
        PENDING[cb.from_user.id] = {"bytes": doc_bytes, "data": d}
        _lang_r = d.get("ui_lang","tk")
        _ready_hdr = {"tk":"✅ <b>Işiňiz taýar boldy!</b>","ru":"✅ <b>Работа готова!</b>","en":"✅ <b>Work is ready!</b>"}.get(_lang_r,"✅ <b>Taýar!</b>")
        await bot.edit_message_text(f"{_ready_hdr}\n\n{t_summary(d)}",
                                     chat_id=cid, message_id=mid, parse_mode="HTML")
        await deliver(cb.from_user.id, bot)
    except Exception as exc:
        import traceback
        tb = traceback.format_exc()
        log.error(f"Generate: {exc}\n{tb}")
        _lang_err = d.get("ui_lang","tk")
        _err_lbl = {"tk":"Ýalňyşlyk","ru":"Ошибка","en":"Error"}.get(_lang_err,"Ýalňyşlyk")
        _restart = {"tk":"Täzeden başlamak: /start","ru":"Начать заново: /start","en":"Start again: /start"}.get(_lang_err,"")
        await bot.edit_message_text(
            f"❌ <b>{_err_lbl}!</b>\n\n<code>{str(exc)[:300]}</code>\n\n{_restart}",
            chat_id=cid, message_id=mid, parse_mode="HTML")
    finally:
        ACTIVE_GENERATES.discard(uid_gen)
        await state.clear()

@router.callback_query(St.s13, F.data.startswith("src:"))
async def h13_src(cb: CallbackQuery, state: FSMContext):
    n = int(cb.data.split(":")[1]); await state.update_data(sources=n)
    d    = await state.get_data()
    lang = d.get("ui_lang","tk")
    # Çeşme sany tassykla we göni generate başla (diňe rusça)
    await state.update_data(doc_lang="ru")
    d = await state.get_data()
    await _run_generate(cb, state, d)

@router.callback_query(St.s13, F.data.startswith("reflang:"))
async def h13_generate(cb: CallbackQuery, state: FSMContext):
    doc_lang = cb.data.split(":")[1]
    await state.update_data(doc_lang=doc_lang)
    d = await state.get_data()
    await _run_generate(cb, state, d)


@router.message(F.video | F.video_note)
async def h_video(msg: Message):
    uid = msg.from_user.id
    if uid in ADMIN_IDS:
        fid = msg.video.file_id if msg.video else msg.video_note.file_id
        await msg.answer(
            f"✅ <b>Wideo file_id:</b>\n\n<code>{fid}</code>\n\n"
            f"Şony <code>INTRO_VIDEO_URL</code>-e ýaz.",
            parse_mode="HTML")
        return
    if uid not in PAYMENT_PENDING:
        return
    await msg.answer(
        "❌ <b>Wideo kabul edilmeýär!</b>\n\n"
        "Çek hökmünde diňe şulary iberiň:\n"
        "📸 Surat (screenshot)\n"
        "📄 Faýl (PDF, PNG, JPG)\n\n"
        "Töleg screenshotyny ýa-da çek suratyny iberiň:",
        parse_mode="HTML")

@router.message(F.text.regexp(r"^/send\s+\d+"))
async def a_send(msg: Message, bot: Bot):
    if msg.from_user.id not in ADMIN_IDS: return
    uid = int(msg.text.split()[1]); await deliver(uid, bot)
    await msg.answer(f"✅ {uid} iberildi.")

@router.message(F.text == "/orders")
async def a_orders(msg: Message):
    if msg.from_user.id not in ADMIN_IDS: return
    if not PENDING: await msg.answer("📭 Garaşyp duran sargyt ýok."); return
    lines = ["📋 <b>Sargytlar:</b>\n"]
    for uid, info in PENDING.items():
        d = info["data"]
        lines.append(f"👤 <code>{uid}</code> | <i>{d.get('theme','?')[:25]}</i> | {d.get('pages','?')} sah.\n   /send {uid}\n")
    await msg.answer("\n".join(lines), parse_mode="HTML")

@router.message(F.text == "/admin")
async def a_help(msg: Message):
    if msg.from_user.id not in ADMIN_IDS: return
    await msg.answer(
        "🔧 <b>Admin:</b>\n"
        "<code>/orders</code> — garaşyp duran sargytlar\n"
        "<code>/send &lt;id&gt;</code> — faýly iber\n"
        "<code>/stats</code> — umumy statistika\n"
        "<code>/users</code> — /start basan ulanyjylar sanawy\n"
        "<code>/getfileid</code> — wideo iberip file_id al",
        parse_mode="HTML")

@router.message(F.text == "/users")
async def a_users(msg: Message):
    if msg.from_user.id not in ADMIN_IDS: return
    try:
        conn = sqlite3.connect(DB_PATH)
        conn.row_factory = sqlite3.Row
        cur = conn.cursor()
        cur.execute("""
            SELECT user_id, username, full_name, free_used, total_orders,
                   datetime(created_at, 'unixepoch') as reg_date
            FROM users
            ORDER BY created_at DESC
            LIMIT 50
        """)
        rows = cur.fetchall()
        conn.close()

        if not rows:
            await msg.answer("📭 Entek ulanyjy ýok."); return

        lines = [f"👥 <b>Ulanyjylar ({len(rows)} soňky):</b>\n"]
        for r in rows:
            uname = f"@{r['username']}" if r['username'] else "—"
            fname = r['full_name'] or "—"
            free  = "✅ mugt aldy" if r['free_used'] else "🎁 mugt bar"
            lines.append(
                f"👤 <code>{r['user_id']}</code> | {uname} | {fname}\n"
                f"   {free} | Sargyt: {r['total_orders']} | {r['reg_date'][:10]}\n"
            )

        # Telegram 4096 simwol çägi — uly bolsa böl
        text = "\n".join(lines)
        if len(text) <= 4096:
            await msg.answer(text, parse_mode="HTML")
        else:
            # Bölek-bölek iber
            chunk = ""
            for line in lines:
                if len(chunk) + len(line) > 3800:
                    await msg.answer(chunk, parse_mode="HTML")
                    chunk = line
                else:
                    chunk += line
            if chunk:
                await msg.answer(chunk, parse_mode="HTML")
    except Exception as e:
        await msg.answer(f"❌ DB ýalňyşlyk: {e}")

@router.message(F.text == "/stats")
async def a_stats(msg: Message):
    if msg.from_user.id not in ADMIN_IDS: return
    try:
        s = _db_stats()
        await msg.answer(
            "📊 <b>STATISTIKA</b>\n\n"
            f"👥 Ulanyjylar (jemi)    : <b>{s['users']}</b>\n"
            f"📋 Sargytlar (jemi)     : <b>{s['orders']}</b>\n"
            f"🎁 Mugt sargytlar      : <b>{s['free']}</b>\n"
            f"💰 Tölegli (taýar)     : <b>{s['paid']}</b>\n"
            f"💵 Umumy girdeji       : <b>{s['revenue']} ₽</b>",
            parse_mode="HTML"
        )
    except Exception as e:
        await msg.answer(f"❌ DB ýalňyşlyk: {e}")

@router.message(F.text == "/getfileid")
async def a_getfileid(msg: Message):
    if msg.from_user.id not in ADMIN_IDS: return
    await msg.answer("📹 Indi wideoňy şu söhbete iber — file_id-ni bererin.")




# ════════════════════════════════════════════════════════════════════
# PREZENTASIÝA BÖLÜMI
# ════════════════════════════════════════════════════════════════════

KB_PPTX_LANG = InlineKeyboardMarkup(inline_keyboard=[
    [InlineKeyboardButton(text="🇷🇺 Rusça (Русский)",   callback_data="plang:ru")],
    [InlineKeyboardButton(text="🇬🇧 İngilizce (English)", callback_data="plang:en")],
    [InlineKeyboardButton(text="🇹🇷 Türkçe",            callback_data="plang:tr")],
    [InlineKeyboardButton(text="🔙  Yza — başa gaýt",   callback_data="back:start")],
])
KB_PPTX_SLIDES = InlineKeyboardMarkup(inline_keyboard=[
    [InlineKeyboardButton(text="7",callback_data="sl:7"),
     InlineKeyboardButton(text="10",callback_data="sl:10"),
     InlineKeyboardButton(text="12",callback_data="sl:12")],
    [InlineKeyboardButton(text="🔙  Yza — başa gaýt",callback_data="back:start")],
])
def kb_pay(lang="tk"):
    labels = {
        "tk": "⭐ Telegram Stars bilen töle",
        "ru": "⭐ Оплатить Telegram Stars",
        "en": "⭐ Pay with Telegram Stars",
    }
    premium_labels = {
        "tk": "👑 Premium Bot",
        "ru": "👑 Premium Bot",
        "en": "👑 Premium Bot",
    }
    stars_bot_labels = {
        "tk": "⭐ Stars satyn al",
        "ru": "⭐ Купить Stars",
        "en": "⭐ Buy Stars",
    }
    channel_labels = {
        "tk": "📢 IT_turkmen kanalymyz",
        "ru": "📢 Наш канал IT_turkmen",
        "en": "📢 Our channel IT_turkmen",
    }
    return InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text=labels.get(lang,labels["tk"]), callback_data="pay:stars")],
        [InlineKeyboardButton(text=premium_labels.get(lang,"👑 Premium Bot"), url="https://t.me/PremiumBot"),
         InlineKeyboardButton(text=stars_bot_labels.get(lang,"⭐ Buy Stars"), url="https://t.me/TheDurovStars_bot")],
        [InlineKeyboardButton(text=channel_labels.get(lang,"📢 IT_turkmen"), url="https://t.me/IT_turkmen")],
    ])

KB_PAY = kb_pay()  # default

import random as _random

# Her prezentasiýada başga reňk palitrasы
_PALETTES = [
    # 1. Lacivert-Gyzyl (koyu lacivert + kızıl)
    {"dark":PRGBColor(0x1A,0x1A,0x2E),"accent":PRGBColor(0xE9,0x4F,0x37),
     "acc2":PRGBColor(0x16,0x21,0x3E),
     "bg_light":PRGBColor(0xF8,0xF9,0xFA),"white":PRGBColor(0xFF,0xFF,0xFF),
     "dark_t":PRGBColor(0x1A,0x1A,0x2E),"light_t":PRGBColor(0xCC,0xCC,0xDD)},
    # 2. Gök-Altyn (koyu mavi + altın sarısı)
    {"dark":PRGBColor(0x04,0x2A,0x4A),"accent":PRGBColor(0xF4,0xB9,0x42),
     "acc2":PRGBColor(0x06,0x3A,0x5F),
     "bg_light":PRGBColor(0xEF,0xF5,0xFB),"white":PRGBColor(0xFF,0xFF,0xFF),
     "dark_t":PRGBColor(0x04,0x2A,0x4A),"light_t":PRGBColor(0xCC,0xDF,0xFF)},
    # 3. Ýaşyl-Limon (koyu yeşil + limon)
    {"dark":PRGBColor(0x14,0x3D,0x23),"accent":PRGBColor(0x7C,0xC1,0x27),
     "acc2":PRGBColor(0x1A,0x52,0x2E),
     "bg_light":PRGBColor(0xF1,0xFA,0xF2),"white":PRGBColor(0xFF,0xFF,0xFF),
     "dark_t":PRGBColor(0x14,0x3D,0x23),"light_t":PRGBColor(0xBE,0xE6,0xC8)},
    # 4. Benewşe-Apelsin (koyu mor + turuncu)
    {"dark":PRGBColor(0x27,0x10,0x4A),"accent":PRGBColor(0xFF,0x7A,0x00),
     "acc2":PRGBColor(0x37,0x18,0x62),
     "bg_light":PRGBColor(0xFA,0xF3,0xFF),"white":PRGBColor(0xFF,0xFF,0xFF),
     "dark_t":PRGBColor(0x27,0x10,0x4A),"light_t":PRGBColor(0xDF,0xC8,0xFF)},
    # 5. Gara-Sian (gece mavisi + açık mavi)
    {"dark":PRGBColor(0x0A,0x14,0x28),"accent":PRGBColor(0x00,0xC2,0xE0),
     "acc2":PRGBColor(0x0D,0x1F,0x3C),
     "bg_light":PRGBColor(0xEE,0xFA,0xFF),"white":PRGBColor(0xFF,0xFF,0xFF),
     "dark_t":PRGBColor(0x0A,0x14,0x28),"light_t":PRGBColor(0xAA,0xEE,0xFF)},
    # 6. Goňur-Krem (koyu kahve + krem)
    {"dark":PRGBColor(0x2C,0x1A,0x0E),"accent":PRGBColor(0xE0,0xA0,0x30),
     "acc2":PRGBColor(0x3E,0x25,0x14),
     "bg_light":PRGBColor(0xFD,0xF8,0xEE),"white":PRGBColor(0xFF,0xFF,0xFF),
     "dark_t":PRGBColor(0x2C,0x1A,0x0E),"light_t":PRGBColor(0xF0,0xDC,0xBB)},
    # 7. Gyzyl-Ak (nar kırmızısı + beyaz)
    {"dark":PRGBColor(0x7B,0x0E,0x0E),"accent":PRGBColor(0xFF,0xC2,0xC2),
     "acc2":PRGBColor(0x9B,0x12,0x12),
     "bg_light":PRGBColor(0xFD,0xEF,0xEF),"white":PRGBColor(0xFF,0xFF,0xFF),
     "dark_t":PRGBColor(0x5A,0x08,0x08),"light_t":PRGBColor(0xFF,0xDD,0xDD)},
    # 8. Çal-Ýaşyl (antrasit + neon yeşil)
    {"dark":PRGBColor(0x1C,0x1C,0x2B),"accent":PRGBColor(0x39,0xD3,0x53),
     "acc2":PRGBColor(0x26,0x26,0x3A),
     "bg_light":PRGBColor(0xF2,0xFF,0xF4),"white":PRGBColor(0xFF,0xFF,0xFF),
     "dark_t":PRGBColor(0x1C,0x1C,0x2B),"light_t":PRGBColor(0xBB,0xEE,0xCC)},
]

def _get_palette():
    """Her çagyrylanda başga reňk palitrasы gaýtarýar"""
    import random
    return random.choice(_PALETTES)
_IMG_POS = ["left","right","bottom","left","right","bottom","left","right","bottom","bottom"]

def _pbg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def _prect(slide, l, t, w, h, color):
    sh = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()
    return sh

def _ptx(slide, l, t, w, h):
    return slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))

def _prun(tf, text, size, bold=False, color=None, italic=False, first=False, sp=0, align=PP_ALIGN.LEFT):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    r = p.add_run(); r.text = text
    r.font.name = "Times New Roman"; r.font.size = PPtx(size)
    r.font.bold = bold; r.font.italic = italic
    if color: r.font.color.rgb = color
    p.space_before = PPtx(sp); p.alignment = align
    return p

def _pimg(slide, img, l, t, w, h):
    if not img:
        return False
    try:
        from PIL import Image as _PIL_Image
        import io as _io2
        # Suraty JPEG hökmünde convert et (python-pptx üçin)
        pil_img = _PIL_Image.open(_io2.BytesIO(img))
        if pil_img.mode in ("RGBA", "LA", "P"):
            pil_img = pil_img.convert("RGB")
        out_buf = _io2.BytesIO()
        pil_img.save(out_buf, format="JPEG", quality=85)
        out_buf.seek(0)
        slide.shapes.add_picture(out_buf, Inches(l), Inches(t), Inches(w), Inches(h))
        return True
    except Exception as _e:
        import logging as _log2
        _log2.getLogger(__name__).warning(f"_pimg şowsuz: {_e}")
        return False

def _pdot(slide, l, t, color=None):
    sh = slide.shapes.add_shape(9, Inches(l), Inches(t), Inches(0.22), Inches(0.22))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color if color else PRGBColor(0xE9,0x4F,0x37)
    sh.line.fill.background()

def _set_transparency(shape, alpha: float):
    """Şekilin aýdyňlygyny XML arkaly sazlaýar. alpha: 0.0=görünmeýär, 1.0=doly aýdyň"""
    from pptx.oxml.ns import qn as _qn
    from lxml import etree as _etree
    el = shape._element
    # spPr dürli usullar bilen tap
    spPr = None
    for attr in ['spPr', 'get_or_add_spPr']:
        try:
            if hasattr(el, attr):
                spPr = getattr(el, attr)
                if callable(spPr): spPr = spPr()
                break
        except: pass
    if spPr is None:
        spPr = el.find(_qn('p:spPr'))
    if spPr is None:
        return
    solidFill = spPr.find(_qn('a:solidFill'))
    if solidFill is None:
        return
    alpha_val = int(alpha * 100000)
    for clr_tag in ['a:srgbClr','a:schemeClr','a:sysClr','a:prstClr']:
        clr = solidFill.find(_qn(clr_tag))
        if clr is not None:
            for old_a in clr.findall(_qn('a:alpha')):
                clr.remove(old_a)
            a_el = _etree.SubElement(clr, _qn('a:alpha'))
            a_el.set('val', str(alpha_val))
            break

def _add_chart_pptx(slide, chart_data: dict, l, t, w, h):
    try:
        from pptx.chart.data import ChartData
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.enum.chart import XL_LEGEND_POSITION
        from pptx.oxml.ns import qn as _cqn
        from lxml import etree as _etree

        labels  = chart_data.get("labels", ["A","B","C","D"])
        values  = chart_data.get("values", [40,35,25,20])
        title   = chart_data.get("title",  "")
        ctype   = chart_data.get("type",   "bar")
        caption = chart_data.get("caption","")
        y_label = chart_data.get("y_label","")
        x_label = chart_data.get("x_label","")

        # Başlyk - grafikiň ýokarynda, galyň
        _title_h = 0.38
        _cap_h   = 0.52
        _chart_t = t + _title_h
        _chart_h = h - _title_h - _cap_h

        if title:
            tx_t = slide.shapes.add_textbox(
                Inches(l), Inches(t), Inches(w), Inches(_title_h))
            tf_t = tx_t.text_frame; tf_t.word_wrap = True
            p_t = tf_t.paragraphs[0]; p_t.alignment = PP_ALIGN.CENTER
            r_t = p_t.add_run(); r_t.text = title
            r_t.font.name = "Times New Roman"
            r_t.font.size = PPtx(13); r_t.font.bold = True

        # Grafik
        cd = ChartData()
        cd.categories = labels
        cd.add_series("", values)
        ct_map = {
            "bar":  XL_CHART_TYPE.COLUMN_CLUSTERED,  # dik sütün (excel bar)
            "pie":  XL_CHART_TYPE.PIE,
            "line": XL_CHART_TYPE.LINE_MARKERS,
        }
        chart_shape = slide.shapes.add_chart(
            ct_map.get(ctype, XL_CHART_TYPE.COLUMN_CLUSTERED),
            Inches(l), Inches(_chart_t), Inches(w), Inches(_chart_h), cd)
        chart_obj = chart_shape.chart

        # Grafik içindäki başlyk öç
        chart_obj.has_title = False

        # ── Data labels (sütündäki sanlar) ──
        try:
            plot = chart_obj.plots[0]
            plot.has_data_labels = True
            dls = plot.data_labels
            dls.font.size = PPtx(10)
            dls.font.bold = True
            if ctype == "pie":
                dls.show_percentage   = True
                dls.show_category_name= True
                dls.show_value        = False
                dls.number_format     = "0%"
            else:
                dls.show_value        = True
                dls.show_category_name= False
                dls.number_format     = "0"
        except: pass

        # ── Value axis (Y oku) ──
        try:
            if ctype in ("bar","line"):
                v_ax = chart_obj.value_axis
                v_ax.tick_labels.font.size = PPtx(9)
                v_ax.tick_labels.font.bold = False
                if y_label:
                    v_ax.has_title = True
                    v_ax.axis_title.text_frame.text = y_label
                    for para in v_ax.axis_title.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.size = PPtx(9)
                            run.font.bold = False
        except: pass

        # ── Category axis (X oku) ──
        try:
            if ctype in ("bar","line"):
                c_ax = chart_obj.category_axis
                c_ax.tick_labels.font.size = PPtx(9)
                c_ax.tick_labels.font.bold = False
                if x_label:
                    c_ax.has_title = True
                    c_ax.axis_title.text_frame.text = x_label
                    for para in c_ax.axis_title.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.size = PPtx(9)
        except: pass

        # ── Legend ──
        try:
            chart_obj.has_legend = True
            chart_obj.legend.include_in_layout = False
            if ctype == "pie":
                chart_obj.legend.position = XL_LEGEND_POSITION.RIGHT
            else:
                chart_obj.legend.position = XL_LEGEND_POSITION.BOTTOM
            for para in chart_obj.legend.text_frame.paragraphs if hasattr(chart_obj.legend,"text_frame") else []:
                for run in para.runs:
                    run.font.size = PPtx(9)
        except: pass

        # ── Caption (asagynda düşündiriş) ──
        _cap_y = _chart_t + _chart_h + 0.04
        if caption:
            tx_c = slide.shapes.add_textbox(
                Inches(l), Inches(_cap_y), Inches(w), Inches(_cap_h))
            tf_c = tx_c.text_frame; tf_c.word_wrap = True
            p_c = tf_c.paragraphs[0]; p_c.alignment = PP_ALIGN.CENTER
            r_c = p_c.add_run(); r_c.text = caption
            r_c.font.name = "Times New Roman"
            r_c.font.size = PPtx(10); r_c.font.italic = True

        return True
    except Exception as e:
        log.warning(f"Grafik goşulmady: {e}"); return False


def build_pptx(slides_data: list, theme: str, images: list, student_info: dict = None) -> bytes:
    PC = _get_palette()   # Her gezek täze reňk!
    prs = Presentation()
    prs.slide_width = Inches(13.33); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]; total = len(slides_data)

    for idx, sld in enumerate(slides_data):
        slide   = prs.slides.add_slide(blank)
        ttl     = sld.get("title", "")
        points  = sld.get("points", [])
        img     = images[idx] if idx < len(images) else None
        has_chart = sld.get("has_chart", False)
        chart_dt  = sld.get("chart_data", {})
        is_first = (idx == 0); is_last = (idx == total - 1)

        if is_first:
            _pbg(slide, PC["dark"])
            # Doly fon suraty — bütin slaýd
            if img:
                _pimg(slide, img, 0, 0, 13.33, 7.5)
            # Garaňky örtük doly slaýd (tekst okalmagy üçin)
            ov = _prect(slide, 0, 0, 13.33, 7.5, PRGBColor(0,0,0))
            _set_transparency(ov, 0.45)  # 45% aýdyň
            # Sol tarap reňkli zolak
            _prect(slide, 0, 0, 0.35, 7.5, PC["accent"])
            # Tema ady - merkez
            tx = _ptx(slide, 0.7, 1.5, 12.0, 2.8)
            tf = tx.text_frame; tf.word_wrap = True
            r0 = tf.paragraphs[0].add_run(); r0.text = ttl
            r0.font.name = "Times New Roman"; r0.font.size = PPtx(38)
            r0.font.bold = True; r0.font.color.rgb = PC["white"]
            tf.paragraphs[0].alignment = PP_ALIGN.LEFT
            # Aşakdaky çyzygy
            _prect(slide, 0.7, 4.55, 11.0, 0.07, PC["accent"])
            # Subtitle teksti
            if points:
                tx2 = _ptx(slide, 0.7, 4.8, 11.5, 2.3)
                tf2 = tx2.text_frame; tf2.word_wrap = True
                p2  = tf2.paragraphs[0]
                r2  = p2.add_run()
                # Points-daky emoji/ikon aýyr - diňe tekst
                import re as _re
                r2.text = _re.sub(r"^[🔹📊✅💡📌🎯🔑📈🔷▸•]+\s*","",points[0])
                r2.font.name = "Times New Roman"; r2.font.size = PPtx(16)
                r2.font.italic = True; r2.font.color.rgb = PC["light_t"]
                p2.alignment = PP_ALIGN.LEFT
            # Slaýd san
            tx_n = _ptx(slide, 12.2, 6.9, 0.9, 0.45)
            rn = tx_n.text_frame.paragraphs[0].add_run(); rn.text = f"1 / {total}"
            rn.font.name = "Times New Roman"; rn.font.size = PPtx(9)
            rn.font.color.rgb = PC["light_t"]
            tx_n.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

            # Student maglumaty — sag tarap, vertikal merkez
            if student_info:
                _sl = student_info.get("pres_lang","ru")
                _nm = student_info.get("name","")
                _cr = student_info.get("course","")
                _gr = student_info.get("group","")
                # Dile görä ýazgy
                if _sl == "ru":
                    _line1 = f"Выполнил: {_nm}"
                    _line2 = f"{_cr} курс, {_gr}"
                elif _sl == "en":
                    _line1 = f"Prepared by: {_nm}"
                    _line2 = f"Year {_cr}, Group {_gr}"
                elif _sl == "tr":
                    _line1 = f"Hazırlayan: {_nm}"
                    _line2 = f"{_cr}. sınıf, {_gr}"
                else:
                    _line1 = f"Taýarlan: {_nm}"
                    _line2 = f"{_cr} kurs, {_gr}"
                # Aşak-sag burçda, accent fon içinde, ak tekst
                # Fon gutusy
                _si_bg = _prect(slide, 6.0, 5.9, 7.0, 1.35, PC["accent"])
                # Ýuka ak çyzgy ýokarynda
                _prect(slide, 6.0, 5.9, 7.0, 0.04, PRGBColor(0xFF,0xFF,0xFF))
                tx_si = _ptx(slide, 6.1, 6.0, 6.8, 1.2)
                tf_si = tx_si.text_frame; tf_si.word_wrap = True
                # 1-nji setir: at-familiýa (ak, uly)
                p_si1 = tf_si.paragraphs[0]
                p_si1.alignment = PP_ALIGN.CENTER
                r_si1 = p_si1.add_run(); r_si1.text = _line1
                r_si1.font.name = "Times New Roman"; r_si1.font.size = PPtx(15)
                r_si1.font.bold = True
                r_si1.font.color.rgb = PRGBColor(0xFF,0xFF,0xFF)
                # 2-nji setir: kurs, gruppa (ýagty, kursiv)
                p_si2 = tf_si.add_paragraph()
                p_si2.alignment = PP_ALIGN.CENTER
                p_si2.space_before = PPtx(4)
                r_si2 = p_si2.add_run(); r_si2.text = _line2
                r_si2.font.name = "Times New Roman"; r_si2.font.size = PPtx(13)
                r_si2.font.color.rgb = PRGBColor(0xFF,0xFF,0xFF)
                r_si2.font.italic = True

        elif is_last:
            _pbg(slide, PC["dark"])
            # Doly fon suraty
            if img:
                _pimg(slide, img, 0, 0, 13.33, 7.5)
            ov2 = _prect(slide, 0, 0, 13.33, 7.5, PRGBColor(0,0,0))
            _set_transparency(ov2, 0.40)  # 40% aýdyň
            # Sag tarap reňkli zolak
            _prect(slide, 13.0, 0, 0.35, 7.5, PC["accent"])
            # Başlyk setiri ýok (AI belli bolmasyn)
            # Tema
            tx_jt = _ptx(slide, 0.7, 1.1, 11.5, 1.4)
            rjt   = tx_jt.text_frame.paragraphs[0].add_run()
            rjt.text = ttl
            rjt.font.name = "Times New Roman"; rjt.font.size = PPtx(26)
            rjt.font.bold = True; rjt.font.color.rgb = PC["white"]
            tx_jt.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            _prect(slide, 0.7, 2.65, 11.0, 0.06, PC["accent"])
            # Nokatlar - 2 sütün görnüşinde
            import re as _re2
            _n_pts = min(len(points), 4)
            _rows  = (_n_pts + 1) // 2          # 1-2 → 1 setir, 3-4 → 2 setir
            _row_h = min(1.9, 4.2 / _rows)      # Her setiriň beýikligi
            _txt_w = 5.8                         # Tekst gutusynyň ini
            _fsize = 11                          # Font ölçegi

            for pi, pt in enumerate(points[:4]):
                col = pi % 2
                row = pi // 2
                xl  = 0.6 + col * 6.5
                yt  = 2.85 + row * _row_h
                # Reňkli san belgisi
                sh_num = slide.shapes.add_shape(9, Inches(xl), Inches(yt), Inches(0.42), Inches(0.42))
                sh_num.fill.solid(); sh_num.fill.fore_color.rgb = PC["accent"]; sh_num.line.fill.background()
                tf_n2 = sh_num.text_frame; tf_n2.word_wrap = False
                pn2 = tf_n2.paragraphs[0]; pn2.alignment = PP_ALIGN.CENTER
                rn2 = pn2.add_run(); rn2.text = str(pi+1)
                rn2.font.name = "Times New Roman"; rn2.font.size = PPtx(10)
                rn2.font.bold = True; rn2.font.color.rgb = PC["white"]
                # Tekst doly, kesilmeýär
                pt_clean = _re2.sub(r"^[🔹📊✅💡📌🎯🔑📈🔷▸•⚡🌟🔶💎🏆📝🔷]+\s*","",pt)
                _th = _row_h - 0.12
                tx_pt = _ptx(slide, xl+0.5, yt, _txt_w, _th)
                tf_pt = tx_pt.text_frame; tf_pt.word_wrap = True
                rpt = tf_pt.paragraphs[0].add_run()
                rpt.text = pt_clean           # DOLY tekst, kesilmeýär
                rpt.font.name = "Times New Roman"; rpt.font.size = PPtx(_fsize)
                rpt.font.color.rgb = PC["light_t"]
                tf_pt.paragraphs[0].alignment = PP_ALIGN.LEFT
            # Slaýd san
            tx_n = _ptx(slide, 12.2, 6.9, 0.9, 0.45)
            rn = tx_n.text_frame.paragraphs[0].add_run(); rn.text = f"{total} / {total}"
            rn.font.name = "Times New Roman"; rn.font.size = PPtx(9)
            rn.font.color.rgb = PC["light_t"]
            tx_n.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

        else:
            _pbg(slide, PC["bg_light"])
            _prect(slide, 0, 0, 13.33, 1.25, PC["acc2"])
            _prect(slide, 0, 1.23, 13.33, 0.07, PC["accent"])
            tx_h = _ptx(slide, 0.35, 0.1, 11.8, 1.0)
            tf_h = tx_h.text_frame; tf_h.word_wrap = True
            rh = tf_h.paragraphs[0].add_run(); rh.text = ttl
            rh.font.name = "Times New Roman"; rh.font.size = PPtx(24)
            rh.font.bold = True; rh.font.color.rgb = PC["white"]
            tx_n = _ptx(slide, 12.5, 0.1, 0.7, 0.5)
            rn = tx_n.text_frame.paragraphs[0].add_run(); rn.text = f"{idx+1}/{total}"
            rn.font.name = "Times New Roman"; rn.font.size = PPtx(9)
            rn.font.color.rgb = PC["light_t"]
            tx_n.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

            import random as _rnd2
            pos = _rnd2.choice(_IMG_POS)
            if has_chart and chart_dt:
                # Grafik sag tarapda — başlyk (0.35) we caption (0.4) üçin ýer
                tl,tt,tw,th = 0.3,1.35,6.3,5.9
                _add_chart_pptx(slide, chart_dt, 6.7, 1.75, 6.3, 4.85)
            elif img:
                if pos == "left":
                    _pimg(slide, img, 0, 1.31, 5.1, 6.19); tl,tt,tw,th = 5.3,1.35,7.7,5.9
                elif pos == "right":
                    _pimg(slide, img, 8.0, 1.31, 5.33, 6.19); tl,tt,tw,th = 0.3,1.35,7.5,5.9
                elif pos in ("bottom_left","bottom_right","bottom"):
                    # Aşakda doly giňlik, 10px açyk
                    _GAP = 0.1  # inç ~9px
                    _img_w = 13.33 - 2*_GAP
                    _img_h = 2.8
                    _img_t = 7.5 - _img_h - _GAP
                    _pimg(slide, img, _GAP, _img_t, _img_w, _img_h)
                    tl,tt,tw,th = 0.3, 1.35, 12.7, _img_t - 1.45
                else:
                    _GAP = 0.1
                    _img_w = 13.33 - 2*_GAP
                    _img_h = 2.8
                    _img_t = 7.5 - _img_h - _GAP
                    _pimg(slide, img, _GAP, _img_t, _img_w, _img_h)
                    tl,tt,tw,th = 0.3, 1.35, 12.7, _img_t - 1.45
            else:
                tl,tt,tw,th = 0.3,1.35,12.7,5.9

            # Tekst sygýan mukdary hasapla
            _avail_h = th - 0.2        # elýeterli beýiklik inç
            _max_pts = min(len(points), 4)
            # Font ölçegi beýiklige görä
            if _avail_h >= 4.5:
                _fsize_m = 13; _line_sp_m = 1.25
            elif _avail_h >= 3.0:
                _fsize_m = 12; _line_sp_m = 1.1
            else:
                _fsize_m = 11; _line_sp_m = 1.0
            # Bir nokat üçin beýiklik (inç)
            _pt_h    = _avail_h / _max_pts
            step     = _pt_h

            tx_b = _ptx(slide, tl+0.38, tt+0.08, tw-0.4, _avail_h)
            tf_b = tx_b.text_frame; tf_b.word_wrap = True

            for pi, pt in enumerate(points[:_max_pts]):
                dy = tt + 0.18 + pi * step
                if dy + 0.2 < tt + th:
                    _pdot(slide, tl+0.05, dy, PC["accent"])
                pp = tf_b.paragraphs[0] if pi == 0 else tf_b.add_paragraph()
                rr = pp.add_run()
                # Her nokat üçin tekst çäkle - 1 setirlere sygdyr
                _chars_per_line = int((tw - 0.6) * 14)  # takmynan
                _max_chars_pt   = _chars_per_line * int(_pt_h / (_fsize_m / 72.0) * _line_sp_m + 0.5)
                rr.text = pt[:min(len(pt), max(80, _max_chars_pt))]
                rr.font.name = "Times New Roman"; rr.font.size = PPtx(_fsize_m)
                rr.font.color.rgb = PC["dark_t"]
                pp.space_before = PPtx(int(step * 72 * 0.35))
                pp.alignment = PP_ALIGN.LEFT

    buf = io.BytesIO(); prs.save(buf); buf.seek(0)
    return buf.getvalue()


async def gen_image(prompt: str):
    if not STABILITY_API_KEY or len(STABILITY_API_KEY) < 10:
        return None
    try:
        headers = {"Authorization": f"Bearer {STABILITY_API_KEY}", "Accept": "image/*"}
        files   = {
            "prompt":        (None, prompt + ", professional, photorealistic"),
            "aspect_ratio":  (None, "16:9"),
            "output_format": (None, "jpeg"),
        }
        async with httpx.AsyncClient(timeout=httpx.Timeout(connect=30,read=120,write=30,pool=10)) as cl:
            r = await cl.post(STABILITY_URL, headers=headers, files=files)
            log.info(f"Stability jogap: {r.status_code}")
            if r.status_code == 200:
                log.info(f"✅ Surat taýar ({len(r.content)} baýt)")
                return r.content
            log.error(f"❌ Stability {r.status_code}: {r.text[:300]}")
            return None
    except Exception as e:
        log.error(f"❌ Surat şowsuz: {type(e).__name__}: {e}"); return None


def _fix_json_pptx(raw: str) -> str:
    objects = []; depth = 0; start = None
    for i, ch in enumerate(raw):
        if ch == "{":
            if depth == 0: start = i
            depth += 1
        elif ch == "}":
            depth -= 1
            if depth == 0 and start is not None:
                try: objects.append(json.loads(raw[start:i+1]))
                except: pass
                start = None
    if not objects: raise RuntimeError("JSON çykmady")
    return json.dumps(objects, ensure_ascii=False)


async def _pptx_one_batch(theme: str, slide_nums: list, total_slides: int,
                           pres_lang: str, system: str, lang_str: str) -> list:
    """Bir toparda birnäçe slaýd üçin JSON al (max 6 slaýd)"""
    batch_size = len(slide_nums)
    is_first   = (1 in slide_nums)
    is_last    = (total_slides in slide_nums)
    lang_instr = {
        "ru": "ВСЕ тексты пиши на русском языке.",
        "en": "Write ALL texts in English language.",
        "tr": "Tüm metinleri Türkiye Türkçesiyle yaz.",
        "tk": "Ähli tekstleri türkmen dilinde ýaz.",
    }.get(pres_lang, "ВСЕ тексты пиши на русском языке.")

    slides_desc = ", ".join(str(x) for x in slide_nums)
    prompt = (
        f"Тема презентации: «{theme}»\n"
        f"Создай ТОЛЬКО слайды №{slides_desc} (из {total_slides} всего).\n"
        f"{lang_instr}\n\n"
        f"Верни JSON массив ровно из {batch_size} объектов.\n"
        f"Формат каждого объекта:\n"
        f'{{"slide_num":N,"title":"Заголовок (4-7 слов)","points":["🔹 Пункт минимум 25 слов содержательно","📊 Пункт минимум 25 слов с фактами","✅ Пункт минимум 25 слов","💡 Пункт минимум 25 слов"],"image_prompt":"English photorealistic description 15 words","has_chart":false,"chart_data":{{}}}}\n\n'
        f"Правила:\n"
    )
    if is_first:
        prompt += f"- Слайд 1: title=тема презентации, points=['Развёрнутый вводный подзаголовок 30-40 слов'], has_chart=false\n"
    if is_last:
        prompt += f"- Слайд {total_slides} (ПОСЛЕДНИЙ): title='Заключение/Выводы', 4 вывода минимум 30 слов каждый, has_chart=false\n"
    prompt += (
        f"- Слайды с графиком: has_chart=true, chart_data={{labels:[реальные категории на {lang_str} языке],values:[числа],title:'заголовок',type:'bar'/'pie'/'line',caption:'что показывает 1-2 предложения на {lang_str}',y_label:'единица'}}\n"
        f"- Остальные слайды: 3-4 пункта минимум 25 слов каждый\n"
        f"- Каждый объект ПОЛНЫЙ, JSON не обрывается\n"
        f"- Только JSON массив"
    )
    headers = {"Authorization": f"Bearer {DEEPSEEK_API_KEY}",
               "Content-Type": "application/json", "Accept-Encoding": "identity"}
    body = {"model": DEEPSEEK_MODEL,
            "messages": [{"role":"system","content":system}, {"role":"user","content":prompt}],
            "max_tokens": min(8000, batch_size * 600 + 500),
            "temperature": 0.1}
    async with httpx.AsyncClient(timeout=httpx.Timeout(connect=60,read=300,write=60,pool=30)) as cl:
        r = await cl.post(DEEPSEEK_URL, headers=headers, json=body)
        if r.status_code != 200: raise RuntimeError(f"DeepSeek {r.status_code}: {r.text[:200]}")
        raw = r.json()["choices"][0]["message"]["content"].strip()
    raw = re.sub(r"```json\s*","",raw); raw = re.sub(r"```\s*","",raw); raw = raw.strip()
    s = raw.find("["); e = raw.rfind("]")
    if s == -1 or e == -1:
        log.warning(f"Batch {slide_nums}: JSON tapylmady, fallback")
        return []
    try: data = json.loads(raw[s:e+1])
    except json.JSONDecodeError:
        try: data = json.loads(_fix_json_pptx(raw))
        except: data = []
    # slide_num meýdanyny aýyr, tertibe görä gaýtar
    result = []
    for item in data:
        if isinstance(item, dict):
            item.pop("slide_num", None)
            result.append(item)
    return result


async def call_deepseek_pptx(theme: str, slides: int, pres_lang: str) -> list:
    lang_str = {"ru":"русском","en":"English","tk":"туркменском","tr":"турецком"}.get(pres_lang,"русском")
    if pres_lang == "en":
        system = ("You are a presentation expert. Reply ONLY with a valid JSON array, no markdown, no extra text. "
                  "Write ALL slide texts ONLY in English.")
    elif pres_lang == "tr":
        system = ("Sen sunum uzmanısın. SADECE geçerli JSON dizisi döndür, markdown yok. "
                  "Tüm metinleri YALNIZCA Türkçe yaz.")
    elif pres_lang == "tk":
        system = ("Sen prezentasiýa hünärmeni. Diňe dogry JSON massiwi gaýtar, markdown ýok. "
                  "Ähli tekstleri diňe türkmen dilinde ýaz.")
    else:
        system = ("Ты эксперт по презентациям. Отвечай ТОЛЬКО валидным JSON массивом, без markdown, без лишнего текста. "
                  "Пиши ВСЕ тексты ТОЛЬКО на русском языке.")

    # Slaýdlary toparlara böl: max 5 slaýd bir toparda
    BATCH = 5
    all_nums = list(range(1, slides + 1))
    batches  = [all_nums[i:i+BATCH] for i in range(0, len(all_nums), BATCH)]

    all_data = []
    for batch_nums in batches:
        batch_data = await _pptx_one_batch(theme, batch_nums, slides, pres_lang, system, lang_str)
        # Eger az geldi - fallback goş
        while len(batch_data) < len(batch_nums):
            idx = len(all_data) + len(batch_data) + 1
            batch_data.append({
                "title": f"{theme} — {idx}",
                "points": [
                    f"🔹 {theme} barada {idx}-nji bölüm.",
                    f"📊 Bu bölümde esasy maglumatlar beýan edilýär.",
                    f"✅ Görkezijiler we netijelar seljerilýär.",
                ],
                "image_prompt": f"{theme} professional photography high quality",
                "has_chart": False,
                "chart_data": {}
            })
        all_data.extend(batch_data[:len(batch_nums)])
        log.info(f"PPTX batch {batch_nums[0]}-{batch_nums[-1]}: {len(batch_data)} slaýd alyndy")

    return all_data[:slides]


# ─── PPTX HANDLERS ───────────────────────────────────────────────────

@router.message(St.sp1)
async def hp1_theme(msg: Message, state: FSMContext):
    if len((msg.text or "").strip()) < 5:
        d = await state.get_data(); lang = d.get("ui_lang","tk")
        errs = {"tk":"❌ Iň az 5 harp ýazyň!","ru":"❌ Минимум 5 символов!","en":"❌ At least 5 chars!"}
        await msg.answer(errs.get(lang,"❌")); return
    await state.update_data(pptx_theme=msg.text.strip())
    d = await state.get_data(); lang = d.get("ui_lang","tk")
    ok = {"tk":"✅ Tema kabul edildi!","ru":"✅ Тема принята!","en":"✅ Topic accepted!"}
    q = {"tk":"📌 <b>2/5:</b> Adyňyzy we Familiýaňyzy ýazyň\n<i>Mysal: Myrat Mämmedow</i>",
         "ru":"📌 <b>2/5:</b> Введите ваше имя и фамилию\n<i>Пример: Иван Иванов</i>",
         "en":"📌 <b>2/5:</b> Enter your full name\n<i>Example: John Smith</i>"}
    await msg.answer(ok.get(lang,"✅") + "\n\n" + q.get(lang,""), parse_mode="HTML")
    await state.set_state(St.sp1b)

@router.message(St.sp1b)
async def hp1b_name(msg: Message, state: FSMContext):
    d = await state.get_data(); lang = d.get("ui_lang","tk")
    if len((msg.text or "").strip()) < 3:
        errs = {"tk":"❌ Iň az 3 harp ýazyň!","ru":"❌ Минимум 3 символа!","en":"❌ At least 3 chars!"}
        await msg.answer(errs.get(lang,"❌")); return
    await state.update_data(pptx_fullname=msg.text.strip())
    q = {"tk":"📌 <b>3/5:</b> Haýsy kursda okaýarsyňyz?",
         "ru":"📌 <b>3/5:</b> На каком курсе вы учитесь?",
         "en":"📌 <b>3/5:</b> Which year are you in?"}
    ok = {"tk":"✅ Kabul edildi!","ru":"✅ Принято!","en":"✅ Accepted!"}
    await msg.answer(ok.get(lang,"✅") + "\n\n" + q.get(lang,""),
                     parse_mode="HTML", reply_markup=kb_crs(lang))
    await state.set_state(St.sp1c)

@router.callback_query(St.sp1c, F.data.startswith("crs:"))
async def hp1c_course(cb: CallbackQuery, state: FSMContext):
    n = cb.data.split(":")[1]
    await state.update_data(pptx_course=n)
    d = await state.get_data(); lang = d.get("ui_lang","tk")
    ok = {"tk":f"✅ {n}-nji kurs!","ru":f"✅ {n}-й курс!","en":f"✅ Year {n}!"}
    q = {"tk":"📌 <b>4/6:</b> Toparыňyzyň adyny ýazyň\n<i>Mysal: EHM-22</i>",
         "ru":"📌 <b>4/6:</b> Введите название вашей группы\n<i>Пример: ЭВМ-22</i>",
         "en":"📌 <b>4/6:</b> Enter your group name\n<i>Example: CS-22</i>"}
    await ask(cb, ok.get(lang,"✅") + "\n\n" + q.get(lang,""))
    await state.set_state(St.sp1d); await cb.answer()

@router.message(St.sp1d)
async def hp1d_group(msg: Message, state: FSMContext):
    d = await state.get_data(); lang = d.get("ui_lang","tk")
    if len((msg.text or "").strip()) < 2:
        errs = {"tk":"❌ Iň az 2 harp!","ru":"❌ Минимум 2 символа!","en":"❌ At least 2 chars!"}
        await msg.answer(errs.get(lang,"❌")); return
    await state.update_data(pptx_group=msg.text.strip())
    q = {"tk":"📌 <b>5/6:</b> Prezentasiýa haýsy dilde bolsun?",
         "ru":"📌 <b>5/6:</b> На каком языке должна быть презентация?",
         "en":"📌 <b>5/6:</b> In which language should the presentation be?"}
    ok = {"tk":"✅ Kabul edildi!","ru":"✅ Принято!","en":"✅ Accepted!"}
    await msg.answer(ok.get(lang,"✅") + "\n\n" + q.get(lang,""),
                     parse_mode="HTML", reply_markup=KB_PPTX_LANG)
    await state.set_state(St.sp2)

@router.callback_query(St.sp2, F.data.startswith("plang:"))
async def hp2_lang(cb: CallbackQuery, state: FSMContext):
    pl = cb.data.split(":")[1]
    await state.update_data(pptx_lang=pl)
    d = await state.get_data(); lang = d.get("ui_lang","tk")
    lang_name = {"ru":"🇷🇺 Rusça / Русский","en":"🇬🇧 Iňlisçe / English","tr":"🇹🇷 Türkçe"}.get(pl,"?")
    q = {"tk":f"✅ {lang_name} saýlandy!\n\n📌 <b>5/5:</b> Näçe slaýd bolmaly?",
         "ru":f"✅ {lang_name} выбран!\n\n📌 <b>5/5:</b> Сколько слайдов?",
         "en":f"✅ {lang_name} selected!\n\n📌 <b>5/5:</b> How many slides?"}
    await ask(cb, q.get(lang,""), KB_PPTX_SLIDES)
    await state.set_state(St.sp3); await cb.answer()

@router.callback_query(St.sp3, F.data.startswith("sl:"))
async def hp3_slides(cb: CallbackQuery, state: FSMContext, bot: Bot):
    n = int(cb.data.split(":")[1])
    await state.update_data(pptx_slides=n)
    d   = await state.get_data()
    uid = cb.from_user.id
    theme      = d.get("pptx_theme","")
    pres_lang  = d.get("pptx_lang","ru")
    pptx_name  = d.get("pptx_fullname","")
    pptx_course= d.get("pptx_course","")
    pptx_group = d.get("pptx_group","")

    if uid in ACTIVE_GENERATES:
        await cb.answer("⏳ Eýýäm taýarlanýar!", show_alert=True); return
    ACTIVE_GENERATES.add(uid); CANCELLED_GENERATES.discard(uid)

    _ui_lang_init = d.get("ui_lang","tk")
    _start_txt = {"tk":"Başlanýar...","ru":"Начинаем...","en":"Starting..."}.get(_ui_lang_init,"Başlanýar...")
    _hdr_init  = {"tk":"Prezentasiýa taýarlanýar","ru":"Презентация создаётся","en":"Presentation is being created"}.get(_ui_lang_init,"Prezentasiýa taýarlanýar")
    _name_line = f"\n👤 {pptx_name}" if pptx_name else ""
    prog = await cb.message.edit_text(
        f"⚙️ <b>{_hdr_init}...</b>\n\n"
        f"📝 <i>{theme}</i>{_name_line}\n\n"
        f"<code>[░░░░░░░░░░]</code> <b>0%</b>\n<i>{_start_txt}</i>",
        parse_mode="HTML")
    await cb.answer(); await state.clear()
    cid = cb.message.chat.id; mid = prog.message_id

    _ui_lang = d.get("ui_lang","tk")
    _pptx_hdr = {"tk":"Prezentasiýa taýarlanýar","ru":"Презентация создаётся","en":"Presentation is being created"}.get(_ui_lang,"Prezentasiýa taýarlanýar")
    async def upd(pct, status, tot=0, don=0):
        bar = "█"*(pct//10)+"░"*(10-pct//10)
        sl  = f"\n🖼 {don}/{tot}" if tot > 0 else ""
        _nl = f"\n👤 {pptx_name}" if pptx_name else ""
        try:
            await bot.edit_message_text(
                f"⚙️ <b>{_pptx_hdr}...</b>\n\n"
                f"📝 <i>{theme}</i>{_nl}{sl}\n\n"
                f"<code>[{bar}]</code> <b>{pct}%</b>\n<i>{status}</i>",
                chat_id=cid, message_id=mid, parse_mode="HTML")
        except: pass

    try:
        BATCH = 5
        n_batches = (n + BATCH - 1) // BATCH
        _s1 = {"tk":f"📝 Slaýd mazmuny ýazylýar (0/{n_batches} blok)...",
               "ru":f"📝 Создание слайдов (0/{n_batches} блок)...",
               "en":f"📝 Writing slides (0/{n_batches} block)..."}.get(_ui_lang,"📝")
        await upd(10, _s1, n, 0)

        # Batch progress - her toparda upd
        all_nums  = list(range(1, n + 1))
        batches   = [all_nums[i:i+BATCH] for i in range(0, len(all_nums), BATCH)]
        slides_data = []
        import httpx as _httpx2, re as _re3, json as _json3
        lang_str2 = {"ru":"русском","en":"English","tk":"туркменском","tr":"турецком"}.get(pres_lang,"русском")
        if pres_lang == "en":
            system2 = ("You are a presentation expert. Reply ONLY with a valid JSON array, no markdown, no extra text. Write ALL slide texts ONLY in English.")
        elif pres_lang == "tr":
            system2 = ("Sen sunum uzmanısın. SADECE geçerli JSON dizisi döndür, markdown yok. Tüm metinleri YALNIZCA Türkçe yaz.")
        elif pres_lang == "tk":
            system2 = ("Sen prezentasiýa hünärmeni. Diňe dogry JSON massiwi gaýtar, markdown ýok. Ähli tekstleri diňe türkmen dilinde ýaz.")
        else:
            system2 = ("Ты эксперт по презентациям. Отвечай ТОЛЬКО валидным JSON массивом, без markdown, без лишнего текста. Пиши ВСЕ тексты ТОЛЬКО на русском языке.")

        for _bi, batch_nums in enumerate(batches):
            _bs_lbl = {"tk":f"📝 {_bi+1}/{n_batches} blok ýazylýar...",
                       "ru":f"📝 Блок {_bi+1}/{n_batches}...",
                       "en":f"📝 Block {_bi+1}/{n_batches}..."}.get(_ui_lang,"📝")
            await upd(10 + int(40*_bi/n_batches), _bs_lbl, n, 0)
            batch_res = await _pptx_one_batch(theme, batch_nums, n, pres_lang, system2, lang_str2)
            while len(batch_res) < len(batch_nums):
                idx = len(slides_data) + len(batch_res) + 1
                batch_res.append({
                    "title": f"{theme} — {idx}",
                    "points": ["🔹 Esasy maglumatlar.","📊 Seljerme we görkezijiler.","✅ Netijeler we çözgütler."],
                    "image_prompt": f"{theme} professional photography",
                    "has_chart": False, "chart_data": {}
                })
            slides_data.extend(batch_res[:len(batch_nums)])
            log.info(f"Batch {batch_nums[0]}-{batch_nums[-1]}: {len(batch_res)} slaýd")

        # Slaýd sanyna görä surat sany
        import random as _rnd

        # Slaýd sany boýunça paýlaşyk:
        #  5 slaýd → 1 surat + 2 grafik + 2 fon (1-nji,soňky)
        #  7 slaýd → 2 surat + 3 grafik + 2 fon
        # 10 slaýd → 4 surat + 4 grafik + 2 fon
        # 12 slaýd → 5 surat + 5 grafik + 2 fon
        _PLAN = {
            7:  {"img": 4, "chart": 1},
            10: {"img": 6, "chart": 2},
            12: {"img": 8, "chart": 2},
        }
        # Iň ýakyn plany tap
        _plan_key = min(_PLAN.keys(), key=lambda k: abs(k - n))
        _plan = _PLAN[_plan_key]
        n_img   = min(_plan["img"],   max(0, n-2))
        n_chart = min(_plan["chart"], max(0, n-2-n_img))

        # Orta slaýd indeksleri (1-nji we soňky aýrylýar)
        _mid = list(range(1, n-1))
        _rnd.shuffle(_mid)
        _img_slots   = set(_mid[:n_img])
        _chart_slots = set(_mid[n_img:n_img+n_chart])
        # img_slots we chart_slots kesişmesin
        _chart_slots = _chart_slots - _img_slots
        # Galanlar - diňe ikon we tekst

        # Grafik tiplerini tertipli belirle: bar → pie → line → line...
        _ch_types_seq = ["bar","pie","line","line","line"]
        for _seq_i, _ci in enumerate(sorted(_chart_slots)):
            _forced_type = _ch_types_seq[min(_seq_i, len(_ch_types_seq)-1)]

            _sld_title = slides_data[_ci].get("title","")
            # Dile görä caption we y_label
            _cap_tmpl = {
                "ru": f"Сравнительный анализ: {_sld_title[:40]}. Значения за исследуемый период.",
                "en": f"Comparative analysis: {_sld_title[:40]}. Values for the study period.",
                "tr": f"Karşılaştırmalı analiz: {_sld_title[:40]}. Araştırma dönemi verileri.",
                "tk": f"Deňeşdirme seljerişi: {_sld_title[:40]}. Döwür boýunça görkezijiler.",
            }.get(pres_lang, f"Analysis: {_sld_title[:40]}")
            _dyn_tmpl = {
                "ru": f"Динамика показателей: {slides_data[_ci].get('chart_data',{}).get('title','')[:40]}",
                "en": f"Indicator dynamics: {slides_data[_ci].get('chart_data',{}).get('title','')[:40]}",
                "tr": f"Gösterge dinamiği: {slides_data[_ci].get('chart_data',{}).get('title','')[:40]}",
                "tk": f"Görkezijileriň dinamikasy: {slides_data[_ci].get('chart_data',{}).get('title','')[:40]}",
            }.get(pres_lang, "")
            _y_lbl_def = {"ru":"%","en":"%","tr":"%","tk":"%"}.get(pres_lang,"%")
            _line_labels = {
                "ru":["2020","2021","2022","2023"],
                "en":["2020","2021","2022","2023"],
                "tr":["2020","2021","2022","2023"],
                "tk":["2020","2021","2022","2023"],
            }.get(pres_lang,["2020","2021","2022","2023"])
            # Fallback labels - tema boýunça anyk sözler (tema sözlerinden al)
            _theme_words = [w for w in _sld_title.split() if len(w) > 3][:4]
            if len(_theme_words) >= 2:
                _bar_lbs = _theme_words[:4]
                while len(_bar_lbs) < 4:
                    _bar_lbs.append(f"{_bar_lbs[-1]}+")
            else:
                _bar_lbs = {
                    "ru":["2020 г.","2021 г.","2022 г.","2023 г."],
                    "en":["2020","2021","2022","2023"],
                    "tr":["2020","2021","2022","2023"],
                    "tk":["2020","2021","2022","2023"],
                }.get(pres_lang,["2020","2021","2022","2023"])
            _bar_labels = _bar_lbs

            if not slides_data[_ci].get("chart_data") or not slides_data[_ci]["chart_data"].get("labels"):
                _vals = [_rnd.randint(35,95) for _ in range(4)]
                _x_lbl_def = {
                    "ru": "Период" if _forced_type=="line" else "Категории",
                    "en": "Period"  if _forced_type=="line" else "Categories",
                    "tr": "Dönem"   if _forced_type=="line" else "Kategoriler",
                    "tk": "Döwür"   if _forced_type=="line" else "Kategoriýalar",
                }.get(pres_lang,"")
                slides_data[_ci]["chart_data"] = {
                    "labels": _line_labels if _forced_type=="line" else _bar_labels,
                    "values": _vals,
                    "title": _sld_title,
                    "type": _forced_type,
                    "y_label": _y_lbl_def,
                    "x_label": _x_lbl_def,
                    "caption": _cap_tmpl,
                }
            else:
                slides_data[_ci]["chart_data"]["type"] = _forced_type
                if not slides_data[_ci]["chart_data"].get("caption"):
                    slides_data[_ci]["chart_data"]["caption"] = _dyn_tmpl
                if not slides_data[_ci]["chart_data"].get("y_label"):
                    slides_data[_ci]["chart_data"]["y_label"] = _y_lbl_def
                # Labels-lary hem dile görä düzelt
                existing_labels = slides_data[_ci]["chart_data"].get("labels",[])
                # Eger labels Kiril harply bolsa we dil EN/TR/TK bolsa — çalyş
                if pres_lang in ("en","tr","tk") and existing_labels:
                    import re as _re_ch
                    if any(_re_ch.search(r"[а-яА-Я]", str(lb)) for lb in existing_labels):
                        slides_data[_ci]["chart_data"]["labels"] = _line_labels if _forced_type=="line" else _bar_labels
            # has_chart diňe chart_slots üçin
            slides_data[_ci]["has_chart"] = True

        # img_slots slaýdlarynda has_chart aýyr (surat bolmaly)
        for _ii in _img_slots:
            slides_data[_ii]["has_chart"] = False

        # Ikon saýlawy — random
        _ICONS = ["📌","📊","✅","💡","🔹","🎯","🔑","📈","🔷","⚡","🌟","🔶","💎","🏆","📝"]

        images = []
        img_count = 0
        for i, sld in enumerate(slides_data):
            # Random ikonlar goş
            pts = sld.get("points", [])
            _icons = _rnd.sample(_ICONS, min(len(pts), len(_ICONS)))
            new_pts = []
            for pi, pt in enumerate(pts):
                import re as _re2
                pt_clean = _re2.sub(r"^[^\w\s]+\s*","",pt)
                icon = _icons[pi] if pi < len(_icons) else "▪"
                new_pts.append(f"{icon} {pt_clean}")
            sld["points"] = new_pts

            if i == 0 or i == n-1:
                # Fon slaýdlar — hökman surat
                img_count += 1
                await upd(20+int(55*i/len(slides_data)),
                          {"tk":f"🖼 Fon surat {img_count} döredilýär...","ru":f"🖼 Фон изображение {img_count}...","en":f"🖼 Background image {img_count}..."}.get(_ui_lang,f"🖼 {img_count}"),
                          len(slides_data), img_count)
                images.append(await gen_image(sld.get("image_prompt", theme)))
            elif i in _img_slots:
                img_count += 1
                await upd(20+int(55*i/len(slides_data)),
                          {"tk":f"🖼 Surat {img_count} döredilýär ({i+1}/{n})...","ru":f"🖼 Изображение {img_count} ({i+1}/{n})...","en":f"🖼 Image {img_count} ({i+1}/{n})..."}.get(_ui_lang,f"🖼 {img_count}"),
                          len(slides_data), img_count)
                images.append(await gen_image(sld.get("image_prompt", theme)))
            else:
                images.append(None)

        _s85 = {"tk":"📊 Prezentasiýa ýygnalyp durýar...","ru":"📊 Сборка презентации...","en":"📊 Building presentation..."}.get(_ui_lang,"📊")
        await upd(85,_s85,len(slides_data),len(slides_data))
        _student_info = {
            "name": pptx_name,
            "course": pptx_course,
            "group": pptx_group,
            "pres_lang": pres_lang,
        }
        pptx_bytes = await asyncio.get_running_loop().run_in_executor(
            None, build_pptx, slides_data, theme, images, _student_info)

        safe = re.sub(r"[^\w]","_",theme)[:15]
        fname = f"presentation_{safe}.pptx"
        is_free = _db_is_first_free(uid)

        if is_free:
            _db_add_order(uid, "pptx", theme, 0, is_free=True, status="delivered")
            _db_mark_free_used(uid)
            lang_f = d.get("ui_lang","tk")
            free_hdr = {"tk":f"🎁 <b>Birinji sargydyňyz MUGT!</b>\n\n📝 <i>{theme}</i>",
                        "ru":f"🎁 <b>Первый заказ БЕСПЛАТНО!</b>\n\n📝 <i>{theme}</i>",
                        "en":f"🎁 <b>First order FREE!</b>\n\n📝 <i>{theme}</i>"}
            _name_cap = f"\n👤 {pptx_name}" + (f" | {pptx_course}" + {"tk":"-nji kurs","ru":"-й курс","en":"th year"}.get(lang_f,"") if pptx_course else "") if pptx_name else ""
            free_cap = {"tk":f"✅ <b>Prezentasiýaňyz taýar!</b>\n\n📊 {n} slaýd | 📝 {theme}{_name_cap}\n\n❓ Sorag: <code>{CONTACT_PHONE}</code>\nTäze sargyt: /start",
                        "ru":f"✅ <b>Ваша презентация готова!</b>\n\n📊 {n} слайдов | 📝 {theme}{_name_cap}\n\n❓ Вопросы: <code>{CONTACT_PHONE}</code>\nНовый заказ: /start",
                        "en":f"✅ <b>Your presentation is ready!</b>\n\n📊 {n} slides | 📝 {theme}{_name_cap}\n\n❓ Questions: <code>{CONTACT_PHONE}</code>\nNew order: /start"}
            await bot.edit_message_text(free_hdr.get(lang_f,free_hdr["tk"]),
                chat_id=cid, message_id=mid, parse_mode="HTML")
            await bot.send_document(uid, BufferedInputFile(pptx_bytes, filename=fname),
                caption=free_cap.get(lang_f,free_cap["tk"]), parse_mode="HTML",
                request_timeout=120)
        else:
            oid = _db_add_order(uid, "pptx", theme, PRICE_STARS["pptx"], is_free=False, status="pending")
            lang_p = d.get("ui_lang","tk")
            PAYMENT_PENDING[uid] = {
                "bytes": pptx_bytes,
                "data":  {"service":"pptx","theme":theme,"ui_lang":lang_p},
                "fname": fname,
                "order_id": oid,
            }
            pay_msgs = {
                "tk": (f"✅ <b>Prezentasiýaňyz taýar boldy!</b>\n\n"
                       f"📝 <i>{theme}</i>\n📊 {n} slaýd\n\n"
                       f"⭐ <b>Töleg üçin aşakdaky düwmä basyň:</b>\n"
                       f"💰 Baha: <b>{PRICE_STARS['pptx']} Telegram Stars</b>"),
                "ru": (f"✅ <b>Ваша презентация готова!</b>\n\n"
                       f"📝 <i>{theme}</i>\n📊 {n} слайдов\n\n"
                       f"⭐ <b>Для оплаты нажмите кнопку ниже:</b>\n"
                       f"💰 Стоимость: <b>{PRICE_STARS['pptx']} Telegram Stars</b>"),
                "en": (f"✅ <b>Your presentation is ready!</b>\n\n"
                       f"📝 <i>{theme}</i>\n📊 {n} slides\n\n"
                       f"⭐ <b>Press the button below to pay:</b>\n"
                       f"💰 Price: <b>{PRICE_STARS['pptx']} Telegram Stars</b>"),
            }
            await bot.edit_message_text(
                pay_msgs.get(lang_p, pay_msgs["tk"]),
                chat_id=cid, message_id=mid,
                parse_mode="HTML", reply_markup=kb_pay(lang_p))

    except Exception as exc:
        import traceback; log.error(f"PPTX: {exc}\n{traceback.format_exc()}")
        try:
            await bot.edit_message_text(
                f"❌ <b>Ýalňyşlyk!</b>\n\n<code>{str(exc)[:200]}</code>\n\n/start",
                chat_id=cid, message_id=mid, parse_mode="HTML")
        except: pass
    finally:
        ACTIVE_GENERATES.discard(uid)




@router.callback_query(F.data == "pay:stars")
async def h_pay_stars(cb: CallbackQuery, bot: Bot):
    uid   = cb.from_user.id
    if uid not in PAYMENT_PENDING:
        await cb.answer("❌ Sargyt tapylmady!", show_alert=True); return
    info  = PAYMENT_PENDING[uid]
    d     = info.get("data", {})
    svc   = d.get("service", "pptx")
    theme = d.get("theme", "")
    lang  = d.get("ui_lang", "tk")
    stars_amt = PRICE_STARS.get(svc, 2)
    titles = {"tk":"Akademik iş 📄","ru":"Академическая работа 📄","en":"Academic work 📄"}
    descs  = {"tk":f"Tema: {theme}","ru":f"Тема: {theme}","en":f"Topic: {theme}"}
    if svc == "pptx":
        titles = {"tk":"Prezentasiýa 📊","ru":"Презентация 📊","en":"Presentation 📊"}
    log.info(f"Stars invoice: uid={uid} svc={svc} amount={stars_amt}")
    try:
        await bot.send_invoice(
            chat_id=uid,
            title=titles.get(lang, titles["tk"]),
            description=descs.get(lang, descs["tk"])[:100],
            payload=f"{svc}_{uid}",
            currency="XTR",
            prices=[LabeledPrice(label=titles.get(lang,"Iş"), amount=stars_amt)]
        )
        await cb.answer()
    except Exception as e:
        log.error(f"Stars invoice error: {e}")
        await cb.answer(f"❌ Stars töleg şowsuz!", show_alert=True)
        await cb.message.answer(f"❌ Ýalňyşlyk: {e}\n\nKart bilen töläp bilersiňiz.")


@router.pre_checkout_query()
async def h_pre_checkout(pcq: PreCheckoutQuery, bot: Bot):
    await bot.answer_pre_checkout_query(pcq.id, ok=True)


@router.message(F.successful_payment)
async def h_stars_paid(msg: Message, bot: Bot):
    uid  = msg.from_user.id
    info = PAYMENT_PENDING.pop(uid, None)
    if not info: return
    d     = info.get("data", {})
    fb    = info.get("bytes", b"")
    svc   = d.get("service","referat")
    theme = d.get("theme","")
    lang  = d.get("ui_lang","tk")
    stars = PRICE_STARS.get(svc, 2)  # synag üçin 2, hakyky: 159
    oid   = info.get("order_id")
    if oid: _db_update_order_status(oid,"delivered")
    _db_mark_paid(uid, stars)
    # faýl adyny kesgitle
    if svc == "pptx":
        fname = info.get("fname", f"presentation_{theme[:10]}.pptx")
    else:
        safe  = __import__("re").sub(r"[^\w]","_",theme)[:15]
        fname = f"{svc}_{safe}.docx"
    lang_sp = info.get("data",{}).get("ui_lang","tk") if info else lang
    del_msg = {"tk":f"✅ <b>Taýar!</b>\n\n📝 {theme}\n\n❓ Sorag: <code>{CONTACT_PHONE}</code>\nTäze sargyt: /start",
               "ru":f"✅ <b>Готово!</b>\n\n📝 {theme}\n\n❓ Вопросы: <code>{CONTACT_PHONE}</code>\nНовый заказ: /start",
               "en":f"✅ <b>Ready!</b>\n\n📝 {theme}\n\n❓ Questions: <code>{CONTACT_PHONE}</code>\nNew order: /start"}
    await bot.send_document(uid, BufferedInputFile(fb, filename=fname),
        caption=del_msg.get(lang_sp, del_msg["tk"]), parse_mode="HTML",
        request_timeout=120)
    log.info(f"⭐ Stars: uid={uid} {stars}★ {svc}")

# ════════════════════════════════════════════════════════════════════
async def main():
    from aiogram.client.default import DefaultBotProperties
    from aiogram.client.session.aiohttp import AiohttpSession

    session = AiohttpSession(timeout=600)
    bot = Bot(token=BOT_TOKEN, session=session,
              default=DefaultBotProperties(parse_mode="HTML"))
    me  = await bot.get_me()
    log.info(f"✅ @{me.username} işe başlady!")

    storage = None
    if REDIS_URL and _HAS_REDIS:
        try:
            storage = RedisStorage.from_url(REDIS_URL)
            log.info("✅ Redis storage ulanylýar")
        except Exception as e:
            log.warning(f"⚠️ Redis ulanylmady ({e}), MemoryStorage-a geçýär")
            storage = MemoryStorage()
    else:
        storage = MemoryStorage()
        log.info("⚠️  MemoryStorage ulanylýar (bot ýapylsa state ýitýär)")

    dp  = Dispatcher(storage=storage)
    dp.include_router(router)

    while True:
        try:
            await dp.start_polling(
                bot,
                allowed_updates=["message","callback_query","pre_checkout_query"],
                polling_timeout=30,
            )
        except Exception as e:
            log.error(f"Polling error: {e} — 5 sek garaşyp täzeden başlanýar")
            await asyncio.sleep(5)


if __name__ == "__main__":
    asyncio.run(main())
